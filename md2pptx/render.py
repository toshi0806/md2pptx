#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""IR（ir.py）→ pptx 描画（md2pptx ステージ2 / Phase 1）．

ステージ0（thmx2pptx.py）が生成した base pptx を土台に開き，IR の Deck を
走査して 1 つのプレゼンテーションへ描画する．配色・フォントはテーマ
（thmx）任せで，スクリプト側で色・フォントをハードコードしない（図形のみ
テーマのアクセント色を参照する）．

``参照スクリプト`` のモジュール大域に依存したヘルパ群（no_bullet /
add_slide_number / set_autonum / fit_body / box / note / block_arrow 等）を
Renderer のメソッドへ移植し，self.prs / self.layouts / テーマ色エイリアスから
状態を解決する（DESIGN.md §6）．移植したうちパイプラインが通らなくなった
4 つ（arrow / add_bullets / enum_items / content_slide）は削除した（Issue #75）．

使い方::

    from thmx2pptx import thmx_to_pptx
    from ir import Deck
    from render import build

    base = thmx_to_pptx("theme.thmx")
    build(deck, base, "out.pptx")
"""
from __future__ import annotations

import copy
import math
import os
import struct
import sys
from typing import TYPE_CHECKING, Any, Callable, NamedTuple

from pptx import Presentation
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt
# python-pptx の Slide は ir.Slide と名前がぶつかる．**外来のほうに印を付ける**
# ——``Slide`` は parser / cli でも ir の意味で使っており，この 1 ファイルのために
# 反転させると同じ識別子がファイルごとに別の型を指すことになる（Issue #35）．
# ``pptx.Presentation`` は**関数**（ファイルを開くファクトリ）で型ではない．
# 注釈にはクラスのほうが要る．
from pptx.presentation import Presentation as PptxPresentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.connector import Connector
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.slide import Slide as PptxSlide, SlideLayout
from pptx.text.text import TextFrame

from .colors import parse_color
from .ir import (
    TITLE_LAYOUT, Arrow, Block, Crop, Deck, Flow, Image, Length, Line,
    is_object_block, ObjectBlock, Seq, Slide, Table, TitleSlide,
)
from .flow import FlowNode, plan_flow
from .parser import parse_content_line
from .seq import plan_seq
from . import workdir

if TYPE_CHECKING:
    # ``_Paragraph`` は python-pptx の**私有クラス**（段落に公開の別名が無い）．
    # 実行時に import しないのは，改名・削除されたときに**注釈のためだけの
    # 名前で起動ごと落とさない**ため——``from __future__ import annotations``
    # により注釈は文字列のままなので、実行には要らない．CI の typecheck が
    # 気づかせてくれる．
    from pptx.text.text import _Paragraph, _Run


# Table.aligns の寄せ名 → PowerPoint の段落水平アラインメント．
_TABLE_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

# コードブロックの既定の等幅フォント（front matter の ``mono_font`` で変えられる）．
# Windows / macOS の Office に同梱されていて，日本語混在時は欧文だけに効く
# （和文はテーマの東アジア用フォントへ落ちる．``font.name`` は latin にだけ書くため）．
DEFAULT_MONO_FONT = "Consolas"


def resolve_image_path(src: str, base_dir: str | None) -> str:
    """``Image.src`` の宣言パスを実際のファイルパスへ解決する．

    絶対パスはそのまま，相対パスは ``base_dir``（Markdown の置き場）基準．
    **存在確認はしない**——``--watch`` は「まだ置かれていない画像」も監視対象に
    入れたいので（置かれたら作り直す），解決と存在確認を分けてある．描画側は
    ``Renderer._resolve_image_path`` で解決したうえで存在を確かめる．
    """
    return src if os.path.isabs(src) else os.path.join(base_dir or ".", src)


def _read_image_size(path: str) -> tuple[int, int]:
    """画像（png / jpg）のピクセル寸法 (width, height) をヘッダ解析で取得する．

    python-pptx の内部クラスや Pillow に依存せず，標準ライブラリだけでファイル
    ヘッダを読む（対応形式は png / jpeg）．寸法が読めない場合は ValueError．
    """
    with open(path, "rb") as f:
        head = f.read(8)
        if head[:8] == b"\x89PNG\r\n\x1a\n":
            # PNG: シグネチャ(8B)の直後が IHDR チャンク．チャンク長(4B)＋チャンクタイプ
            # "IHDR"(4B) を読み飛ばすと，先頭に width/height（ビッグエンディアン 4B×2）．
            f.read(4 + 4)  # チャンク長(4) + チャンクタイプ"IHDR"(4)
            w, h = struct.unpack(">II", f.read(8))
            return w, h
        if head[:2] == b"\xff\xd8":  # JPEG（SOI）．SOF マーカーまで走査する．
            f.seek(2)
            while True:
                b = f.read(1)
                if not b:
                    break
                if b != b"\xff":
                    continue
                marker = f.read(1)
                while marker == b"\xff":  # 連続する 0xFF（フィルバイト）を読み飛ばす
                    marker = f.read(1)
                if not marker:
                    break
                m = marker[0]
                # SOF0..SOF15（0xC0..0xCF）に寸法．ただし DHT/JPG/DAC は除く．
                if 0xC0 <= m <= 0xCF and m not in (0xC4, 0xC8, 0xCC):
                    f.read(2 + 1)  # segment length(2) + precision(1)
                    h, w = struct.unpack(">HH", f.read(4))
                    return w, h
                seg = f.read(2)
                if len(seg) < 2:
                    break
                length = struct.unpack(">H", seg)[0]
                f.seek(length - 2, os.SEEK_CUR)  # このセグメントを読み飛ばす
    raise ValueError(f"cannot read image dimensions (png/jpeg only): {path}")


def is_dark(rgb: str | None) -> bool:
    """``"RRGGBB"`` が「白文字を載せるべき濃さ」か．

    WCAG の相対輝度（sRGB をガンマ戻ししてから重みづけ）で見る．素の
    ``0.299R+0.587G+0.114B`` だと cn2026-theme の accent2（#3B812F）と
    accent5（#E2CAAA）のような、緑が効いた色の判定を外す．

    閾値 0.35 は「白と黒のどちらがコントラスト比を稼げるか」の分かれ目
    （相対輝度 0.179）より明るめ．**WCAG の判定そのものではなく、読みやすさを
    採った経験則**——0.18〜0.35 は黒でも 4.5:1 を満たすが、投影して見ると
    白抜きのほうが読める（cn2026-theme の accent2 #3B812F がここに入る）．

    **色が引けなかったら濃いものとして扱う**（``None``）．塗ってあるのに黒文字だと
    読めない——読めるかもしれない側へ倒す．
    """
    if not rgb or len(rgb) != 6:
        return True
    try:
        ch = [int(rgb[i:i + 2], 16) / 255.0 for i in (0, 2, 4)]
    except ValueError:
        return True
    lin = [c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4
           for c in ch]
    return 0.2126 * lin[0] + 0.7152 * lin[1] + 0.0722 * lin[2] < 0.35


# ひらがな（U+3041-309F）とカタカナ（U+30A0-30FF）．``ー``（U+30FC）や小書きの
# かなも含む——並びの一部なので、同じ詰め方をする（Issue #156）．
# 全角の句読点（U+3001-3002）はこの外．そちらは 1em のままにしておく．
_KANA_MIN, _KANA_MAX = 0x3041, 0x30FF


class SpaceBefore(NamedTuple):
    """段落前のアキ 1 つぶん（``a:spcBef``）．

    OOXML は同じ「アキ」を 2 通りで書く——``spcPct`` は**フォントサイズに対する
    割合**、``spcPts`` は絶対値（pt）．読んだ場所ではサイズが分からないので、
    どちらなのかを持ったまま返し、pt に直すのは ``Renderer._para_height``．
    """

    value: float
    percent: bool


class Renderer:
    """IR を pptx へ描画するレンダラ．

    base pptx（テーマのみを持つ 0 枚構成）を開き，レイアウトとテーマ色
    エイリアスを初期化する．スライドはすべて新規追加で生成する
    （thmx 由来の base は本文スライドを持たない）．
    """

    # フロー box の幅見積もりに掛ける安全係数（_text_width_pt の楽観的見積もりを補正）．
    _BOX_W_SAFETY = 1.15

    # 相対フォントサイズ 1 段あたりの倍率（≈12.5%）．拡大は ×，縮小は ÷．
    # 絶対 pt はハードコードせず，テーマ既定サイズ（_body_font_levels）からの相対比のみ持つ．
    _SIZE_STEP_RATIO = 1.125
    # 相対サイズの下限・上限（極小化／巨大化を防ぐ安全クランプ．段数の暴走対策）．
    _SIZE_MIN_PT = 8.0
    _SIZE_MAX_PT = 96.0

    # run に付ける言語（_apply_text_language）．PowerPoint はこれで禁則処理を選ぶ．
    # 2 つは**対で意味を持つ**（lang が主たる言語，altLang がもう一方の字種の言語）．
    # 日本語版 PowerPoint も ja-JP／en-US の組で書き出す．片方だけ変えないこと．
    _LANG = "ja-JP"
    _ALT_LANG = "en-US"

    def __init__(self, base_pptx_path: str, base_dir: str | None = None) -> None:
        self.prs = Presentation(base_pptx_path)
        # 画像などの相対パスを解決する基準ディレクトリ（既定は Markdown ファイルの
        # 置き場．cli が渡す）．None なら実行時のカレントを基準にする．
        self.base_dir = base_dir
        # テーマに pptx を渡した場合，元々入っているスライド（テンプレート用の
        # プレースホルダ枚）が先頭に残らないよう，常に 0 枚から描画を始める．
        self._clear_slides()
        # python-pptx はスライド寸法を Optional で返す（寸法未設定の pptx がありうる）．
        # 以降の座標計算は全てこの 2 つを基準にするので，ここで一度だけ確定させる．
        if self.prs.slide_width is None or self.prs.slide_height is None:
            raise ValueError("theme has no slide size (slide_width/height)")
        self.SW = self.prs.slide_width
        self.SH = self.prs.slide_height
        # コードブロックの等幅フォント（§5.12）．render() が front matter の
        # ``mono_font`` で上書きする．Renderer を直接使う経路でも既定が要る．
        self._mono_font = DEFAULT_MONO_FONT

        # テーマのアクセント色（図形用）．テキスト色・フォントはテーマ任せ．
        # Phase 1 では box/block_arrow/note/table（Phase 2/3）が使うため保持のみ．
        self.A2 = MSO_THEME_COLOR.ACCENT_2       # 緑
        self.A6 = MSO_THEME_COLOR.ACCENT_6       # 緑（濃）
        self.T2 = MSO_THEME_COLOR.TEXT_2         # 緑系テキスト
        self.GOLD = MSO_THEME_COLOR.ACCENT_1     # 金
        self.BG = MSO_THEME_COLOR.BACKGROUND_1   # 背景（白）
        self.TX = MSO_THEME_COLOR.TEXT_1         # 本文色（黒）

        # フロー図 box の自動配色（テーマアクセント色を順番に割当）．
        self._box_palette = [self.T2, self.A6, self.A6, self.GOLD, self.A2]
        # テーマ色名（DSL の {accent6} 等）→ MSO_THEME_COLOR の対応表．
        self._theme_map = {
            "accent1": MSO_THEME_COLOR.ACCENT_1, "accent2": MSO_THEME_COLOR.ACCENT_2,
            "accent3": MSO_THEME_COLOR.ACCENT_3, "accent4": MSO_THEME_COLOR.ACCENT_4,
            "accent5": MSO_THEME_COLOR.ACCENT_5, "accent6": MSO_THEME_COLOR.ACCENT_6,
            "tx1": MSO_THEME_COLOR.TEXT_1, "tx2": MSO_THEME_COLOR.TEXT_2,
            "bg1": MSO_THEME_COLOR.BACKGROUND_1, "bg2": MSO_THEME_COLOR.BACKGROUND_2,
        }
        # マスターの txStyles 由来レベル別サイズ（pt）のキャッシュ（style 名 → 一覧）．
        self._master_levels: dict[str, list[float]] = {}
        self._indent_cache: list[int] | None = None
        self._spc_cache: list[SpaceBefore] | None = None
        self._theme_rgb_cache: dict[str, str] | None = None

        # レイアウト解決．title=0 / content=1 / section=2．
        layouts = self.prs.slide_layouts
        self.layouts = layouts
        self.title_layout = layouts[TITLE_LAYOUT]
        self.L1 = layouts[1] if len(layouts) > 1 else layouts[0]
        self.section_layout = layouts[2] if len(layouts) > 2 else self.L1

    # ------------------------------------------------------------ helpers
    def _clear_slides(self) -> None:
        """base pptx に既存のスライドがあれば取り除く（0 枚から描画するため）．"""
        sldIdLst = self.prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(qn("r:id"))
            if rId:
                try:
                    self.prs.part.drop_rel(rId)
                except KeyError:
                    pass
            sldIdLst.remove(sldId)

    def no_bullet(self, para: _Paragraph) -> None:
        """段落の行頭記号を消す（結論行など）．"""
        pPr = para._p.get_or_add_pPr()
        for tag in ("a:buChar", "a:buAutoNum"):
            e = pPr.find(qn(tag))
            if e is not None:
                pPr.remove(e)
        if pPr.find(qn("a:buNone")) is None:
            pPr.append(pPr.makeelement(qn("a:buNone"), {}))

    def add_slide_number(self, slide: PptxSlide) -> None:
        """スライド自身のレイアウトの番号プレースホルダ（idx==12）を複製して有効化．

        セクションスライド（レイアウト2）など L1 以外の上でも，そのスライドの
        実レイアウトから番号プレースホルダを取得する．
        """
        for lph in slide.slide_layout.placeholders:  # type: ignore[misc]
            if lph.placeholder_format.idx == 12:
                slide.shapes._spTree.append(copy.deepcopy(lph._element))
                return

    def set_autonum(self, p: _Paragraph, fmt: str = "arabicPeriod",
                    color: str | None = None, start: int | None = None) -> None:
        """段落の行頭記号を自動採番（1. 2. 3. …）に切り替える（enumerate 相当）．

        color にテーマ色名（例 "tx1"）を渡すと採番記号の色を指定する．
        start を渡すとその番号から数え始める（buAutoNum の startAt）．
        PowerPoint の自動採番は**プレースホルダごとに 1 から数え直す**ので，
        2 カラムに割った採番リストの続きを右カラムで書くにはこれが要る（Issue #107）．

        **呼び出し側は採番段落すべてに start を渡す．** PowerPoint は startAt の
        付いた段落の**次から数え直す**ので，先頭にだけ渡すと "8. 1. 2. 3. …" になる．
        番号を数えるのは ``_append_lines`` の ``counters``．
        """
        pPr = p._p.get_or_add_pPr()
        if color:
            for tag in ("a:buClrTx", "a:buClr"):
                el = pPr.find(qn(tag))
                if el is not None:
                    pPr.remove(el)
            buClr = pPr.makeelement(qn("a:buClr"), {})
            # 色名の解決は本文の行内装飾と同じリゾルバを使う（Issue #105）．
            # @autonum-color でも CSS の色名と16進が書けるのはこのため．
            kind, value = parse_color(color)
            tag = "a:schemeClr" if kind == "theme" else "a:srgbClr"
            buClr.append(buClr.makeelement(qn(tag), {"val": value}))
            pPr.insert(0, buClr)  # buClr は採番記号より前に置く
        attrs = {"type": fmt}
        if start is not None:
            attrs["startAt"] = str(start)
        bu = pPr.makeelement(qn("a:buAutoNum"), attrs)
        for tag in ("a:buChar", "a:buNone", "a:buAutoNum"):
            el = pPr.find(qn(tag))
            if el is not None:
                pPr.replace(el, bu)
                break
        else:
            pPr.append(bu)

    def fit_body(self, tf: TextFrame, scale: float | None = None) -> None:
        """本文プレースホルダに自動調整（normAutofit）を設定する．
        scale を与えると縮小率（%）を明示的に焼き込む．"""
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if scale is not None:
            bodyPr = tf._txBody.bodyPr
            na = bodyPr.find(qn("a:normAutofit"))
            if na is None:
                # auto_size 設定で生成されない環境向けのフォールバック
                na = bodyPr.makeelement(qn("a:normAutofit"), {})
                bodyPr.append(na)
            na.set("fontScale", str(int(scale * 1000)))

    def _body_placeholder(self, slide: PptxSlide) -> SlidePlaceholder | None:
        """本文プレースホルダ（idx==1）を返す．無ければ None．"""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                return ph
        return None

    def _warn_no_body(self, lines: list[Line]) -> None:
        """地の文の置き場所が無いことを知らせる（Issue #67）．

        本文プレースホルダは ``idx==1`` で探すが，その番号は PowerPoint が内部で
        振るもので画面に出ない．標準レイアウトを動かす分には変わらないものの，
        枠を消して作り直すと別の番号になりうる——テーマを自作する人が普通にやる
        操作で，そうなると**地の文がどこにも入らない**．

        止めはしない（表・図とタイトルは出せるし，pptx の保存自体は成功する）が，
        **黙るのが最悪**なので何を落としたかを言う．消えたのが自分の書いた文章だと
        気づけなければ，テーマを直しようがない．
        """
        if not lines:
            return
        # 先頭 30 字は「自分のどの文章か」を思い出すための手掛かりで，全文ではない
        # （全部出すと 1 行に収まらず，警告としてかえって読みにくい）．
        head = lines[0].text.strip().replace("\n", " ")[:30]
        sys.stderr.write(
            f"md2pptx: warning: no body placeholder (idx 1) on this layout; "
            f"dropped {len(lines)} line(s) of body text starting {head!r}\n"
        )

    def _master_style_levels(self, style: str) -> list[float]:
        """マスターの ``p:txStyles/p:<style>`` のレベル別サイズ（pt）を返す（lvl1 始まり）．

        style は ``"bodyStyle"`` / ``"titleStyle"``．取得できなければ空リストを返す
        （既定値は呼び出し側が当てる——本文と表題では妥当な既定が違う）．
        """
        cached = self._master_levels.get(style)
        if cached is not None:
            return cached
        levels: list[float] = []
        try:
            master = self.prs.slide_masters[0]
            root = master.element.find(
                qn("p:txStyles") + "/" + qn("p:" + style))
            if root is not None:
                for lvl in range(1, 10):
                    el = root.find(
                        qn("a:lvl%dpPr" % lvl) + "/" + qn("a:defRPr"))
                    if el is not None and el.get("sz"):
                        levels.append(int(el.get("sz")) / 100.0)
        except Exception:
            pass
        self._master_levels[style] = levels
        return levels

    def _indent_for(self, level: int) -> int:
        """そのレベルの左インデント（EMU）．取れなければ 0．"""
        ind = self._body_indents()
        return ind[min(level, len(ind) - 1)] if ind else 0

    def _body_indents(self) -> list[int]:
        """マスター本文スタイルのレベル別の左インデント（EMU．lvl1 始まり）．

        ``{box}`` の枠を文字の始まりに合わせるために要る（Issue #133）．
        取れなければ空リストを返し，呼び出し側が 0 を当てる——枠が少し左へ出るだけで、
        位置の見積もりそのものは壊れない．
        """
        if self._indent_cache is not None:
            return self._indent_cache
        levels: list[int] = []
        try:
            master = self.prs.slide_masters[0]
            root = master.element.find(
                qn("p:txStyles") + "/" + qn("p:bodyStyle"))
            if root is not None:
                # レベルが飛んでいても打ち切らない．テーマが lvl2 だけ書かない
                # ことはありうるので、欠けたレベルは**直前の値で埋める**
                # （0 に落とすと、そのレベルだけ枠が左へ飛び出す）．
                last = 0
                for lvl in range(1, 10):
                    el = root.find(qn("a:lvl%dpPr" % lvl))
                    marL = el.get("marL") if el is not None else None
                    if marL:
                        last = int(marL)
                    levels.append(last)
        except Exception:
            levels = []
        self._indent_cache = levels
        return levels

    def _theme_rgb(self, name: str) -> str | None:
        """テーマ色名（DSL の ``bg2`` / ``accent2`` …）を ``"RRGGBB"`` へ解決する．

        **``p:clrMap`` を必ず踏む**——``bg2`` が指すのはテーマの ``lt2`` で、
        名前がそのまま clrScheme のタグ名になっているのは accent1〜6 と
        hlink / folHlink だけ．踏まずに引くと別の色を見る（Issue #148）．

        引けなければ None．テーマによっては色が ``srgbClr`` ではなく
        ``sysClr``（``lastClr`` に実効値）で書いてあるので、そちらも見る．
        """
        if self._theme_rgb_cache is None:
            self._theme_rgb_cache = self._read_color_scheme()
        return self._theme_rgb_cache.get(name)

    def _read_color_scheme(self) -> dict[str, str]:
        """テーマの clrScheme を、DSL の色名で引ける辞書にして返す．

        見るのは ``slide_masters[0]`` だけ．md2pptx はテーマ 1 つ・マスター 1 つを
        前提に描いており（``_master_style_levels`` なども同じ）、複数マスターの
        pptx を作る手立ては無い．
        """
        out: dict[str, str] = {}
        try:
            master = self.prs.slide_masters[0]
            part = master.part.part_related_by(RT.THEME)
            root = parse_xml(part.blob)
            scheme = root.find(
                qn("a:themeElements") + "/" + qn("a:clrScheme"))
            if scheme is None:
                return out
            raw: dict[str, str] = {}
            for el in scheme:
                key = str(el.tag).rsplit("}", 1)[-1]
                for tag, attr in ((qn("a:srgbClr"), "val"),
                                  (qn("a:sysClr"), "lastClr")):
                    child = el.find(tag)
                    if child is not None and child.get(attr):
                        raw[key] = str(child.get(attr)).upper()
                        break
            # clrMap（bg1="lt1" 等）で DSL の名前へ引き直す．属性が無いテーマは
            # 恒等写像とみなす（accent1〜6 / hlink は元から同名）．
            cmap = master.element.find(qn("p:clrMap"))
            for name in self._theme_map:
                mapped = cmap.get(name) if cmap is not None else None
                val = raw.get(str(mapped) if mapped else name)
                if val:
                    out[name] = val
        except Exception:
            return out
        return out

    def _body_space_before(self) -> list[SpaceBefore]:
        """マスター本文スタイルのレベル別の段落前アキ（lvl1 始まり）．

        **帯の高さはここを数えないと足りない**（Issue #145）．cn2026-theme は
        全レベルに 20% を持ち、30pt の段落なら 0.21cm、4段落で 0.85cm ずれる．

        ``spcPct``（フォントサイズに対する％）は**その場では pt に直せない**ので、
        ``SpaceBefore`` の ``percent`` を立てて返し、サイズを掛けるのは
        ``_para_height`` に任せる．``spcPts``（絶対値）はそのまま pt．

        アキを書いていないレベルは**直前のレベルの値を引き継ぐ**．テーマは
        上位レベルだけ書いて下位を省くことがあり、そこを 0 と読むと深い階層の
        段落だけ高さが足りなくなる．取得に失敗したら空リスト——呼び出し側が
        0 を当てる（アキを数えない従来の見積もりに戻るだけ）．
        """
        if self._spc_cache is not None:
            return self._spc_cache
        levels: list[SpaceBefore] = []
        try:
            master = self.prs.slide_masters[0]
            root = master.element.find(
                qn("p:txStyles") + "/" + qn("p:bodyStyle"))
            last = SpaceBefore(0.0, False)
            if root is not None:
                for lvl in range(1, 10):
                    el = root.find(qn("a:lvl%dpPr" % lvl))
                    spc = el.find(qn("a:spcBef")) if el is not None else None
                    if spc is not None:
                        pct = spc.find(qn("a:spcPct"))
                        pts = spc.find(qn("a:spcPts"))
                        if pct is not None and pct.get("val"):
                            last = SpaceBefore(
                                int(pct.get("val")) / 100000.0, True)
                        elif pts is not None and pts.get("val"):
                            last = SpaceBefore(int(pts.get("val")) / 100.0,
                                               False)
                        else:
                            last = SpaceBefore(0.0, False)
                    levels.append(last)
        # OOXML の探りは軒並みこの形（_master_style_levels / _layout_level_sizes
        # と同じ）．テーマの作りは千差万別で、読めなければ既定へ落ちれば済む．
        except Exception:
            levels = []
        self._spc_cache = levels
        return levels

    # 行送りの倍率．**PowerPoint で実測した値**（Issue #152）．
    # cn2026-theme（BIZ UDPゴシック・30pt）で ``spcBef`` を 0 にして段落の送りを
    # 測ると 36.0pt ちょうど＝ 30 × 1.20．以前は保守的な 1.32 を置いていたが、
    # ``{box}`` は**位置**を出すのに使うので、大きめに見ることに意味が無い
    # （枠が段落の下半分から次の項目へ掛かっていた）．
    _LINE = 1.20

    # 枠（``{box}``）を行の箱より少し上へずらす割合（Issue #160）．
    # PowerPoint は行の箱の中で字を**上寄り**に置き、下に descent ぶんの空きを
    # 残す．行の箱にそのまま合わせると、字に対して枠が下がって見える——実測で
    # 上 4.8pt / 下 11pt（30pt の行・箱 36pt）．差の半分だけ持ち上げて、
    # 上下の空きを揃える．位置によらず一定なので、積もりはしない．
    _BOX_LIFT = 0.086

    # 折り返し判定の許容幅．PowerPoint は日本語を詰めて改行を避けるので、
    # 幅を数％超えただけでは折り返らない（Issue #156）．
    _WRAP_SLACK = 1.05

    @classmethod
    def _line_height(cls, size_pt: float) -> int:
        """文字 1 行ぶんの高さ（EMU）．"""
        return int(Pt(size_pt * cls._LINE))

    def _space_before(self, level: int, size_pt: float) -> int:
        """そのレベルの段落前アキ（EMU）．

        ``level`` は 0 始まり（``Line.level``）、``_body_space_before`` のリストは
        lvl1 始まり——先頭が level 0 に当たるので添字はそのまま。テーマが書いて
        いない深さは末尾で頭打ち。``_body_font_levels`` を引くときと同じ数え方
        （render 全体で揃えてある）。
        """
        spc = self._body_space_before()
        raw = spc[min(level, len(spc) - 1)] if spc else SpaceBefore(0.0, False)
        if not raw.percent:
            return int(Pt(raw.value))          # spcPts はそのまま pt
        # **``spcPct`` はフォントサイズではなく「行の高さ」に対する割合**
        # （Issue #152）．30pt・20% の段落を測ると 43.2pt ＝ 36.0 × 1.20 で、
        # 増えぶんは 7.2 ＝ 0.20 × 36.0．フォントに掛けると 6.0 で足りない．
        return int(raw.value * self._line_height(size_pt))

    def _para_height(self, level: int, size_pt: float, lines: int = 1) -> int:
        """段落 1 つぶんの高さ（EMU）．行の高さ × 行数 ＋ 段落前アキ．

        アキを足さないと帯が上へずれ、図が地の文に食い込む（Issue #145）。
        **アキは段落に 1 回だけ**——折り返した行ごとに足すと、2 行の項目を囲む
        ``{box}`` が 1 行ぶん下へ伸びて次の項目に掛かる（Issue #150）。
        """
        return lines * self._line_height(size_pt) + self._space_before(
            level, size_pt)

    def _body_font_levels(self) -> list[float]:
        """マスター本文スタイルのレベル別フォントサイズ（pt）を返す（lvl1 始まり．既定 [18]）．

        テーマが持つ「本文サイズの梯子」で，表・図が標準サイズで収まらないときに
        下位レベルの小さいサイズへ段階的に切り替えるために用いる（``_fit_font``）．
        **プレースホルダへ流す文字の基点にはこれを使わないこと**——レイアウトによる
        上書きを含む ``_frame_font_levels`` が実際に描かれるサイズを返す（Issue #83）．
        """
        return self._master_style_levels("bodyStyle") or [18.0]

    def _layout_level_sizes(self, layout: SlideLayout, idx: int) -> dict[int, float]:
        """レイアウトの idx プレースホルダが ``a:lstStyle`` で上書きするサイズ（pt）．

        戻り値はレベル番号（1 始まり）→ pt の辞書．上書きの無いレベルは含まない
        （テーマは一部のレベルだけ上書きすることがあるため，欠けは呼び出し側が
        マスター側の値で埋める）．
        """
        sizes: dict[int, float] = {}
        try:
            # python-pptx の LayoutPlaceholders は __iter__ を Callable 属性として
            # 宣言しており，mypy が "self を取らない" と誤検出する（_effective_geom 同様）．
            for lph in layout.placeholders:  # type: ignore[misc]
                if lph.placeholder_format.idx != idx:
                    continue
                lst = lph._element.find(
                    qn("p:txBody") + "/" + qn("a:lstStyle"))
                if lst is not None:
                    for lvl in range(1, 10):
                        el = lst.find(
                            qn("a:lvl%dpPr" % lvl) + "/" + qn("a:defRPr"))
                        if el is not None and el.get("sz"):
                            sizes[lvl] = int(el.get("sz")) / 100.0
                break
        except Exception:
            pass
        return sizes

    def _frame_font_levels(self, tf: TextFrame) -> list[float]:
        """tf に流す文字のレベル別既定サイズ（pt）を返す（lvl1 始まり）．

        相対サイズ（``{±n}``）の基点は「その文字が実際に出るサイズ」でなければ
        意味を持たない（Issue #83）．PowerPoint の継承順に合わせ，**レイアウトの
        プレースホルダの ``a:lstStyle`` をマスターの ``txStyles`` より優先**する．
        テーマは一部のレベルだけ上書きすることがあるのでレベル単位で重ねる．

        マスター側プレースホルダの ``a:lstStyle`` は見ていない．継承順では
        レイアウトとマスター txStyles の間に入るが，手元のテーマはどれもそこに
        サイズを持たず，足しても通らない経路が増えるだけになる．
        プレースホルダでない図形（テキストボックス等）は本文スタイルを基点とする．
        """
        ptype, idx = self._frame_placeholder(tf)
        is_title = ptype in ("title", "ctrTitle")
        base = self._master_style_levels(
            "titleStyle" if is_title else "bodyStyle")
        if not base:
            base = [42.0] if is_title else [18.0]
        if idx is None:
            return base
        try:
            # tf.part はこのテキストフレームを含む SlidePart．そこから所属スライドの
            # レイアウトを引く（スライド側の描画中にしか呼ばれない）．
            layout = tf.part.slide.slide_layout  # type: ignore[attr-defined]
        except Exception:
            return base
        over = self._layout_level_sizes(layout, idx)
        if not over:
            return base
        # 必要なレベル数＝マスター側の段数と，上書きが触れる最大レベルの大きいほう．
        # over のキーはレベル番号（1 始まり）なので len(base) と同じ土俵で比べられる．
        n = max(len(base), max(over.keys()))
        return [over.get(i + 1, base[min(i, len(base) - 1)]) for i in range(n)]

    @staticmethod
    def _frame_placeholder(tf: TextFrame) -> tuple[str, int | None]:
        """tf を持つ図形のプレースホルダ種別と idx を返す（プレースホルダ以外は idx=None）．

        ``p:ph`` は属性を省略できる（``type`` 既定 ``body`` / ``idx`` 既定 0）ので，
        欠けている場合は既定を当てる．
        """
        try:
            sp = tf._txBody.getparent()
            el = sp.find(
                qn("p:nvSpPr") + "/" + qn("p:nvPr") + "/" + qn("p:ph"))
        except Exception:
            el = None
        if el is None:
            return "body", None
        return el.get("type") or "body", int(el.get("idx") or 0)

    def _body_font_size(self) -> float:
        """本文プレースホルダの標準フォントサイズ（pt．lvl1）を返す．"""
        return self._body_font_levels()[0]

    def _apply_size_delta(self, p: _Paragraph, level: int,
                          delta: int | None, levels: list[float]) -> None:
        """段落 p に相対フォントサイズ（delta 段）を適用する．

        基点はその段落の level に対応するテーマ既定サイズ．levels は段落を置く
        枠の実効サイズ（``_frame_font_levels``）で，呼び出し側が枠ごとに 1 度
        解決して渡す（Issue #83．行ごとに引き直さないための引数でもある）．
        実サイズ = round(base × 1.125**delta) を [_SIZE_MIN_PT, _SIZE_MAX_PT] で
        クランプする（大きな段数指定でも極小・巨大化しない）．p.level（インデント）は
        変更しない＝段落の既定文字書式（defRPr＝p.font）にサイズを設定するため，
        run が無い空段落でも有効で，bullet/採番記号も本文と同じサイズになる．

        - delta is None（未指定）: 何もしない（スライド既定もテーマ既定もそのまま）．
        - delta == 0（テーマ既定に固定）: サイズ指定を明示的に外しテーマ継承へ戻す．
        """
        if delta is None:
            return
        if delta == 0:
            # スライド既定（@body-size）を無効化し，テーマ既定サイズへ戻す．
            p.font.size = None
            return
        base = levels[min(level, len(levels) - 1)]
        size: float = round(base * self._SIZE_STEP_RATIO ** delta)
        size = min(self._SIZE_MAX_PT, max(self._SIZE_MIN_PT, size))
        p.font.size = Pt(size)

    def _write_spans(self, p: _Paragraph, blk: Line) -> None:
        """Line を段落へ書く．装飾があれば run を分けて属性を付ける（§5.13）．

        装飾の無い行（``blk.spans`` が空）は従来どおり ``p.text`` へ一括で入れる——
        python-pptx が ``\\v`` を ``a:br`` に割ってくれるので，触る理由がない．
        **装飾を使わない原稿の出力はこの経路のままで 1 ビットも変わらない．**

        装飾がある行は run を自分で並べる．セグメントの境目（``segment`` が変わる位置）
        には ``a:br`` を挟む——``p.text`` に任せられないのは，1 セグメントが
        複数 run に割れるため．
        """
        if not blk.spans:
            p.text = blk.text
            return
        p.text = ""
        prev_seg = blk.spans[0].segment
        for i, span in enumerate(blk.spans):
            if i and span.segment != prev_seg:
                self._append_break(p)
                prev_seg = span.segment
            run = p.add_run()
            run.text = span.text
            if span.bold:
                run.font.bold = True
            if span.mono:
                run.font.name = self._mono_font
            if span.color:
                self._set_run_color(run, span.color)
            if span.link:
                run.hyperlink.address = span.link
            if span.script:
                self._set_run_script(run, span.script)

    @staticmethod
    def _append_break(p: _Paragraph) -> None:
        """段落の末尾（``a:endParaRPr`` の前）へ ``a:br`` を足す．"""
        br = p._p.makeelement(qn("a:br"), {})
        end = p._p.find(qn("a:endParaRPr"))
        if end is None:
            p._p.append(br)
        else:
            end.addprevious(br)

    def _set_run_color(self, run: _Run, name: str) -> None:
        """run の文字色を設定する（テーマ色名／CSS の色名／16進）．

        テーマ色は **RGB へ潰さず** ``theme_color`` で指定する——
        テーマを差し替えたときに追従させたいため．

        ``Span.color`` はパーサが正規化済みだが，ここでも ``parse_color`` を通す．
        冪等（テーマ色名も "#RRGGBB" もそのまま返る）なうえ，**色名の語彙を知る
        場所を 1 つに保てる**——先頭の "#" で振り分けると，render が
        「正規化済みである」という書かれざる前提に依存することになる．
        """
        kind, value = parse_color(name)
        if kind == "theme":
            run.font.color.theme_color = self._theme_map[value]
        else:
            run.font.color.rgb = RGBColor.from_string(value)

    @staticmethod
    def _set_run_script(run: _Run, script: str) -> None:
        """run を上付き／下付きにする（``a:rPr/@baseline``）．

        python-pptx に API が無いので属性を直接書く．値は OOXML の千分率で，
        PowerPoint の UI が付けるのと同じ ±30% にしている．
        """
        run.font._rPr.set("baseline", "30000" if script == "sup" else "-25000")

    def _apply_segment_deltas(self, p: _Paragraph, deltas: list[int | None],
                              base_pt: float,
                              segments: list[int] | None = None) -> None:
        """段落 p の run へ，セグメントごとの相対サイズを適用する（Issue #82）．

        segments を渡すと run とセグメントの対応をそれで決める（Issue #105）．
        **行内装飾があると 1 セグメントが複数 run に割れる**ので，位置で対応させると
        別のセグメントへサイズが付く．渡さないときは従来どおり位置で対応させる．

        ``<br>`` を含む段落は python-pptx が ``\\v`` ごとに run を分けて ``a:br`` で
        つなぐので，run と IR のセグメントが順に 1 対 1 で対応する．行や段落を分けずに
        **一部だけサイズを変えられる**のはこの経路だけ——タイトル枠に副題を収める
        （Issue #82）ような使い方がこれに当たる．

        deltas は IR 側で**セグメント数と同じ長さに正規化済み**（``__post_init__``）で，
        run 数もセグメント数に一致するので，``zip`` は常に全要素を回る．
        短いほうで止まるのは，万一ずれたときに**別のセグメントへサイズを付けない**
        ための安全側の性質——見た目の崩れは目に付くが，1 つずれたサイズは
        正しく見えてしまうので，付けないほうが害が小さい．

        本文が空の段落（``- `` 由来）は run が 0 個で，このループは回らない．
        """
        idx = segments if segments is not None else range(len(deltas))
        for run, seg in zip(p.runs, idx):
            delta = deltas[seg] if seg < len(deltas) else None
            if delta is None:
                continue
            size = round(base_pt * self._SIZE_STEP_RATIO ** delta)
            run.font.size = Pt(
                min(self._SIZE_MAX_PT, max(self._SIZE_MIN_PT, size)))

    def _title_font_size(self) -> float:
        """マスター表題スタイルの既定フォントサイズ（pt．lvl1）を返す（既定 42）．

        **タイトルの描画には使っていない**——``title:`` の段落に ``sz`` は書かず，
        実サイズはテーマ（レイアウトの CENTER_TITLE）が決める．唯一の用途は
        front matter の ``subtitle:`` の基点で，そこが凍結されている理由は
        ``render_title_slide`` のコメントを参照（Issue #83 / #82）．
        """
        try:
            master = self.prs.slide_masters[0]
            el = master.element.find(
                qn("p:txStyles") + "/" + qn("p:titleStyle")
                + "/" + qn("a:lvl1pPr") + "/" + qn("a:defRPr"))
            if el is not None and el.get("sz"):
                return int(el.get("sz")) / 100.0
        except Exception:
            pass
        return 42.0

    def _size_from_delta(self, base_pt: float, delta: int | None) -> float:
        """基点サイズ base_pt（pt）に相対段数 delta を適用した pt 値を返す（範囲クランプ）．

        本文の _apply_size_delta と同じ 1.125 倍/段の比率．delta が None なら base をそのまま返す．
        """
        if delta is None:
            return base_pt
        size = round(base_pt * self._SIZE_STEP_RATIO ** delta)
        return min(self._SIZE_MAX_PT, max(self._SIZE_MIN_PT, size))

    @staticmethod
    def _text_width_pt(text: str | None, font_pt: float) -> float:
        """テキストの概算表示幅（pt）．

        cn2026-theme（BIZ UDPゴシック）・30pt で PDF から実測した文字送り：

        =========== ========= ======
        種類        1文字      em 比
        =========== ========= ======
        カタカナ    26.53pt   0.885
        漢字        29.87pt   0.996
        ASCII       15.65pt   0.522
        =========== ========= ======

        **かなだけが 1em より狭い**（Issue #156）．日本語のプロポーショナル
        フォント（Yu Gothic / Meiryo / BIZ UDP…）はかなを詰めるため．全角を
        一律 1em で数えていた頃は折り返しを多く見積もり、``{box}`` の枠が
        2 行ぶんの高さになって次の項目まで覆っていた．

        等幅の日本語フォント（MS ゴシックなど）では逆に 1 割ぶん小さく見積もる．
        いま広く使われている日本語フォントに寄せた**概算**．
        """
        w = 0.0
        for ch in text or "":
            c = ord(ch)
            # 0x2E80（CJK 部首補助の先頭）より前は半角として扱う．ここには
            # ギリシャ文字や記号（U+2600〜 など）も入り、フォントによっては
            # 全角で出るが、講義スライドで使うのは ASCII がほとんどなので
            # この境目のままにしてある（**概算**．Issue #156）．
            if c <= 0x2E80:
                w += font_pt * 0.55         # 半角
            elif _KANA_MIN <= c <= _KANA_MAX:
                w += font_pt * 0.885        # かな（長音符・小書きを含む）
            else:
                w += font_pt                # 漢字・全角記号
        return w

    @classmethod
    def _wrapped_lines(cls, text: str | None, font_pt: float,
                       avail_pt: float) -> int:
        """``avail_pt`` の幅に流したときの行数（1 以上）．

        **僅差の超過では折り返さない**（Issue #156）．PowerPoint は日本語を
        少し詰めて改行を避ける（``eaLnBrk`` / ``hangingPunct``）ので、幅を
        1〜数％超えただけの行は 1 行のまま出る．cn2026-02 のシラバスの
        「ネットワークコミュニケーション」は実測 397.9pt・使える幅 382.6pt
        （4.0% 超過）で、1 行に収まっている．

        取り違える向きが問題になるのは ``{box}`` で、**折り返すと決めつけると
        枠が次の項目まで覆う**．収まると見て外したときは枠が短くなるだけなので、
        迷ったら「折り返さない」へ倒す．
        """
        if avail_pt <= 0:
            return 1
        w = cls._text_width_pt(text, font_pt)
        return max(1, math.ceil(w / (avail_pt * cls._WRAP_SLACK)))

    def _fit_font(self, fits_at: Callable[[float], bool]) -> float:
        """レベル別サイズを大きい順に試し，``fits_at(size)`` が真の最大サイズを返す．

        どのレベルでも収まらなければ最小レベル（最後）のサイズを返す（ベストエフォート）．
        """
        levels = self._body_font_levels()
        for sz in levels:
            if fits_at(sz):
                return sz
        return levels[-1]

    # --------------------------------------------------------- title slide
    def render_title_slide(self, ts: TitleSlide) -> PptxSlide:
        """タイトルスライドを 1 枚目として追加する（番号は付けない）．

        base は 0 枚構成のため，タイトルレイアウト上に新規スライドを追加し，
        CENTER_TITLE（idx==0）に title を，SUBTITLE（idx==1）に
        subtitle＋author＋affiliation を流し込む．
        """
        s = self.prs.slides.add_slide(self.title_layout)

        title_ph = self._find_placeholder(s, 0)
        if title_ph is not None and ts.title:
            tf = title_ph.text_frame
            lines = ts.title.split("\n")
            tf.paragraphs[0].text = lines[0]
            for ln in lines[1:]:
                tf.add_paragraph().text = ln
            # 副題はタイトル枠内に，本文より少し小さめの文字で入れる．
            if ts.subtitle:
                sp = tf.add_paragraph()
                sp.text = ts.subtitle
                sp.space_before = Pt(6)
                # 副題だけは基点が「実際に描かれるサイズ」になっていない．
                # _title_font_size() はマスターの titleStyle（42pt 等）で，
                # タイトルが実際に出るサイズ（レイアウト 0 の上書き．60pt / 50pt）
                # ではない．**承知のうえで現状の値のまま凍結する**（Issue #83）．
                # front matter の表紙記述は Issue #82 で本文記法へ移して非推奨に
                # するので，この経路は消える．非推奨の間は既存原稿が同じ見た目で
                # 動き続けるほうが価値があり，消える直前に見た目を変える益がない．
                # 直すのではなく，#82 で「副題＝タイトル枠に置く行」になったとき，
                # 他の行と同じ _frame_font_levels 基点へ自然に合流させること．
                sub_sz = self._size_from_delta(
                    self._title_font_size() * 0.8, ts.subtitle_delta)
                for r in sp.runs:
                    r.font.size = Pt(sub_sz)

        # 副題プレースホルダには著者・所属のみを入れる（副題はタイトル枠へ移動）．
        sub_ph = self._find_placeholder(s, 1)
        if sub_ph is not None:
            # 各行の相対サイズ段数（{-1} 等）を行と 1 対 1 で持ち回る（None＝未指定）．
            # affiliation_deltas は TitleSlide.__post_init__ で affiliation と同長が保証される．
            sub_lines = []
            sub_deltas = []
            if ts.author:
                sub_lines.append(ts.author)
                sub_deltas.append(ts.author_delta)
            for aff, delta in zip(ts.affiliation or [], ts.affiliation_deltas):
                sub_lines.append(aff)
                sub_deltas.append(delta)
            if sub_lines:
                # 所属行の折り返しを抑えるため右方向へ枠を広げる（左位置は維持）．
                # 継承ジオメトリの場合は 4 辺すべてを実効値で明示する
                # （一部だけ設定すると top/height が 0 に落ちて枠が移動するため）．
                left, top, width, height = self._effective_geom(sub_ph, s)
                if None not in (left, top, width, height):
                    new_w = self.SW - left - Inches(0.2)
                    if new_w > width:
                        sub_ph.left = left
                        sub_ph.top = top
                        sub_ph.height = height
                        sub_ph.width = new_w
                tf = sub_ph.text_frame
                tf.paragraphs[0].text = sub_lines[0]
                for ln in sub_lines[1:]:
                    tf.add_paragraph().text = ln
                # {-1}/{+1} 指定のある行だけ副題既定サイズを基点に段階調整する．
                # 基点は本文行と同じ解決を通す（同じ枠に出る同じ記法が経路で
                # 違うサイズになっていた．Issue #83）．
                base = self._frame_font_levels(tf)[0]
                for para, delta in zip(tf.paragraphs, sub_deltas):
                    if delta is not None:
                        para.font.size = Pt(self._size_from_delta(base, delta))
        self._set_notes(s, ts.notes)
        return s

    def _set_notes(self, slide: PptxSlide, notes: str | None) -> None:
        """発表者ノート（```note 由来）を notes slide へ書き込む（§5.10）．

        notes_slide への初回アクセスで python-pptx がノートスライドを生成する．
        None・空文字なら何もしない（ノートスライド自体を作らない）．
        text 代入は既存段落をすべて消して置き換える（python-pptx の setter は
        clear_content 後に "\\n" 区切りで段落を再構築する）．md2pptx は常に
        0 枚の base から描画する（_clear_slides）ため，ここで消える既存ノートはない．
        """
        if notes:
            # ``notes_text_frame`` は**ノート用プレースホルダが無いとき** None を
            # 返す（python-pptx は ``notes_placeholder is None`` でそう決めている）．
            # 止めはしない——ノートは補助情報で，ここで落とすとスライド自体が
            # 出せなくなる．ただし**黙らない**（Issue #67 と同じ扱い）．書いた
            # ノートが出ないことは pptx を開くまで分からず，開いても原因が
            # テーマ側だとは思い当たらない．
            tf = slide.notes_slide.notes_text_frame
            if tf is None:
                head = notes.strip().replace("\n", " ")[:30]
                sys.stderr.write(
                    "md2pptx: warning: this theme's notes master has no "
                    "placeholder for notes text; dropped the speaker notes "
                    f"starting {head!r}\n")
                return
            tf.text = notes

    def _effective_geom(self, ph: SlidePlaceholder, slide: PptxSlide) -> tuple[
            int | None, int | None, int | None, int | None]:
        """プレースホルダの実効ジオメトリ (left, top, width, height) を返す．

        スライド上で未指定（継承）の値はレイアウトの同 idx プレースホルダで補う．
        """
        left, top, width, height = ph.left, ph.top, ph.width, ph.height
        if None in (left, top, width, height):
            idx = ph.placeholder_format.idx
            try:
                for lph in slide.slide_layout.placeholders:  # type: ignore[misc]
                    if lph.placeholder_format.idx == idx:
                        left = left if left is not None else lph.left
                        top = top if top is not None else lph.top
                        width = width if width is not None else lph.width
                        height = height if height is not None else lph.height
                        break
            except Exception:
                pass
        return left, top, width, height

    def _find_placeholder(self, slide: PptxSlide,
                          idx: int) -> SlidePlaceholder | None:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    # ----------------------------------------------------------- 図形（flow）
    def box(self, slide: PptxSlide, l: int, t: int, w: int, h: int, text: str,
            tc: MSO_THEME_COLOR, sub: str | None = None,
            fsize: float | None = None, ssize: float | None = None) -> Shape:
        """角丸四角ノードを描く（塗りはテーマ色 tc，文字は背景色 BG）．

        fsize / ssize（pt）を省略するとテーマ既定のフォントサイズを継承する．
        """
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
        shp.fill.solid()
        shp.fill.fore_color.theme_color = tc
        shp.line.color.theme_color = self.TX
        shp.line.width = Pt(0.5)
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
        pa = tf.paragraphs[0]
        pa.alignment = PP_ALIGN.CENTER
        pa.text = text
        for r in pa.runs:
            r.font.color.theme_color = self.BG
            if fsize is not None:
                r.font.size = Pt(fsize)
            r.font.bold = True
        if sub:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.text = sub
            for r in p2.runs:
                r.font.color.theme_color = self.BG
                if ssize is not None:
                    r.font.size = Pt(ssize)
        return shp

    def line(self, slide: PptxSlide, x1: int, y1: int, x2: int, y2: int,
             color: MSO_THEME_COLOR | None = None, width_pt: float = 1.0,
             dashed: bool = False, arrow: bool = False) -> Connector:
        """2 点を結ぶ線を引く（Issue #108）．

        ``block_arrow`` はノード間の**すき間に収まる塗り矢印**で，box に食い込ませない
        ための道具．こちらは任意の 2 点を結ぶ細い線で，シーケンス図のライフライン
        （縦線）やメッセージ（斜めの矢印）のように**すき間ではなく図の骨格**を描くのに使う．

        ``arrow`` を立てると終点に矢じりが付き，``dashed`` で破線になる
        （時間の経過や省略を表す線に使う）．色はテーマ任せ（既定は本文色）．
        """
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Emu(int(x1)), Emu(int(y1)),
            Emu(int(x2)), Emu(int(y2)))
        conn.line.color.theme_color = color or self.TX
        conn.line.width = Pt(width_pt)
        if dashed:
            conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        if arrow:
            # 矢じりは python-pptx に API が無いので ``a:ln`` へ直接書く．
            ln = conn.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:tailEnd"),
                                     {"type": "triangle", "w": "med", "len": "med"}))
        return conn

    def block_arrow(self, slide: PptxSlide, x1: int, y1: int, x2: int, y2: int,
                    thickness: int,
                    color: MSO_THEME_COLOR | None = None) -> Shape:
        """ノード間のすき間に塗りつぶしのブロック矢印を置く．

        太い直線＋三角矢じりは box に食い込み見栄えが悪いため，すき間に収まる
        塗り矢印シェイプ（RIGHT/DOWN_ARROW）を用いる．色はテーマ任せ（既定はアクセント）．
        """
        inset = Inches(0.05)
        if abs(x2 - x1) >= abs(y2 - y1):       # 横向き
            left = min(x1, x2) + inset
            width = abs(x2 - x1) - 2 * inset
            if width <= 0:
                left, width = min(x1, x2), abs(x2 - x1)
            height = thickness
            top = y1 - height // 2
            shape = MSO_SHAPE.RIGHT_ARROW if x2 >= x1 else MSO_SHAPE.LEFT_ARROW
        else:                                   # 縦向き
            top = min(y1, y2) + inset
            height = abs(y2 - y1) - 2 * inset
            if height <= 0:
                top, height = min(y1, y2), abs(y2 - y1)
            width = thickness
            left = x1 - width // 2
            shape = MSO_SHAPE.DOWN_ARROW if y2 >= y1 else MSO_SHAPE.UP_ARROW
        shp = slide.shapes.add_shape(shape, Emu(int(left)), Emu(int(top)),
                                     Emu(int(width)), Emu(int(height)))
        shp.fill.solid()
        shp.fill.fore_color.theme_color = color or self.TX
        shp.line.fill.background()
        return shp

    def note(self, slide: PptxSlide, l: int, t: int, w: int, h: int, text: str,
             size: float, tc: MSO_THEME_COLOR | None = None, bold: bool = False,
             align: PP_ALIGN = PP_ALIGN.LEFT,
             anchor: MSO_ANCHOR | None = None, wrap: bool = True) -> Shape:
        """注記用テキストボックスを描く（キャプション・矢印ラベル・省略記号）．

        ``wrap=False`` は折り返さず、左右の余白も取らない．矢印ラベルのように
        **短くて折れては困る**文字に使う——枠に収まらないと単語の途中で割れて
        読めなくなる（Issue #111）．はみ出す先は box の上の何も無い場所なので害が無い．
        """
        tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        if not wrap:
            tf.margin_left = tf.margin_right = 0
        if anchor is not None:
            tf.vertical_anchor = anchor
        pa = tf.paragraphs[0]
        pa.alignment = align
        pa.text = text
        for r in pa.runs:
            if size is not None:
                r.font.size = Pt(size)
            if tc is not None:
                r.font.color.theme_color = tc
            if bold:
                r.font.bold = True
        return tb

    def _theme_color(self, name: object) -> MSO_THEME_COLOR | None:
        """テーマ色名を MSO_THEME_COLOR へ解決する（未知なら None）．"""
        if not name:
            return None
        return self._theme_map.get(str(name).lower())

    def _box_fits(self, node: FlowNode, bw: int, bh: int,
                  font_pt: float) -> bool:
        """box（主ラベル＋副ラベル）が指定フォントサイズで収まるか概算判定する．

        幅見積もりは安全係数 ``_BOX_W_SAFETY`` を掛けて保守的に評価する（``theme.thmx``
        のような半角主体ラベルが実 PowerPoint で 1 字あふれて折り返すのを防ぐ）．
        """
        line_h = font_pt * 1.2
        inner_w = max(1.0, bw / 12700.0 - 8)   # 左右マージン約 Pt(4)×2
        inner_h = bh / 12700.0 - 4             # 上下マージン約
        safe = self._BOX_W_SAFETY
        lines = max(1, math.ceil(
            safe * self._text_width_pt(node.label, font_pt) / inner_w))
        if node.sublabel:
            lines += max(1, math.ceil(
                safe * self._text_width_pt(node.sublabel, font_pt) / inner_w))
        return lines * line_h <= inner_h

    def render_flow(self, slide: PptxSlide, flow: Flow, left: int, top: int,
                    width: int, height: int) -> None:
        """Flow ブロックを矩形領域へ描画する（flow.plan_flow の座標プランを使用）．

        図中の文字サイズは本文プレースホルダの標準サイズに揃える．
        """
        plan = plan_flow(flow, left, top, width, height)
        # box が標準サイズで収まらなければ，全 box 一律で下位レベルへ切り替える．
        boxes = plan.boxes
        if boxes:
            bsz = self._fit_font(
                lambda sz: all(
                    self._box_fits(b.node, b.rect.width, b.rect.height, sz)
                    for b in boxes))
        else:
            bsz = self._body_font_size()
        for bi, box in enumerate(boxes):
            r = box.rect
            tc = self._theme_color(box.node.color) or \
                self._box_palette[bi % len(self._box_palette)]
            self.box(slide, r.left, r.top, r.width, r.height, box.node.label,
                     tc, sub=box.node.sublabel or None, fsize=bsz, ssize=bsz)
        for ell in plan.ellipses:
            # 省略記号スロットは固定幅／固定高（flow.py）なので上下中央に置き，
            # 隣接する矢印との重なりを避ける．
            r = ell.rect
            self.note(slide, r.left, r.top, r.width, r.height, ell.text, bsz,
                      tc=self.T2, bold=True, align=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)
        # ノード間のすき間に塗りつぶしのブロック矢印を置く（太さは box 高に比例）．
        box_h_emu = boxes[0].rect.height if boxes else Inches(1.0)
        thick = int(box_h_emu * 0.34)
        for arrow in plan.arrows:
            self.block_arrow(slide, arrow.x1, arrow.y1, arrow.x2, arrow.y2,
                             thick)
        # 隣り合わないノードを結ぶ線は**細い矢印**で引く（Issue #109）．
        # 上の塗り矢印はすき間を埋めるための形なので、離れた 2 点を結ぶと
        # box に食い込む．
        for ln in plan.lines:
            self.line(slide, ln.x1, ln.y1, ln.x2, ln.y2,
                      width_pt=1.5, dashed=ln.dashed, arrow=True)
        for lab in plan.labels:
            # 矢印ラベルは**折り返さない**．短い語なので、折れると単語の途中で
            # 割れて読めなくなる（"NAMEPREP" が "NAM / EPRE / P" になっていた）．
            r = lab.rect
            self.note(slide, r.left, r.top, r.width, r.height, lab.text, bsz,
                      tc=self.T2, bold=True, align=PP_ALIGN.CENTER, wrap=False)
        for cap in plan.captions:
            # 図に付くのは caption だけ（note_top / note_bottom は地の文なので
            # 本文プレースホルダ側で描く——plan にも入っていない）．
            r = cap.rect
            self.note(slide, r.left, r.top, r.width, r.height, cap.text, bsz,
                      tc=self.T2, align=PP_ALIGN.CENTER)

    def _table_col_widths(self, ncols: int, width: int,
                          col_ratios: list[float] | None) -> list[int]:
        """表の列幅（EMU）リストを返す（均等 or 比率指定）．"""
        if col_ratios and len(col_ratios) == ncols and sum(col_ratios) > 0:
            tot = float(sum(col_ratios))
            return [int(width * r / tot) for r in col_ratios]
        cw = int(width / ncols)
        return [cw] * ncols

    def _table_height_emu(self, data: list[list[str]], col_w: list[int],
                          font_pt: float) -> int:
        """指定フォントサイズでの表の概算総高（EMU）を見積もる（折り返し考慮）．

        実際の PowerPoint レンダリングは行間・最小行高などで見積りより伸びがち
        なため，安全係数を掛けて保守的（やや大きめ）に見積もる．
        """
        line_h = font_pt * 1.32          # 行間込みの行高
        cell_pad_pt = 6                  # セル上下マージン＋最小余白（約）
        side_pad_pt = 18                 # セル左右マージン合計（約．Pt(10)+Pt(6)＋余裕）
        safety = 1.15                    # 折り返し・最小行高ぶんの安全係数
        total_pt = 0.0
        for row in data:
            row_h = line_h + cell_pad_pt
            for ci, cw in enumerate(col_w):
                text = row[ci] if ci < len(row) else ""
                inner_pt = max(1.0, cw / 12700.0 - side_pad_pt)
                lines = max(1, math.ceil(self._text_width_pt(text, font_pt) / inner_pt))
                row_h = max(row_h, lines * line_h + cell_pad_pt)
            total_pt += row_h
        return int(total_pt * safety * 12700)

    def render_table(self, slide: PptxSlide, table: Table, left: int, top: int,
                     width: int, height: int,
                     col_ratios: list[float] | None = None,
                     overflow: bool = False,
                     has_prose_after: bool = False) -> GraphicFrame | None:
        """Table ブロックを座標指定で 1 つ描画する（ヘッダ行をアクセント色で着色）．

        ``参照スクリプト`` の表描画を移植・一般化したもの．列幅は既定で均等，
        ``col_ratios`` を与えると比率配分する．配色はテーマ任せ（ヘッダのみ
        アクセント色 A2＋背景色 BG の文字）．

        overflow=True（@overflow）の場合はフォント縮小（_fit_font）を行わず，
        本文標準（lvl1）サイズのまま必要な高さで描画する．上端は帯上端に固定し，
        はみ出しは下（結論文・罫線側）のみ（画像の overflow と同じ規約）．
        帯に収まる表は通常配置と同じで，はみ出しは発生しない．
        """
        nrows = len(table.rows) + (1 if table.header else 0)
        ncols = max(
            len(table.header) if table.header else 0,
            max((len(r) for r in table.rows), default=0),
        )
        if nrows == 0 or ncols == 0:
            return None

        col_w = self._table_col_widths(ncols, width, col_ratios)
        data = ([table.header] if table.header else []) + list(table.rows)

        if overflow:
            # 縮小せず本文標準サイズを維持し，収まらない分は下へはみ出す．
            fsize = self._body_font_size()
            est_h = self._table_height_emu(data, col_w, fsize)
            if est_h > height:
                if has_prose_after:
                    sys.stderr.write(
                        "md2pptx: warning: overflowing table extends into the "
                        "concluding text below its band\n")
                if top + est_h > self.SH:
                    sys.stderr.write(
                        "md2pptx: warning: overflowing table extends beyond "
                        "the slide bottom edge\n")
            # 帯に収まるなら従来どおり帯高で描く（overflow は no-op）．
            height = max(height, est_h)
        else:
            # フォントは本文標準（lvl1）を基本に，収まらなければ下位レベルへ切り替える．
            fsize = self._fit_font(
                lambda sz: self._table_height_emu(data, col_w, sz) <= height)
            if self._table_height_emu(data, col_w, fsize) > height:
                # 最小レベルまで縮小しても収まらない見積もり．PowerPoint は行を
                # 最小行高以上へ自動拡張するため，黙って帯を超過しうる（従来は
                # 無警告）．気づけるよう警告し，@overflow への誘導も添える．
                sys.stderr.write(
                    "md2pptx: warning: table does not fit its band even at the "
                    "smallest body font size and may overlap following content "
                    "(consider '<!-- @overflow: true -->')\n")

        gf = slide.shapes.add_table(
            nrows, ncols, Emu(left), Emu(top), Emu(width), Emu(height))
        tbl = gf.table

        for ci, cw in enumerate(col_w):
            tbl.columns[ci].width = cw

        for ri, row in enumerate(data):
            is_header = bool(table.header) and ri == 0
            for ci in range(ncols):
                cell = tbl.cell(ri, ci)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                # 列の水平寄せ（区切り行のコロン由来）．未指定は左寄せ扱い．
                al = table.aligns[ci] if ci < len(table.aligns) else "left"
                cell.margin_left = Pt(10)
                # 右寄せ列は既定 margin_right(6pt)だと数字が右壁へ貼りつくため広げる．
                cell.margin_right = Pt(12) if al == "right" else Pt(6)
                cell.margin_top = Pt(2)
                cell.margin_bottom = Pt(2)
                pa = cell.text_frame.paragraphs[0]
                pa.text = row[ci] if ci < len(row) else ""
                if al != "left":
                    pa.alignment = _TABLE_ALIGN[al]
                # セルごとの背景色（§5.4）．**ヘッダの既定より優先する**——
                # 書いたものがそのまま出るほうが説明しやすい．
                fill_name: str | None = None
                if is_header:
                    if ci < len(table.header_fills):
                        fill_name = table.header_fills[ci]
                else:
                    bi = ri - (1 if table.header else 0)
                    if bi < len(table.fills) and ci < len(table.fills[bi]):
                        fill_name = table.fills[bi][ci]
                if fill_name is None and is_header:
                    fill_name = "accent2"       # ヘッダの既定
                fill_rgb: str | None = None
                if fill_name:
                    cell.fill.solid()
                    kind, value = parse_color(fill_name)
                    if kind == "theme":
                        cell.fill.fore_color.theme_color = self._theme_map[value]
                        fill_rgb = self._theme_rgb(value)
                    else:
                        cell.fill.fore_color.rgb = RGBColor.from_string(value)
                        fill_rgb = value
                for run in pa.runs:
                    run.font.size = Pt(fsize)
                    if is_header:
                        run.font.bold = True
                    # 塗ったセルは文字色も塗りに合わせる（Issue #148）．濃い塗りに
                    # 黒文字を載せると読めない——見出し行だけの特別扱いだったものを
                    # 全セルへ広げた．明るい塗りでは触らない（テーマの本文色のまま）．
                    if fill_name and is_dark(fill_rgb):
                        run.font.color.theme_color = self.BG
        return gf

    def _resolve_image_path(self, src: str) -> str:
        """画像パスを base_dir 基準で解決し，存在しなければ fail fast する．"""
        path = resolve_image_path(src, self.base_dir)
        if not os.path.isfile(path):
            raise FileNotFoundError(f"image not found: {src}")
        return path

    @staticmethod
    def _crop_fractions(crop: Crop | None, W: int, H: int) -> tuple[
            float, float, float, float, float, float]:
        """Crop（残す矩形）を PowerPoint のクロップ割合と可視画素サイズへ換算する．

        戻り値 (cl, ct, cr, cb, vis_w_px, vis_h_px)．cl 等は各辺で削る割合（0..1）．
        """
        if crop is None:
            return 0.0, 0.0, 0.0, 0.0, float(W), float(H)
        if crop.unit == "px":
            x, y, w, h = crop.x, crop.y, crop.w, crop.h
        else:  # ソース画像サイズに対する割合
            x, y = crop.x / 100.0 * W, crop.y / 100.0 * H
            w, h = crop.w / 100.0 * W, crop.h / 100.0 * H
        # 各辺で削る割合（0..1）に正規化して検証する．単位（px / %）に依らず
        # 分数で評価できるうえ，許容誤差も相対値（eps）で一貫して扱える．
        cl, ct = x / W, y / H
        cr, cb = (W - (x + w)) / W, (H - (y + h)) / H
        eps = 1e-6  # 割合換算の丸め誤差の吸収（絶対 px ではなく相対量で判定）
        if (w <= 0 or h <= 0 or cl < -eps or ct < -eps
                or cr < -eps or cb < -eps):
            raise ValueError(
                f"crop rectangle out of bounds for {W}x{H}px source: "
                f"keep x={x:g},y={y:g},w={w:g},h={h:g}")
        clamp = lambda v: min(1.0, max(0.0, v))  # 誤差ぶんを [0,1] に丸め込む
        return clamp(cl), clamp(ct), clamp(cr), clamp(cb), w, h

    @staticmethod
    def _resolve_len(length: Length | None, base_emu: float) -> float | None:
        """Length を EMU（float）へ解決する．割合は base_emu 比，絶対はそのまま．None は None．"""
        if length is None:
            return None
        if length.unit == "percent":
            return length.value / 100.0 * base_emu
        return float(length.value)

    def render_seq(self, slide: PptxSlide, seq: Seq, left: int, top: int,
                   width: int, height: int) -> None:
        """シーケンス図（ラダー図）を矩形領域へ描く（Issue #110）．

        座標は ``plan_seq`` が EMU で出し終えているので、ここは図形を置くだけ．
        ライフラインは**矢尻の無い縦線**、メッセージは**矢尻付きの横線**で、
        描き分けは ``PlacedLine`` と ``PlacedArrow`` の型がそのまま指示になる．
        """
        plan = plan_seq(seq, left, top, width, height)
        if not plan.heads:
            return
        bsz = self._body_font_size()
        for i, head in enumerate(plan.heads):
            # 頭は**塗り箱ではなく折り返さない文字**にする．box() は枠の高さを
            # 固定して文字を折るので、「クライアント」のような長い名前が
            # 縦に潰れて読めなくなる（実 PowerPoint で確認）．
            r = head.rect
            self.note(slide, r.left, r.top, r.width, r.height, head.text, bsz,
                      tc=self.T2, bold=True, align=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        for ln in plan.lines:
            # ライフラインは細く目立たせない——主役は矢印のほう．
            self.line(slide, ln.x1, ln.y1, ln.x2, ln.y2,
                      width_pt=1.0, dashed=ln.dashed)
        for ar in plan.arrows:
            self.line(slide, ar.x1, ar.y1, ar.x2, ar.y2,
                      width_pt=1.75, arrow=True)
        for lab in plan.labels:
            r = lab.rect
            self.note(slide, r.left, r.top, r.width, r.height, lab.text, bsz,
                      tc=self.T2, bold=True, align=PP_ALIGN.CENTER, wrap=False)
        for nt in plan.notes:
            r = nt.rect
            self.note(slide, r.left, r.top, r.width, r.height, nt.text, bsz,
                      tc=self.TX, align=PP_ALIGN.LEFT, wrap=False)
        for cap in plan.captions:
            r = cap.rect
            self.note(slide, r.left, r.top, r.width, r.height, cap.text, bsz,
                      tc=self.TX, align=PP_ALIGN.CENTER)

    def render_image(self, slide: PptxSlide, img: Image, left: int, top: int,
                     width: int, seg_h: int, overflow: bool | None = None,
                     has_prose_after: bool = False) -> Picture:
        """Image ブロックをセグメント矩形 (left, top, width, seg_h) 内に配置する．

        ソース画像のピクセル寸法を読み，crop（残す矩形）を PowerPoint のクロップ割合へ
        換算．width/height はアスペクト維持で解決し（両指定かつ fit=fill のときのみ歪ませ），
        align と縦中央でセグメント内へ収める．caption があれば画像下に描画する．
        overflow=True の場合は最終クランプを行わず，明示サイズのまま下方向への
        はみ出しを許可する（上端はセグメント上端まで．タイトル・導入文に重ねない）．
        overflow はブロック指定とスライド @overflow を解決した実効値を受ける
        （None なら Image.overflow のみ参照＝直接呼び出し時の後方互換）．
        """
        if overflow is None:
            overflow = bool(img.overflow)
        path = self._resolve_image_path(img.src)
        W, H = _read_image_size(path)                   # ソースのピクセル寸法
        cl, ct, cr, cb, vis_w, vis_h = self._crop_fractions(img.crop, W, H)
        aspect = (vis_w / vis_h) if vis_h else 1.0      # クロップ後の可視領域の比

        # キャプション用の高さを確保（1 行分）．
        cap_h = int(Pt(self._body_font_size()) * 1.4) if img.caption else 0
        avail_w = float(width)
        avail_h = float(max(1, seg_h - cap_h))

        # width/height を EMU へ解決（未指定は None）．
        w = self._resolve_len(img.width, avail_w)
        h = self._resolve_len(img.height, avail_h)
        # 入れ子にしているのは，「片方が None なら他方は None でない」が
        # ``w is None and h is None`` の否定からは出てこないため——読む側も
        # 型チェッカも，条件の形になっていない不変条件は追えない（Issue #35）．
        if w is None:
            if h is None:                               # 両省略：領域に内接
                w, h = self._fit_within(avail_w, avail_h, aspect)
            else:                                       # 高さのみ：幅は比で
                w = h * aspect
        elif h is None:                                 # 幅のみ：高さは比で
            h = w / aspect
        elif img.fit != "fill":                         # 両指定・contain：比維持で内接
            w, h = self._fit_within(w, h, aspect)
        # 極端な指定（0% 等）でも非正にならないよう下限を張る（ゼロ除算・負サイズ回避）．
        w, h = max(w, 1.0), max(h, 1.0)
        # セグメントを超えないよう最終クランプ（比維持）．overflow 指定時は
        # クランプせず，明示サイズのまま帯からのはみ出しを許可する．
        if not overflow:
            w, h = self._fit_within(min(w, avail_w), min(h, avail_h), w / h)

        # 水平寄せ（align）と縦中央でセグメント内へ配置．
        x: float
        if img.align == "left":
            x = left
        elif img.align == "right":
            x = left + (avail_w - w)
        else:
            x = left + (avail_w - w) / 2.0
        y = top + (avail_h - h) / 2.0
        if overflow:
            # はみ出しは下（結論文・罫線側）のみ．上端はセグメント上端で止め，
            # タイトル・導入文には重ねない（top も y も同じ EMU 数値）．
            y = max(y, top)

        pic = slide.shapes.add_picture(
            path, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
        if img.crop is not None:
            pic.crop_left, pic.crop_top = cl, ct
            pic.crop_right, pic.crop_bottom = cr, cb

        bottom = int(y + h)
        if img.caption:
            # 画像直下に置く．h ≤ avail_h なので通常 y+h ≤ top+avail_h だが，丸め等で
            # セグメント外へ出ないよう cap 上端を [.., top+seg_h-cap_h] にクランプする．
            # overflow 時は画像に追従してさらに下がる（スライド外に出うる）．
            cap_top = int(y + h)
            if not overflow:
                cap_top = min(cap_top, int(top + seg_h - cap_h))
            self._draw_caption(slide, img.caption, left, cap_top, width, cap_h)
            bottom = cap_top + cap_h
        if overflow and bottom > top + seg_h and has_prose_after:
            sys.stderr.write(
                f"md2pptx: warning: overflowing image/caption ({img.src}) "
                "extends into the concluding text below its band\n"
            )
        if overflow and bottom > self.SH:
            sys.stderr.write(
                f"md2pptx: warning: overflowing image/caption ({img.src}) "
                "extends beyond the slide bottom edge\n"
            )
        return pic

    @staticmethod
    def _fit_within(box_w: float, box_h: float,
                    aspect: float) -> tuple[float, float]:
        """アスペクト比 aspect の矩形を (box_w, box_h) に内接させた (w, h) を返す．

        box_w / box_h / aspect が非正のときは安全側（最低 1 EMU）に倒し，ゼロ除算や
        負サイズを避ける（極端に狭いセグメントでも描画を止めない）．
        """
        if aspect <= 0 or box_w <= 0 or box_h <= 0:
            return max(box_w, 1.0), max(box_h, 1.0)
        if box_w / aspect <= box_h:     # 幅が制約：幅いっぱい
            return box_w, box_w / aspect
        return box_h * aspect, box_h    # 高さが制約：高さいっぱい

    def _draw_caption(self, slide: PptxSlide, text: str, left: int, top: int,
                      width: int, height: int) -> None:
        """図下キャプションを中央寄せの小さめ本文サイズで描く．"""
        tb = slide.shapes.add_textbox(
            Emu(left), Emu(top), Emu(width), Emu(max(height, Pt(12))))
        tf = tb.text_frame
        tf.word_wrap = True
        # キャプションは短文前提（1 行分の高さを確保）．枠を内容で伸ばさない
        # （長文で下方向へはみ出さないよう auto_size を無効化）．必要なら折り返す．
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = text
        # 本文標準より 1 段小さめ（テーマ既定サイズ体系の中で縮小）．
        levels = self._body_font_levels()
        size = levels[1] if len(levels) > 1 else levels[0]
        for r in p.runs:
            r.font.size = Pt(size)

    # ------------------------------------------------------- content slide
    def render_slide(self, slide: Slide, slide_number: bool = True,
                     default_autofit: bool = True) -> PptxSlide:
        """コンテンツスライドを 1 枚追加して返す．

        ``slide.blocks`` を出現順に処理する．表を含まないスライドは Phase 1 と
        同じく本文プレースホルダへ Line を流し込み，表を含むスライドは座標スタック
        配置（テキスト→表→テキスト …）で描画する．Flow は Phase 3 でスキップ．
        """
        directives = slide.directives or {}
        # directives の値はキーごとに型が違う（int/str/bool）．@layout に数値以外が
        # 来たら既定レイアウトへ倒す（parser 側でも int 化を試みている）．
        raw_layout = directives.get("layout", slide.layout)
        layout_idx = raw_layout if isinstance(raw_layout, int) else None
        try:
            layout = self.L1 if layout_idx is None else self.layouts[layout_idx]
        except IndexError:
            layout = self.L1

        s = self.prs.slides.add_slide(layout)
        if slide.title is not None and s.shapes.title is not None:
            tf = s.shapes.title.text_frame
            tf.text = slide.title
            # 基点は lvl1 固定．本文行と違いタイトルには**ネストの概念が無く**，
            # p.level を設定しないので常に lvl1 で描かれる（本文行が
            # levels[blk.level] を引くのはインデントで段が変わるから）．
            # サイズは run の rPr へ入る．タイトルに段落既定（defRPr）を書く経路を
            # 足すなら，**run 側が優先される**ことに注意——ここで書いた段数が
            # 段落側の指定を上書きする（本文行は defRPr と rPr を役割で分けている）．
            self._apply_segment_deltas(tf.paragraphs[0], slide.title_deltas,
                                       self._frame_font_levels(tf)[0])

        # @widths によるプレースホルダ幅の上書きは，本文描画より
        # 前に済ませる（以降の _effective_geom / _content_rect が上書き後を参照）．
        self._apply_placeholder_widths(s, directives,
                                       is_columns=bool(slide.columns))

        # スライド既定の採番色（@autonum-color）．Line.num_color が優先．
        raw_num_color = directives.get("autonum_color")
        default_num_color = (raw_num_color if isinstance(raw_num_color, str)
                             else None)
        # スライド既定の相対サイズ段数（@body-size）．Line.size_delta が優先．
        default_size_delta = self._body_size_delta(directives)
        scale = self._autofit_scale(directives)
        # スライド単位の overflow（@overflow）．表・画像に共通で効き，
        # 画像はブロックの overflow: 明示があればそちらが優先する．
        slide_overflow = bool(directives.get("overflow", False))
        blocks = slide.blocks or []

        if slide.columns:
            self._render_columns(s, slide.columns, default_num_color, scale,
                                 default_autofit, default_size_delta,
                                 self._col_ratios(directives), slide_overflow)
            if directives.get("col_arrow"):
                self.draw_column_arrow(s, len(slide.columns))
        elif any(is_object_block(b) for b in blocks):
            self._render_stacked(s, blocks, default_num_color, scale, default_autofit,
                                 self._col_ratios(directives), default_size_delta,
                                 slide_overflow)
        else:
            line_blocks = [b for b in blocks if isinstance(b, Line)]
            body = self._body_placeholder(s)
            if body is None:
                self._warn_no_body(line_blocks)
            elif line_blocks:
                tf = body.text_frame
                self._fill_lines(tf, line_blocks, default_num_color,
                                 default_size_delta)
                eff = self._autofit_for(body, line_blocks, scale,
                                        default_autofit, default_size_delta)
                self.draw_line_boxes(s, body, line_blocks, default_size_delta,
                                     shrink=eff)

        # 表紙レイアウト（テーマの「タイトル スライド」）には番号を付けない．
        # そのレイアウトを選ぶこと自体が「これは表紙」の宣言なので，
        # 番号の有無をそこに紐づける（Issue #82）．front matter 由来の表紙を
        # 描く render_title_slide も番号を付けず，本文記法で書いた表紙が
        # 同じ扱いになる．slide_number: false は従来どおり全体に効く．
        if slide_number and layout_idx != TITLE_LAYOUT:
            self.add_slide_number(s)
        self._set_notes(s, slide.notes)
        return s

    def _body_size_delta(self, directives: dict[str, Any]) -> int | None:
        """@body-size ディレクティブをスライド既定の相対サイズ段数へ解釈する．

        未指定・非整数値はいずれも None（＝スライド既定なし）を返す．None は
        「未指定」を明示する番兵で，size_delta=None の行はテーマ既定のままになる
        （0＝明示的に 0 段，とは区別する）．

        `@body-size: 0`（0 段）は「スライド既定なし」と同義として None を返す．
        スライド全体に対する 0 段は変化なし＝既定なしと区別する意味がないため．
        （行トークン `{0}` の「テーマ既定へ明示的に戻す」意味は別物で，スライド既定
        が非 0 のとき個別行を素のテーマ既定へ戻す用途に残る．Line.size_delta=0 が
        担い，こちらには波及しない．）

        parser 経由なら body_size は _INT_DIRECTIVES で int 化済みのため int()
        は素通りする．try/except は parser を介さず directives を直接組み立てる
        ケース（テスト・他コードからの呼び出し）に対する防御で，不正値で落とさず
        「スライド既定なし」に倒す．
        """
        val = directives.get("body_size")
        if val is None:
            return None
        try:
            iv = int(val)
        except (TypeError, ValueError):
            sys.stderr.write(
                f"md2pptx: warning: ignoring non-integer @body-size value "
                f"{val!r}\n"
            )
            return None
        return iv if iv != 0 else None

    def _render_columns(self, slide: PptxSlide, columns: list[list[Block]],
                        default_num_color: str | None, scale: float | None,
                        default_autofit: bool,
                        default_size_delta: int | None = None,
                        col_ratios: list[float] | None = None,
                        slide_overflow: bool = False) -> None:
        """多カラム（「2つのコンテンツ」）：各カラムを idx 1, 2 … へ流す．

        columns[i] を プレースホルダ idx=i+1 へ描画する（idx 0 はタイトル）．
        Line（箇条書き・採番・no_bullet）はプレースホルダへ流し込み，Table/Image/
        Flow を含むカラムはそのプレースホルダ矩形へ座標スタック配置する（地の文と
        混在する場合は _render_stacked_into が空行帯で棲み分ける）．
        """
        for ci, col_blocks in enumerate(columns):
            ph = self._find_placeholder(slide, ci + 1)
            if ph is None:
                continue  # レイアウトに該当プレースホルダが無ければスキップ
            if any(is_object_block(b) for b in col_blocks):
                # カラム矩形へ表・図をスタック配置．継承ジオメトリはレイアウトで補う．
                # 通常 layout 3 は idx1/idx2 のジオメトリを持つため，解決失敗はテーマ
                # 異常時のみ．その場合は本文領域へフォールバックする（表が消えるより，
                # 見えて重なる方が原因に気づきやすい）が，複数カラムが重なりうるので警告
                # を出す．
                gl, gt, gw, gh = self._effective_geom(ph, slide)
                if gl is None or gt is None or gw is None or gh is None:
                    sys.stderr.write(
                        f"md2pptx: warning: could not resolve geometry for "
                        f"column {ci}; falling back to the body area "
                        f"(columns may overlap)\n"
                    )
                    left, top, width, height = self._content_rect(slide)
                else:
                    left, top, width, height = gl, gt, gw, gh
                # @table-widths はスライド共通で全カラムの表に適用する．列数が比率の
                # 要素数と一致しない表は _table_col_widths が等幅へフォールバックする．
                self._render_stacked_into(slide, col_blocks, ph, left, top,
                                          width, height, default_num_color, scale,
                                          default_autofit, col_ratios,
                                          default_size_delta, slide_overflow)
                continue
            # Line のみのカラムはプレースホルダへ直接流し込む．_render_stacked_into は
            # objects（表・図）が空だと何も描画せず return する設計なので，ここを通すと
            # 箇条書きが消える．そのため表・図を含むカラムとは意図的に経路を分ける．
            lines = [b for b in col_blocks if isinstance(b, Line)]
            if lines:
                tf = ph.text_frame
                self._fill_lines(tf, lines, default_num_color, default_size_delta)
                eff = self._autofit_for(ph, lines, scale, default_autofit,
                                        default_size_delta)
                self.draw_line_boxes(slide, ph, lines, default_size_delta,
                                     shrink=eff)

    # ----------------------------------------------------- 描画ユーティリティ
    def _append_lines(self, tf: TextFrame, line_blocks: list[Line], first: bool,
                      default_num_color: str | None,
                      default_size_delta: int | None = None,
                      counters: dict[tuple[int, str], int] | None = None) -> bool:
        """Line 列を text_frame に段落として追記する（採番／no_bullet を適用）．

        first=True なら最初の 1 行は既存の paragraphs[0] を使う．残りの行を
        追記しても良いよう，処理後の first 状態を返す．

        相対サイズは行の size_delta を優先し，None の行はスライド既定
        （default_size_delta，@body-size 由来）を継承する．基点は tf の枠が
        実際に持つ既定サイズで，枠ごとに 1 度だけ解決する（Issue #83）．

        counters は (level, 形式) ごとの次の番号（Issue #107）．**全ての採番段落に
        番号を明示して書く**——PowerPoint は startAt の付いた段落の**次から数え直す**
        ので，先頭にだけ書くと 8. の次が 1. に戻る．番号はここで数えて渡す．
        原稿の番号は**リストの先頭の行だけ**を種にし，以降は 1 ずつ増やす
        （CommonMark と同じ規則．"1. 1. 1." と書けば 1・2・3 になる）．
        **同じ枠へ 2 回以上追記する呼び出し側は，この辞書を自分で作って渡す**．
        図表スライドでは地の文が帯の上下に分かれて同じ枠へ 2 回追記されるので，
        渡さないと結論文側の番号が 1 に戻る．1 回で流し込む経路（``_fill_lines``）は
        渡さなくてよく，そのとき既定の ``None`` が枠ごとに新しい辞書になる．
        """
        if counters is None:
            counters = {}
        levels = self._frame_font_levels(tf)
        for blk in line_blocks:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = blk.level
            self._write_spans(p, blk)
            if blk.kind == "autonum":
                fmt = blk.num_style or "arabicPeriod"
                color = blk.num_color or default_num_color
                key = (blk.level, fmt)
                num = counters.get(key)
                if num is None:                     # そのリストの先頭の行
                    num = blk.num_start if blk.num_start is not None else 1
                counters[key] = num + 1
                self.set_autonum(p, fmt, color=color, start=num)
            elif blk.kind == "plain":
                self.no_bullet(p)
            elif blk.kind == "code":
                # 行頭記号を消して等幅にする（§5.12）．書いた桁が揃わないと
                # コードとして読めないので，ここだけはテーマ既定に任せない．
                self.no_bullet(p)
                for run in p.runs:
                    run.font.name = self._mono_font
            # kind == "bullet" はテーマ既定のまま
            delta = blk.size_delta if blk.size_delta is not None else default_size_delta
            # 行トークン {0} は「スライド既定を無効化してテーマ既定へ戻す」意味だが，
            # そもそもスライド既定が無い（default_size_delta is None）なら戻す対象が
            # なく無意味な no-op なので適用しない（テーマの段落サイズに触れない）．
            if delta == 0 and default_size_delta is None:
                delta = None
            self._apply_size_delta(p, blk.level, delta, levels)
            # セグメントの段数は行の段数と**同じ基点**（その level のテーマ既定）から
            # 数える．行が {+1} でもセグメントの {-2} は同じ大きさになる——
            # 「テーマ既定からの相対段数」という記法の意味を段で変えないため．
            self._apply_segment_deltas(
                p, blk.seg_deltas, levels[min(blk.level, len(levels) - 1)],
                [s.segment for s in blk.spans] if blk.spans else None)
        return first

    def draw_line_boxes(self, slide: PptxSlide, ph: SlidePlaceholder,
                        line_blocks: list[Line],
                        default_size_delta: int | None = None,
                        preceding: list[Line] | None = None,
                        blank_paras: int = 0,
                        shrink: float | None = None) -> None:
        """``{box}`` の付いた段落を枠で囲む（Issue #133）．

        枠は**段落に重ねて描く**——PowerPoint が実際にどこへ行を置いたかは
        pptx を書く側からは分からないので、``_render_stacked_into`` の帯計算と
        **同じ見積もり**（枠の上端から「その行のサイズ × 1.32」を積む）で位置を出す．
        折り返しの行数も ``_text_width_pt`` で見積もり、2 行になる項目を 1 つの枠で囲む．

        ``preceding`` は同じ枠へ先に書かれた行（表・図スライドの導入文）、
        ``blank_paras`` はその後ろに挟まる空段落の数（帯を空ける行）．
        **先行ぶんも行ごとのサイズで積む**——一律に本文標準サイズで数えると、
        導入文に ``{+1}`` が付いているだけで結論文側の枠がずれる．
        空段落だけは標準サイズで数える（帯の計算がそう作っているため）．

        ``shrink`` は枠に効いている縮小率（％．無ければ None）．**これを渡さないと
        縮んだ枠で位置がずれる**——字だけが小さくなり、枠は元の大きさのまま
        取り残される（Issue #154）．

        **見積もりなので完全ではない**．テーマがレベルごとに行間を変えている枠では
        ずれうる（SYNTAX.md に明記）．ずれても文字は動かない——枠だけが少し外れる．
        """
        k = (shrink / 100.0) if shrink is not None else 1.0
        if not any(ln.boxed for ln in line_blocks):
            return
        tf = ph.text_frame
        levels = self._frame_font_levels(tf)
        indents = self._body_indents()

        def size_of(level: int) -> float:
            return levels[min(level, len(levels) - 1)] if levels else 18.0

        def indent_of(level: int) -> int:
            return indents[min(level, len(indents) - 1)] if indents else 0

        # python-pptx の margin_* は未設定でも既定値（EMU）を返すので、
        # そのまま使える（None にはならない）．
        pad_l = tf.margin_left
        pad_r = tf.margin_right
        avail_full = max(1, ph.width - pad_l - pad_r)

        def measure(ln: Line) -> tuple[float, int, int]:
            """(実効サイズ pt, 折り返し行数, 左インデント EMU) を返す．"""
            d = ln.size_delta if ln.size_delta is not None else default_size_delta
            sz = self._size_from_delta(size_of(ln.level), d) * k
            ind = indent_of(ln.level)
            avail_pt = max(1.0, (avail_full - ind) / 12700.0)
            text = (ln.text or "").replace("\v", " ")
            return sz, self._wrapped_lines(text, sz, avail_pt), ind

        y = ph.top + tf.margin_top
        # preceding は**位置を数えるためだけ**に使う．そこに枠が付いていても
        # ここでは描かない——呼び出し元が同じ行列で 1 度描いている．
        for ln in preceding or []:
            sz, n, _ = measure(ln)
            y += self._para_height(ln.level, sz, n)
        # 空段落は標準サイズ（帯の計算がそう作っている）．
        y += blank_paras * self._para_height(0, self._body_font_size() * k)

        for ln in line_blocks:
            sz, wrapped, ind = measure(ln)
            para_h = self._para_height(ln.level, sz, wrapped)
            if ln.boxed:
                avail = max(1, avail_full - ind)
                text = (ln.text or "").replace("\v", " ")
                # 枠は文字幅に合わせる（1 行に収まる短い項目まで枠が伸びると、
                # 「ここだけ囲んでいる」ことが伝わらない）．左右に少し余白を取る．
                gap = int(Pt(sz) * 0.35)
                # 左へ食み出しても**スライドの外へは出さない**．
                bl = max(0, ph.left + pad_l + ind - gap)
                w = min(int(self._text_width_pt(text, sz) * 12700) + 2 * gap,
                        ph.left + ph.width - bl)
                # 枠は**アキを含めない**——アキは段落の上に空く隙間で、字の入る
                # ところではない．囲むのは字のほうだけ（Issue #150）．
                # そこから ``_BOX_LIFT`` だけ持ち上げる（Issue #160）．
                line_h = self._line_height(sz)
                box_top = (y + self._space_before(ln.level, sz)
                           - int(line_h * self._BOX_LIFT))
                self.line_box(slide, bl, max(0, box_top),
                              max(w, int(avail * 0.2)),
                              wrapped * line_h, ln.box_color)
            y += para_h

    def line_box(self, slide: PptxSlide, left: int, top: int, w: int, h: int,
                 color: str | None = None) -> Shape:
        """段落を囲む枠（塗りつぶし無しの角丸四角）を描く．

        塗りつぶさないのは、下に文字があるため．色の既定はテーマのアクセント色で、
        ``{box:blue}`` のように指定があればそちらを使う（語彙は行内装飾と共通）．

        ``color`` はパーサが正規化済みだが、``_set_run_color`` と同じくここでも
        ``parse_color`` を通す．冪等（テーマ色名も "#RRGGBB" もそのまま返る）なうえ、
        **色名の語彙を知る場所を 1 つに保てる**．
        """
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(w), Emu(h))
        shp.fill.background()
        if color:
            kind, value = parse_color(color)
            if kind == "theme":
                shp.line.color.theme_color = self._theme_map[value]
            else:
                shp.line.color.rgb = RGBColor.from_string(value)
        else:
            shp.line.color.theme_color = self.A2
        # 元の講義スライドは 6pt．そこまで太らせると枠が主役になるので、
        # 「離れて見ても囲みと分かる」ところで 3pt にしている．
        shp.line.width = Pt(3)
        shp.shadow.inherit = False
        # 図形に文字は入れない（文字はプレースホルダ側にある）．
        shp.text_frame.word_wrap = False
        return shp

    def _fill_lines(self, tf: TextFrame, line_blocks: list[Line],
                    default_num_color: str | None,
                    default_size_delta: int | None = None) -> None:
        """Line 列を text_frame の段落として流し込む（先頭から）．"""
        self._append_lines(tf, line_blocks, True, default_num_color,
                           default_size_delta)

    def _autofit_scale(self, directives: dict[str, Any]) -> float | None:
        """@autofit ディレクティブを縮小率へ解釈する（不正値は警告して None）．

        返すのは**％の数値そのもの**（``@autofit: 90`` → ``90.0``）で、比ではない．
        ``fit_body`` の ``scale`` がそのまま％を取るため——比へ直すのは使う側．

        0 以下は受けない——文字が消えるか裏返るかで、どちらも書き手の意図では
        ありえない．ここで弾いておかないと帯の計算（``shrink``）まで巻き込む．
        """
        autofit = directives.get("autofit")
        if autofit is None:
            return None
        try:
            scale = float(autofit)
        except (TypeError, ValueError):
            sys.stderr.write(
                f"md2pptx: warning: ignoring non-numeric @autofit value "
                f"{autofit!r}\n"
            )
            return None
        if scale <= 0:
            sys.stderr.write(
                f"md2pptx: warning: ignoring non-positive @autofit value "
                f"{autofit!r}\n"
            )
            return None
        return scale

    def _text_height(self, lines: list[Line], levels: list[float],
                     default_size_delta: int | None, avail_w: int,
                     shrink: float = 1.0) -> int:
        """``lines`` を ``avail_w`` の幅へ流したときの総高（EMU）．

        折り返しは ``_text_width_pt`` の概算で数える（`draw_line_boxes` と同じ）．
        """
        total = 0
        for ln in lines:
            d = ln.size_delta if ln.size_delta is not None else default_size_delta
            base = levels[min(ln.level, len(levels) - 1)] if levels else 18.0
            sz = self._size_from_delta(base, d) * shrink
            ind = self._indent_for(ln.level)
            avail_pt = max(1.0, (avail_w - ind) / 12700.0)
            text = (ln.text or "").replace("\v", " ")
            total += self._para_height(
                ln.level, sz, self._wrapped_lines(text, sz, avail_pt))
        return total

    def _fit_scale(self, ph: SlidePlaceholder, lines: list[Line],
                   default_size_delta: int | None) -> float | None:
        """枠に収めるための縮小率（％）．収まるなら None．

        **``normAutofit`` を置くだけでは何も起きない**（Issue #154）．PowerPoint は
        自動調整を開いたときに計算し直さず、保存されている ``fontScale`` で描く。
        だから md2pptx 側で率を出して焼き込む．

        縮めると折り返しが減ってさらに縮むので、数回まわして落ち着かせる．
        """
        tf = ph.text_frame
        # 継承したままのプレースホルダは width / height が None を返す．
        # そこはレイアウト側で補えるが（``_effective_geom``）、スライドを持って
        # いないここでは補えないので、測れなければ縮めない．
        if ph.width is None or ph.height is None or not lines:
            return None
        avail_h = ph.height - tf.margin_top - tf.margin_bottom
        avail_w = ph.width - tf.margin_left - tf.margin_right
        if avail_h <= 0 or avail_w <= 0:
            return None
        levels = self._frame_font_levels(tf)
        need = self._text_height(lines, levels, default_size_delta, avail_w)
        if need <= avail_h:
            return None
        scale = avail_h / need
        for _ in range(3):
            need = self._text_height(lines, levels, default_size_delta,
                                     avail_w, scale)
            if need <= avail_h:
                break
            scale *= avail_h / need
        # 下限は 25%．それ以下まで縮むのは原稿の量がおかしいので、縮めきらずに
        # 止めて（はみ出しは残る）書き手に気づかせる．
        return max(25.0, min(100.0, scale * 100.0))

    def _apply_autofit(self, tf: TextFrame, scale: float | None,
                       default_autofit: bool) -> None:
        """縮小率指定があれば焼き込み，無ければ既定の自動調整を設定する．"""
        if scale is not None:
            self.fit_body(tf, scale=scale)
        elif default_autofit:
            self.fit_body(tf)

    def _autofit_for(self, ph: SlidePlaceholder, lines: list[Line],
                     scale: float | None, default_autofit: bool,
                     default_size_delta: int | None) -> float | None:
        """枠に自動調整を設定し、**実際に効く縮小率**を返す．

        戻り値は ``draw_line_boxes`` へ渡す——縮小率を知らないと ``{box}`` の枠が
        字とずれる（枠だけが元の大きさの位置に取り残される）．
        """
        eff = scale
        if eff is None and default_autofit:
            eff = self._fit_scale(ph, lines, default_size_delta)
        self._apply_autofit(ph.text_frame, eff, default_autofit)
        return eff

    def _col_ratios(self, directives: dict[str, Any]) -> list[float] | None:
        """@table-widths ディレクティブ（"45,55" 等）を表の列幅比リストへ解釈する．"""
        v = directives.get("table_widths")
        if not v:
            return None
        try:
            return [float(x) for x in str(v).replace("，", ",").split(",")]
        except ValueError:
            return None

    def _split_allow_left(self, value: object) -> tuple[str | None, bool]:
        """値末尾の ``!``（左余白の使用許可）を分離して (本体文字列, フラグ) を返す．"""
        v = str(value).strip()
        if v.endswith(("!", "！")):
            return v[:-1].strip(), True
        return v, False

    def _parse_pct_list(self, value: object) -> list[float] | None:
        """百分率リスト（例: "55,45"，"55%,45%"）を float の list へ解釈する．

        全角の区切り（，／％）も受理する．不正値は None を返す．
        """
        try:
            return [float(str(x).strip().rstrip("%％"))
                    for x in str(value).replace("，", ",").split(",")]
        except ValueError:
            return None

    def _override_geom(self, ph: SlidePlaceholder, left: float, top: float,
                       width: float, height: float) -> None:
        """スライド側へ明示ジオメトリを書き，レイアウト継承を上書きする．

        4 値すべて設定するのは，一部のみ明示すると xfrm が不完全になり
        PowerPoint 側の解釈が実装依存になるため（継承値で補って全指定する）．
        """
        ph.left, ph.top = int(left), int(top)
        ph.width, ph.height = int(width), int(height)

    _PH_MARGIN = Inches(0.1)   # プレースホルダ拡幅時にスライド端へ残す余白

    def _apply_placeholder_widths(self, slide: PptxSlide,
                                  directives: dict[str, Any],
                                  is_columns: bool) -> None:
        """@widths をスライド種別（単カラム／多カラム）に応じて適用する．

        いずれも「標準の使用可能幅に対する百分率」で解釈する（詳細は各メソッド）．
        値 1 個は単カラム本文幅（例: "104"），複数はカラムごとの幅（例: "62,40"）．
        拡幅は**左端固定・右余白のみ**が既定（箇条書きの行頭位置がスライド間で
        揃い，遷移時の見た目が安定する）．右余白で収まらない指定はクランプして
        警告する．値の末尾に ``!`` を付けると（例: "108!" / "62,47!"），収まら
        ない分だけ左余白へ逃がすことを許可する．その場合もスライド端は余白
        _PH_MARGIN でクランプし，それでも収まらない指定は警告のうえ比例縮小する．

        スライド種別と値の個数が合わない指定は無視し，警告を出す．
        """
        val = directives.get("widths")
        if val is None:
            return
        if is_columns:
            self._apply_ph_widths(slide, val)
        else:
            self._apply_body_width(slide, val)

    def _apply_ph_widths(self, slide: PptxSlide, val: object) -> None:
        """@widths: "55,45" — 多カラムのプレースホルダ幅を再指定する．

        カラム群の合計スパンからカラム間ギャップを除いた幅を 100% とし，
        各カラム幅を百分率で再指定する（ギャップは維持）．合計が 100 を
        超えると全体が右方向へ広がる（55,50 → 全体が標準の 105%）．
        ジオメトリを解決できない場合は警告して何もしない（従来描画）．
        """
        val, allow_left = self._split_allow_left(val)
        pcts = self._parse_pct_list(val)
        if not pcts or any(p <= 0 for p in pcts):
            sys.stderr.write(
                f"md2pptx: warning: ignoring invalid @widths value {val!r}\n")
            return
        if len(pcts) < 2:
            sys.stderr.write(
                "md2pptx: warning: @widths on a multi-column slide expects "
                f"one value per column, got {val!r}; ignoring\n")
            return
        # md のカラム順＝プレースホルダ idx 順（_render_columns と同じ対応）で集める．
        phs = []
        for i, _pct in enumerate(pcts):
            ph = self._find_placeholder(slide, i + 1)
            if ph is None:
                sys.stderr.write(
                    f"md2pptx: warning: @widths has {len(pcts)} values but "
                    f"column placeholder {i + 1} does not exist; ignoring\n")
                return
            gl, gt, gw, gh = self._effective_geom(ph, slide)
            if gl is None or gt is None or gw is None or gh is None:
                sys.stderr.write(
                    "md2pptx: warning: @widths skipped "
                    "(could not resolve column geometry)\n")
                return
            phs.append((ph, (gl, gt, gw, gh)))
        lefts = [g[0] for _, g in phs]
        rights = [g[0] + g[2] for _, g in phs]
        gaps = [lefts[i + 1] - rights[i] for i in range(len(phs) - 1)]
        if any(g < 0 for g in gaps):
            # 重なったプレースホルダ（負のギャップ）は usable を過大にするため 0 扱い．
            sys.stderr.write(
                "md2pptx: warning: @widths found overlapping column "
                "placeholders; treating the negative gap as 0\n")
            gaps = [max(g, 0) for g in gaps]
        span_l, span_r = lefts[0], rights[-1]
        usable = (span_r - span_l) - sum(gaps)
        widths = [usable * p / 100.0 for p in pcts]
        new_span = sum(widths) + sum(gaps)
        # 既定は左端固定（右余白のみ使用）．"...!" で左余白の使用を許可する．
        max_span = ((self.SW - self._PH_MARGIN - span_l) if not allow_left
                    else (self.SW - 2 * self._PH_MARGIN))
        if new_span > max_span:
            sys.stderr.write(
                "md2pptx: warning: @widths total exceeds the "
                f"{'slide' if allow_left else 'right margin'}; clamping"
                f"{'' if allow_left else ' (append ! to use the left margin)'}\n")
            k = (max_span - sum(gaps)) / float(sum(widths))
            widths = [w * k for w in widths]
            new_span = max_span
        new_left = span_l
        if allow_left:
            overflow = (new_left + new_span) - (self.SW - self._PH_MARGIN)
            if overflow > 0:
                new_left -= overflow
            new_left = max(new_left, self._PH_MARGIN)
        # 幅は百分率から出るので実数．書き込む直前に _override_geom が int 化する．
        x: float = new_left
        for i, ((ph, (_l, t, _w, h)), nw) in enumerate(zip(phs, widths)):
            self._override_geom(ph, x, t, nw, h)
            x += nw + (gaps[i] if i < len(gaps) else 0)

    def _apply_body_width(self, slide: PptxSlide, val: object) -> None:
        """@widths: "105" — 単カラム本文プレースホルダ幅を再指定する．

        継承した本文プレースホルダ幅に対する百分率（% 付き可）．値は 1 個のみ
        （複数値は多カラム用）．ジオメトリを解決できない場合は何もしない（従来描画）．
        """
        val, allow_left = self._split_allow_left(val)
        pcts = self._parse_pct_list(val)
        if pcts is None:
            sys.stderr.write(
                f"md2pptx: warning: ignoring invalid @widths value {val!r}\n")
            return
        if len(pcts) != 1:
            sys.stderr.write(
                "md2pptx: warning: @widths on a single-column slide expects "
                f"exactly 1 value, got {val!r}; ignoring\n")
            return
        pct = pcts[0]
        if pct <= 0:
            sys.stderr.write(
                f"md2pptx: warning: ignoring non-positive @widths value {val!r}\n")
            return
        ph = self._body_placeholder(slide)
        if ph is None:
            return
        left, top, width, height = self._effective_geom(ph, slide)
        if left is None or top is None or width is None or height is None:
            return
        new_w = width * pct / 100.0
        # 既定は左端固定（右余白のみ使用）．"...!" で左余白の使用を許可する．
        max_w = ((self.SW - self._PH_MARGIN - left) if not allow_left
                 else (self.SW - 2 * self._PH_MARGIN))
        if new_w > max_w:
            sys.stderr.write(
                "md2pptx: warning: @widths exceeds the "
                f"{'slide' if allow_left else 'right margin'}; clamping"
                f"{'' if allow_left else ' (append ! to use the left margin)'}\n")
            new_w = max_w
        new_l = left
        if allow_left:
            overflow = (new_l + new_w) - (self.SW - self._PH_MARGIN)
            if overflow > 0:
                new_l -= overflow
            new_l = max(new_l, self._PH_MARGIN)
        self._override_geom(ph, new_l, top, new_w, height)

    def _content_rect(self, slide: PptxSlide) -> tuple[int, int, int, int]:
        """本文領域の矩形 (left, top, width, height) を返す（座標配置の基準）．"""
        ph = self._body_placeholder(slide)
        if ph is not None and None not in (ph.left, ph.top, ph.width, ph.height):
            return (ph.left, ph.top, ph.width, ph.height)
        try:
            for lph in slide.slide_layout.placeholders:  # type: ignore[misc]
                if lph.placeholder_format.idx == 1 and None not in (
                    lph.left, lph.top, lph.width, lph.height
                ):
                    return (lph.left, lph.top, lph.width, lph.height)
        except Exception:
            pass
        # 本文プレースホルダの無いレイアウト（「白紙」「タイトルのみ」）．
        # **タイトルが無ければ、タイトルぶんを空けない**（Issue #138）．
        # 1.7in はタイトルを避けるための値で、タイトルの無いレイアウトでは
        # 根拠が無い——図だけのスライドがそのぶん小さくなる．
        top = Inches(1.7) if self._layout_has_title(slide) else Inches(0.4)
        return (Inches(0.6), top, self.SW - Inches(1.2),
                self.SH - top - Inches(0.6))

    @staticmethod
    def _layout_has_title(slide: PptxSlide) -> bool:
        """このスライドのレイアウトがタイトルの枠を持つか．

        スライド側ではなく**レイアウト**を見る．`###` の見出しはタイトル枠が
        無ければ描かれないので、「書いたかどうか」ではなく「置ける場所があるか」で
        決める．
        """
        try:
            for lph in slide.slide_layout.placeholders:  # type: ignore[misc]
                if lph.placeholder_format.idx == 0:
                    return True
        except Exception:
            return True          # 分からなければ従来どおり空ける（安全側）
        return False

    def _note_to_line(self, text: str) -> Line | None:
        """図（Flow / Seq）の note 文字列を本文プレースホルダ用の Line へ変換する．

        note(top) / note(bottom) は図の一部ではなく**地の文**なので、解釈は
        本文行と同じでなければならない（Issue #129）．行頭マーカーだけを
        自前で見ていた頃は ``[語]{red}`` が生の文字で出ていた——行内装飾は
        本文行が通る ``parse_content_line`` の中で解決される．

        **段落にならない行では None を返す**（``1.`` のように行頭マーカーだけの
        note）．本文では行を作らない書き方なので、ここで空段落を作ると地の文が
        1 行ぶん増え、帯が詰まって図と結論文が近づく．
        """
        return parse_content_line((text or "").strip())

    def _obj_weight(self, obj: ObjectBlock) -> int:
        """オブジェクト（Table / Flow / Seq / Image）の縦方向の重み（高さ配分用）．"""
        if isinstance(obj, Flow):
            return max(4, len(obj.nodes) + 2)
        if isinstance(obj, Seq):
            # ラダー図は**やりとりの本数**で背が決まる（横幅は人数で決まる）．
            return max(4, len(obj.messages) + 2)
        if isinstance(obj, Arrow):
            # 矢印は「流れの向き」を示すだけなので、帯は狭くてよい．
            return 3
        if isinstance(obj, Image):
            # 画像は帯を広めに確保（キャプションぶんを少し足す）．細かな大きさは
            # width/height でセグメント内に調整する．
            return 8 + (1 if obj.caption else 0)
        return max(2, len(obj.rows) + (1 if obj.header else 0))

    # ``` ```arrow ``` の向き → OOXML の図形．型注釈（ir.ArrowDirection）と
    # 別に並べているので、**向きを増やしたらここも足す**（テストが先に落ちる）．
    # 矢じりと軸の比率（OOXML の ``a:avLst``）．**書かないと PowerPoint の既定に
    # 落ち、両端に矢じりのある形（updown / leftright）は箱いっぱいの菱形になる**．
    # 値は元の講義スライドと同じ——軸 50%、矢じりは片側 25%・両側なら 20%．
    _ARROW_ADJ = {"adj1": 50000, "adj2": 25000}
    _ARROW_ADJ_BOTH = {"adj1": 50000, "adj2": 20000}

    _ARROW_SHAPES = {
        "down": MSO_SHAPE.DOWN_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "updown": MSO_SHAPE.UP_DOWN_ARROW,
        "leftright": MSO_SHAPE.LEFT_RIGHT_ARROW,
    }

    def render_arrow(self, slide: PptxSlide, arrow: Arrow,
                     left: int, top: int, width: int, height: int) -> Shape:
        """大きな下向き矢印を、与えられた帯の中央に描く（Issue #134）．

        大きさは**上限を持たせる**．帯の高さに素直に比例させると、地の文が
        少ないスライドで矢印がページの主役になってしまう．

        横向きは長手が横になるので、**長さの基準も帯の幅**に取る．高さから
        測ると、帯が薄いスライドで横向き矢印だけが縮む．
        """
        shape = self._ARROW_SHAPES[arrow.direction]
        horizontal = arrow.direction in ("right", "left", "leftright")
        along = width if horizontal else height     # 矢印が伸びる向きの余地
        across = height if horizontal else width    # それと直交する向きの余地
        # 長手は自分の帯の 8 割．**ここを削ると矢印が向きを失う**——2 つ置いて
        # 帯を分け合うと 1 つあたり 1.8cm ほどしか無く、半分では菱形に見える
        # （Issue #141）．上限は元の講義スライドの大きさに合わせた．
        long_ = min(Inches(1.0), max(Inches(0.3), int(along * 0.8)))
        # 短手には下限（0.35in）を置くが、**長手の 0.8 倍を超えさせない**——
        # 帯が薄いと下限のほうが勝って正方形に近づき、向きが読めなくなる
        # （Issue #141 と同じ症状が、下限の側から出る）．
        floor = min(max(Inches(0.35), int(long_ * 0.75)), int(long_ * 0.8))
        short = min(int(across * 0.9), floor)
        w, h = (long_, short) if horizontal else (short, long_)
        # 明示した大きさは**上限を超えてよい**——書いた人がそう決めたということ
        # （層をまたぐ 1.5×7.6cm の矢印は自動では書けない．Issue #143）．
        ew = self._resolve_len(arrow.width, width)
        eh = self._resolve_len(arrow.height, height)
        if ew is not None:
            w = int(ew)
        if eh is not None:
            h = int(eh)
        x = left + (width - w) // 2
        y = top + (height - h) // 2
        if h > height or w > width:
            # 帯に収まらない大きさを書かれたとき．**上端（左端）に寄せる**——
            # 中央のままだと上へも食い込み、導入文に重なる．はみ出す向きを
            # 下（結論文・罫線側）だけにするのは @overflow と同じ規約
            # （Issue #146）．黙って重ねずに知らせる．
            if h > height:
                y = top
            if w > width:
                x = left
            sys.stderr.write(
                "md2pptx: warning: the arrow is larger than the band it sits "
                "in and will overlap the text below (shorten the prose or "
                "reduce width/height)\n")
        shp = slide.shapes.add_shape(
            shape, Emu(x), Emu(y), Emu(w), Emu(h))
        both = arrow.direction in ("updown", "leftright")
        self._set_shape_adj(shp, self._ARROW_ADJ_BOTH if both
                            else self._ARROW_ADJ)
        shp.fill.solid()
        if arrow.color:
            kind, value = parse_color(arrow.color)
            if kind == "theme":
                shp.fill.fore_color.theme_color = self._theme_map[value]
            else:
                shp.fill.fore_color.rgb = RGBColor.from_string(value)
        else:
            shp.fill.fore_color.theme_color = self.GOLD
        shp.line.fill.background()
        shp.shadow.inherit = False
        shp.text_frame.word_wrap = False
        return shp

    @staticmethod
    def _set_shape_adj(shp: Shape, adj: dict[str, int]) -> None:
        """図形の調整値（``a:avLst/a:gd``）を書き込む．

        python-pptx は ``add_shape`` で空の ``avLst`` しか作らず、PowerPoint は
        そこを既定値で埋める．両端に矢じりのある形はその既定だと**箱いっぱいの
        菱形**になり、矢印に見えない（Issue #143）．
        """
        geom = shp._element.spPr.find(qn("a:prstGeom"))
        av = geom.find(qn("a:avLst")) if geom is not None else None
        if av is None:
            return          # プリセット図形でなければ調整値そのものが無い
        # **既にある調整値は捨ててから書く**．2 回呼んでも同じ結果になるように
        # ——同名の ``a:gd`` が並ぶと、どちらが効くかは実装依存になる．
        for old in list(av.findall(qn("a:gd"))):
            av.remove(old)
        for name, val in adj.items():
            gd = av.makeelement(qn("a:gd"), {"name": name,
                                             "fmla": f"val {val}"})
            av.append(gd)

    def draw_column_arrow(self, slide: PptxSlide, ncols: int) -> Shape | None:
        """カラムとカラムのすき間に、右向きの大きな矢印を描く（``@col: arrow``）．

        置き場は**左カラムの右端の内側**．テーマのカラム間のすき間は 0.5cm ほどしか
        無く、そこへ収めると矢印が糸のように細くなる（元の講義スライドも、すき間では
        なく左カラムの右寄りに置いてある）．

        箇条書きが長くて右端まで届くスライドでは**文字に重なる**．元のスライドも
        同じ作りで、そこは書く側が見て決める（SYNTAX.md に明記）．
        カラムが 2 つ無ければ何もしない．``@col: arrow`` はカラム区切りそのものなので、
        パーサを通る限りカラムは必ず 2 つ以上ある——ここは直接呼ばれたときの防御で、
        起きないことに警告は出さない．
        """
        if ncols < 2:
            return None
        a = self._find_placeholder(slide, 1)
        b = self._find_placeholder(slide, 2)
        if a is None or b is None:
            return None
        h = min(Inches(1.1), max(Inches(0.5), int(a.height * 0.18)))
        w = h
        inset = Inches(0.1)
        # 右端の内側．右カラムの本文には決して掛からない．
        x = min(a.left + a.width - w - inset, b.left - w - inset)
        # 縦は**上寄り**．矢印は「左の並びから右の並びへ」を指すもので、
        # 指す先は列の先頭にある．中央に置くと、項目数が少ないスライドで
        # 矢印だけが下に取り残される（元の講義スライドも上寄り）．
        y = a.top + int(a.height * 0.12)
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, Emu(x), Emu(y), Emu(w), Emu(h))
        self._set_shape_adj(shp, self._ARROW_ADJ)
        shp.fill.solid()
        shp.fill.fore_color.theme_color = self.GOLD
        shp.line.fill.background()
        shp.shadow.inherit = False
        shp.text_frame.word_wrap = False
        return shp

    def _stack_objects(self, slide: PptxSlide, objects: list[ObjectBlock],
                       left: int, top: int, width: int, height: int,
                       col_ratios: list[float] | None,
                       slide_overflow: bool = False,
                       has_prose_after: bool = False) -> None:
        """Table / Flow / Image を矩形領域内に重みづけで縦に積んで座標配置する．

        slide_overflow（@overflow）は表・画像に共通で効く．画像はブロックの
        overflow: 明示（Image.overflow が True/False）を優先する．Flow は帯に
        内接固定のため対象外（将来拡張）．
        """
        weights = [self._obj_weight(o) for o in objects]
        total = float(sum(weights)) or 1.0
        gap = Pt(6)
        avail = height - gap * (len(objects) - 1)
        y = top
        for obj, w in zip(objects, weights):
            seg_h = int(avail * w / total)
            if isinstance(obj, Arrow):
                self.render_arrow(slide, obj, left, y, width, seg_h)
            elif isinstance(obj, Flow):
                self.render_flow(slide, obj, left, y, width, seg_h)
            elif isinstance(obj, Seq):
                self.render_seq(slide, obj, left, y, width, seg_h)
            elif isinstance(obj, Image):
                eff = obj.overflow if obj.overflow is not None else slide_overflow
                self.render_image(slide, obj, left, y, width, seg_h,
                                  overflow=eff, has_prose_after=has_prose_after)
            else:
                self.render_table(slide, obj, left, y, width, seg_h, col_ratios,
                                  overflow=slide_overflow,
                                  has_prose_after=has_prose_after)
            y += seg_h + gap

    def _render_stacked(self, slide: PptxSlide, blocks: list[Block],
                        default_num_color: str | None, scale: float | None,
                        default_autofit: bool,
                        col_ratios: list[float] | None,
                        default_size_delta: int | None = None,
                        slide_overflow: bool = False) -> None:
        """表／図を含むスライドを描画する．

        地の文（Line）は **標準の本文プレースホルダ**へ流し込み，表・図だけを
        座標配置する．プレースホルダには「導入文＋空行スペーサ＋結論文」を入れ，
        確保した中央帯に表・図を重ねる（``参照スクリプト`` の図スライドと同方式）．
        地の文を自由位置のテキストボックスには置かない．
        """
        left, top, width, height = self._content_rect(slide)
        body = self._body_placeholder(slide)
        self._render_stacked_into(slide, blocks, body, left, top, width, height,
                                  default_num_color, scale, default_autofit,
                                  col_ratios, default_size_delta, slide_overflow)

    def _render_stacked_into(self, slide: PptxSlide, blocks: list[Block],
                             body: SlidePlaceholder | None, left: int, top: int,
                             width: int, height: int,
                             default_num_color: str | None,
                             scale: float | None, default_autofit: bool,
                             col_ratios: list[float] | None,
                             default_size_delta: int | None = None,
                             slide_overflow: bool = False) -> None:
        """``blocks`` を矩形 (left, top, width, height) 内へスタック描画する．

        地の文（Line）は ``body`` プレースホルダへ流し込み，表・図は矩形内に
        座標配置する．描画先（プレースホルダ＋矩形）を引数で受けるため，本文領域
        （単一カラム）にも多カラムの各カラム矩形にも使える．``body`` が None の
        場合は地の文を捨て，矩形全体にオブジェクトを積む．
        """
        # 地の文（前後）とオブジェクト（表・図）に分ける．
        # 図（Flow / Seq）の note(top)/note(bottom) も地の文としてプレースホルダ
        # へ回す．**Seq も同じ扱い**——分岐が Flow 限定だった頃、SYNTAX.md に
        # 載っている seq の note は丸ごと落ちていた（Issue #129）．
        prose_before: list[Line] = []
        objects: list[ObjectBlock] = []
        prose_after: list[Line] = []
        seen_obj = False
        for b in blocks:
            if is_object_block(b):
                if isinstance(b, (Flow, Seq)) and b.note_top:
                    bucket = prose_after if seen_obj else prose_before
                    ln = self._note_to_line(b.note_top)
                    if ln is not None:
                        bucket.append(ln)
                objects.append(b)
                seen_obj = True
                if isinstance(b, (Flow, Seq)) and b.note_bottom:
                    ln = self._note_to_line(b.note_bottom)
                    if ln is not None:
                        prose_after.append(ln)
            elif isinstance(b, Line):
                (prose_after if seen_obj else prose_before).append(b)
        if not objects:
            return

        # 相対サイズ（{+1} / @body-size）の警告はここにあったが、**要らなくなった**．
        # 帯の高さを本文標準サイズ固定で見積もっていたから食い違っていたのであって、
        # いまは para_h が行ごとの実サイズ（デルタ込み）で数える（Issue #145）．

        # 地の文が無ければプレースホルダは使わず，領域全体にオブジェクトを置く．
        if not prose_before and not prose_after:
            if body is not None:
                body._element.getparent().remove(body._element)
            self._stack_objects(slide, objects, left, top, width, height, col_ratios,
                                slide_overflow, has_prose_after=False)
            return

        # 地の文あり：プレースホルダに導入文＋空行＋結論文を流して中央帯を確保．
        # 帯と空行数はプレースホルダ矩形から逆算し，地の文＋空行＋結論文が
        # プレースホルダ高を超えないようにする（結論文がスライド外へ出ない）．
        bsz = self._body_font_size()
        # 段落の高さは**その段落のレベルのサイズ＋そのレベルの段落前アキ**で数える
        # （Issue #145）．どの行も lvl1 で数え、``spcBef`` を落としていた頃は
        # 帯が上へずれ、図が地の文に食い込んでいた．
        levels = (self._frame_font_levels(body.text_frame) if body is not None
                  else self._body_font_levels())

        # @autofit は**実際に描かれる字を縮める**ので、帯の計算もそれに合わせる．
        # 見ていなかった頃は、縮めたぶん空いた場所を帯が使えず、図が結論文へ
        # 食い込んでいた（Issue #145）．
        shrink = (scale / 100.0) if scale is not None else 1.0

        # 折り返しは**プレースホルダの幅**で数える（Issue #158）．1 行ずつと
        # 決めつけていた頃は、導入文が折り返すとそのぶん帯が上へずれ、
        # 結論文が下の罫線を越えていた（cn2026-02 p.34）．
        if body is not None:
            tf = body.text_frame
            tf_w = body.width - tf.margin_left - tf.margin_right
        else:
            # 枠が無ければ地の文も描かれない（``_warn_no_body``）ので、ここは
            # 帯の見積もりが 0 除算しないための置きにすぎない．
            tf_w = width

        def para_h(ln: Line) -> int:
            d = ln.size_delta if ln.size_delta is not None else default_size_delta
            base = levels[min(ln.level, len(levels) - 1)] if levels else bsz
            sz = self._size_from_delta(base, d) * shrink
            avail_pt = max(1.0, (tf_w - self._indent_for(ln.level)) / 12700.0)
            text = (ln.text or "").replace("\v", " ")
            return self._para_height(ln.level, sz,
                                     self._wrapped_lines(text, sz, avail_pt))

        before_h = sum(para_h(ln) for ln in prose_before)
        after_h = sum(para_h(ln) for ln in prose_after)
        # 帯を埋める空段落は lvl1 の書式（アキ込み）．
        blank_h = max(1, self._para_height(0, bsz * shrink))
        inset = Pt(4)
        band_h = height - before_h - after_h - 2 * inset
        if band_h < Inches(0.8):
            band_h = Inches(0.8)
        band_top = top + before_h + inset
        blanks = max(1, int(band_h / blank_h))   # 帯を埋める空行数（超過しない）

        if body is None:
            self._warn_no_body(prose_before + prose_after)
        else:
            tf = body.text_frame
            # 採番の状態は 2 回の追記で共有する——同じ枠なので，結論文で
            # 数え直すと番号が 1 に戻る（Issue #107）．
            counters: dict[tuple[int, str], int] = {}
            first = self._append_lines(tf, prose_before, True, default_num_color,
                                       default_size_delta, counters)
            for _ in range(blanks):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
            self._append_lines(tf, prose_after, first, default_num_color,
                               default_size_delta, counters)
            # 帯を持つスライドでは**自分で率を出さない**——空段落が枠を埋めて
            # いるので「収まっていない」ようには見えず、縮める必要が無い．
            # 明示の ``@autofit`` はそのまま効き、枠にも同じ率を渡す．
            self._apply_autofit(tf, scale, default_autofit)
            # 枠は導入文と結論文の両方に付けられる．結論文は空段落のぶんだけ
            # 下から始まるので、その数をそのまま渡す（Issue #133）．
            self.draw_line_boxes(slide, body, prose_before, default_size_delta,
                                 shrink=scale)
            self.draw_line_boxes(slide, body, prose_after, default_size_delta,
                                 preceding=prose_before, blank_paras=blanks,
                                 shrink=scale)

        # 帯の高さは **band_h ではなく空行数から** 求める（Issue #131）．
        # 結論文が描き始められる位置を決めるのは流し込んだ空行の数であって、
        # band_h ではない——blanks は int() で切り捨てるので、band_h をそのまま
        # 使うと最大 1 行ぶん帯のほうが下まで伸び、表が結論文に重なる。
        # しかも黙って重なる（表は帯高で描かれるので警告の条件に掛からない）。
        #   結論文の上端 = top + (nb + blanks) * line_h
        #   帯の上端     = band_top = top + nb * line_h + inset
        # なので帯に使えるのは blanks * line_h - inset まで．そこから
        # 従来どおり Pt(8) を余白として引く．
        _MIN_BAND = Inches(0.8)
        fits = blanks * blank_h - inset - Pt(8)
        obj_h = max(_MIN_BAND, fits)
        # 警告するのは**結論文があるときだけ**——下端に何も無ければ、帯が
        # 最小高まで広がっても重なる相手がいない．
        if prose_after and fits < _MIN_BAND:
            # 最小高（0.8in）に張り付くのは、地の文が枠をほぼ埋めて空行が
            # 1 行しか取れないとき．図を読める大きさに保つため下限は残すが、
            # その結果として結論文へ食い込む——**黙って重ねない**（Issue #131）．
            sys.stderr.write(
                "md2pptx: warning: too much body text for a table/figure "
                "slide; the band hit its minimum height and may overlap the "
                "concluding text (shorten the prose or split the slide)\n")
        self._stack_objects(slide, objects, left, band_top, width, obj_h, col_ratios,
                            slide_overflow, has_prose_after=bool(prose_after))

    # ------------------------------------------------------------- deck
    def render(self, deck: Deck) -> PptxPresentation:
        """Deck 全体を描画し，Presentation を返す．"""
        meta = deck.meta or {}
        # meta の値は object．どちらも描画側では真偽値としてしか使わないので、
        # ここで bool にする（if での評価と同じ結果になる）．
        slide_number = bool(meta.get("slide_number", True))
        default_autofit = bool(meta.get("default_autofit", True))
        # 等幅フォントだけはテーマに委ねきれない（§5.12）．テーマの本文フォントは
        # プロポーショナルで，桁が揃わないとコードとして読めない——見た目の好みでは
        # なく機能なので，既定を持ち，変えたい人は front matter で変える．
        mono = meta.get("mono_font")
        self._mono_font = str(mono) if mono else DEFAULT_MONO_FONT

        if deck.title_slide is not None:
            self.render_title_slide(deck.title_slide)

        for sl in deck.slides:
            self.render_slide(
                sl,
                slide_number=slide_number,
                default_autofit=default_autofit,
            )
        # 言語の付与は描画の最後に一度だけ通す（run を作る経路は本文・タイトル・表・
        # フロー図・ノートと多く，作る側それぞれに足すと必ずどこかが漏れる）．
        self._apply_text_language()
        return self.prs

    def _apply_text_language(self) -> None:
        """描画した全 run に言語（``a:rPr/@lang``）を付ける．**禁則処理の要**（Issue #79）．

        PowerPoint は**行分割の規則を run の言語で選ぶ**．python-pptx は
        ``paragraph.text = …`` で ``<a:r><a:t>…</a:t></a:r>`` を作り ``a:rPr`` を
        **一切書かない**ので，何もしないと言語不明のまま出力される．そうなると
        日本語の禁則処理が適用されず，**行頭に「ー」や句読点が来る**——文字列は
        正しいので pptx を開くまで気づけず，開いても原因が md2pptx だとは思い当たらない．

        効くのは ``lang`` だけで，紛らわしい近縁の属性・要素は**効かない**（実 PowerPoint
        変換で確認済み．固定しているのは ``tests/test_text_language.py``）:

        - ``kumimoji="1"`` のみ付けても禁則は効かない（これは縦書き中の数字の扱い）．
        - ``presentation.xml`` の ``<p:kinsoku>``（禁則文字の定義）を足しても効かない．
          run が何語か決まらないうちは，禁則文字表を引く段階に来ない．

        テーマ側の既存 run（スライド番号フィールド等）は ``lang`` を持っているので
        触らない．**上書きせず未設定のものだけ埋める**——何度通しても結果が変わらず，
        将来 run 単位で言語を決める余地も残る．``altLang``（もう一方の字種の言語）も
        同じ規則で，既に決まっていればそのまま残す．

        レイアウト・マスターは書き換えない（テーマの所有物で，md2pptx が描くのは
        スライドとノートだけ）．ノートも通すのは，発表者ノートが折り返す先も
        同じ理由で崩れるため．
        """
        for slide in self.prs.slides:
            parts = [slide.element]
            if slide.has_notes_slide:
                parts.append(slide.notes_slide.element)
            for part in parts:
                for r in part.iter(qn("a:r")):
                    rPr = r.get_or_add_rPr()
                    if rPr.get("lang") is not None:
                        continue            # 言語の決まっている run は触らない
                    rPr.set("lang", self._LANG)
                    # altLang も**未設定のときだけ**付ける．lang と対で書くのが
                    # 普通だが，片方だけ持つ run が無いとは言えない——その 1 つを
                    # 上書きすると「決まっているものは触らない」が崩れる．
                    if rPr.get("altLang") is None:
                        rPr.set("altLang", self._ALT_LANG)

    def save(self, path: str) -> str:
        """現在の Presentation を保存する（差し替えは**アトミック**）．

        出力先へ直接は書かない．同じディレクトリに使い捨ての作業ディレクトリを作って
        そこへ保存し，``os.replace`` で置き換える．作業場所は出力先ディレクトリの中に
        新しく作るので必ず同一ファイルシステム上にあり，置き換えは常にアトミックになる
        （EXDEV は起こりえない）．``pdf.convert`` と同じ形（Issue #56 / #53）．

        - **保存中も前回の pptx がそのまま読める．** python-pptx の ``save`` は出力先を
          切り詰めてから書くので，直接書くと「開いたら壊れていた」時間が生まれる．
          PowerPoint で開いたまま作り直す／``--watch`` と手で叩くのを併用する，という
          使い方で中途半端な pptx を掴ませない．
        - **失敗したら前回の pptx を残す．** ここは PDF と**逆の契約**で，理由は失敗の
          見え方が違うこと——PDF 変換の失敗は終了コードを変えない（警告だけ）ので，
          古い PDF を残すと「新しい出力」と取り違えられる．pptx の保存失敗は cli が
          ``BuildError`` にして終了コード 1 で終えるため取り違えようがなく，それなら
          主成果物（PowerPoint で開いているかもしれない）を消さない方がよい．
        """
        directory = os.path.dirname(os.path.abspath(path))
        try:
            work = workdir.create(directory)
        except OSError as e:
            # 素の errno を通すと「一時ディレクトリ名が読めない形で出る」だけになるので，
            # 何をしようとして失敗したかを添える（cli が整形して表示する）．
            # ``from e`` で元の例外を __cause__ に残す——errno を見たいときの手掛かり．
            raise OSError(
                f"cannot create a working directory in {directory} ({e})") from e
        try:
            staged = os.path.join(work, os.path.basename(path))
            self.prs.save(staged)
            os.replace(staged, path)
        finally:
            # 成功後は staged が移動済みで work は空，失敗時は書きかけが中に残る．
            # どちらも同じ扱いで片付ける（規則と理由は workdir.discard に 1 か所だけ）．
            workdir.discard(work)
        return path


def build(deck: Deck, base_pptx_path: str, out_path: str,
          base_dir: str | None = None) -> str:
    """Deck を base pptx 上に描画して out_path に保存する（CLI 用エントリ）．

    Args:
        base_dir: 画像などの相対パスを解決する基準ディレクトリ（既定は Markdown の
            置き場）．None なら実行時のカレントを基準にする．

    Returns:
        out_path（保存先パス）．
    """
    r = Renderer(base_pptx_path, base_dir=base_dir)
    r.render(deck)
    r.save(out_path)
    return out_path
