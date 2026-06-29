#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""IR（ir.py）→ pptx 描画（md2pptx ステージ2 / Phase 1）．

ステージ0（thmx2pptx.py）が生成した base pptx を土台に開き，IR の Deck を
走査して 1 つのプレゼンテーションへ描画する．配色・フォントはテーマ
（thmx）任せで，スクリプト側で色・フォントをハードコードしない（図形のみ
テーマのアクセント色を参照する）．

``参照スクリプト`` のモジュール大域に依存したヘルパ群（no_bullet /
add_slide_number / add_bullets / set_autonum / enum_items / fit_body /
content_slide）を Renderer のメソッドへ移植し，self.prs / self.layouts /
テーマ色エイリアスから状態を解決する（DESIGN.md §6）．

使い方::

    from thmx2pptx import thmx_to_pptx
    from ir import Deck
    from render import build

    base = thmx_to_pptx("theme.thmx")
    build(deck, base, "out.pptx")
"""
from __future__ import annotations

import copy
import math
import sys

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

try:  # パッケージ実行・単体実行のどちらでも import できるように
    from .ir import Deck, Slide, TitleSlide, Line, Table, Flow
    from .flow import plan_flow
except ImportError:  # pragma: no cover - 単体実行時のフォールバック
    from ir import Deck, Slide, TitleSlide, Line, Table, Flow
    from flow import plan_flow


class Renderer:
    """IR を pptx へ描画するレンダラ．

    base pptx（テーマのみを持つ 0 枚構成）を開き，レイアウトとテーマ色
    エイリアスを初期化する．スライドはすべて新規追加で生成する
    （thmx 由来の base は本文スライドを持たない）．
    """

    def __init__(self, base_pptx_path):
        self.prs = Presentation(base_pptx_path)
        self.SW = self.prs.slide_width
        self.SH = self.prs.slide_height

        # テーマのアクセント色（図形用）．テキスト色・フォントはテーマ任せ．
        # Phase 1 では box/arrow/note/table（Phase 2/3）が使うため保持のみ．
        self.A2 = MSO_THEME_COLOR.ACCENT_2       # 緑
        self.A6 = MSO_THEME_COLOR.ACCENT_6       # 緑（濃）
        self.T2 = MSO_THEME_COLOR.TEXT_2         # 緑系テキスト
        self.GOLD = MSO_THEME_COLOR.ACCENT_1     # 金
        self.BG = MSO_THEME_COLOR.BACKGROUND_1   # 背景（白）
        self.TX = MSO_THEME_COLOR.TEXT_1         # 本文色（黒）

        # フロー図 box の自動配色（テーマアクセント色を順番に割当）．
        self._box_palette = [self.T2, self.A6, self.A6, self.GOLD, self.A2]
        # テーマ色名（DSL の {accent6} 等）→ MSO_THEME_COLOR の対応表．
        self._theme_map = {
            "accent1": MSO_THEME_COLOR.ACCENT_1, "accent2": MSO_THEME_COLOR.ACCENT_2,
            "accent3": MSO_THEME_COLOR.ACCENT_3, "accent4": MSO_THEME_COLOR.ACCENT_4,
            "accent5": MSO_THEME_COLOR.ACCENT_5, "accent6": MSO_THEME_COLOR.ACCENT_6,
            "tx1": MSO_THEME_COLOR.TEXT_1, "tx2": MSO_THEME_COLOR.TEXT_2,
            "bg1": MSO_THEME_COLOR.BACKGROUND_1, "bg2": MSO_THEME_COLOR.BACKGROUND_2,
        }
        # 本文スタイルのレベル別フォントサイズ（pt）のキャッシュ（None=未取得）．
        self._body_levels = None

        # レイアウト解決．title=0 / content=1 / section=2．
        layouts = self.prs.slide_layouts
        self.layouts = layouts
        self.title_layout = layouts[0]
        self.L1 = layouts[1] if len(layouts) > 1 else layouts[0]
        self.section_layout = layouts[2] if len(layouts) > 2 else self.L1

    # ------------------------------------------------------------ helpers
    def no_bullet(self, para):
        """段落の行頭記号を消す（結論行など）．"""
        pPr = para._p.get_or_add_pPr()
        for tag in ("a:buChar", "a:buAutoNum"):
            e = pPr.find(qn(tag))
            if e is not None:
                pPr.remove(e)
        if pPr.find(qn("a:buNone")) is None:
            pPr.append(pPr.makeelement(qn("a:buNone"), {}))

    def add_slide_number(self, slide):
        """スライド自身のレイアウトの番号プレースホルダ（idx==12）を複製して有効化．

        セクションスライド（レイアウト2）など L1 以外の上でも，そのスライドの
        実レイアウトから番号プレースホルダを取得する．
        """
        for lph in slide.slide_layout.placeholders:
            if lph.placeholder_format.idx == 12:
                slide.shapes._spTree.append(copy.deepcopy(lph._element))
                return

    def add_bullets(self, tf, items):
        """(level, text) の列を本文 text_frame に箇条書きとして流し込む．"""
        first = True
        for lvl, txt in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = lvl
            p.text = txt

    def set_autonum(self, p, fmt="arabicPeriod", color=None):
        """段落の行頭記号を自動採番（1. 2. 3. …）に切り替える（enumerate 相当）．
           color にテーマ色名（例 "tx1"）を渡すと採番記号の色を指定する．"""
        pPr = p._p.get_or_add_pPr()
        if color:
            for tag in ("a:buClrTx", "a:buClr"):
                el = pPr.find(qn(tag))
                if el is not None:
                    pPr.remove(el)
            buClr = pPr.makeelement(qn("a:buClr"), {})
            buClr.append(buClr.makeelement(qn("a:schemeClr"), {"val": color}))
            pPr.insert(0, buClr)  # buClr は採番記号より前に置く
        bu = pPr.makeelement(qn("a:buAutoNum"), {"type": fmt})
        for tag in ("a:buChar", "a:buNone", "a:buAutoNum"):
            el = pPr.find(qn(tag))
            if el is not None:
                pPr.replace(el, bu)
                break
        else:
            pPr.append(bu)

    def enum_items(self, tf, items):
        """(見出し, 説明) を，見出し=自動採番(level0)・説明=通常箇条書き(level1) で並べる．"""
        first = True
        for head, desc in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = 0
            p.text = head
            self.set_autonum(p)
            d = tf.add_paragraph()
            d.level = 1
            d.text = desc

    def fit_body(self, tf, scale=None):
        """本文プレースホルダに自動調整（normAutofit）を設定する．
        scale を与えると縮小率（%）を明示的に焼き込む．"""
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if scale is not None:
            bodyPr = tf._txBody.bodyPr
            na = bodyPr.find(qn("a:normAutofit"))
            if na is None:
                # auto_size 設定で生成されない環境向けのフォールバック
                na = bodyPr.makeelement(qn("a:normAutofit"), {})
                bodyPr.append(na)
            na.set("fontScale", str(int(scale * 1000)))

    def _body_placeholder(self, slide):
        """本文プレースホルダ（idx==1）を返す．無ければ None．"""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                return ph
        return None

    def _body_font_levels(self):
        """本文スタイルのレベル別フォントサイズ（pt）のリストを返す（lvl1 始まり）．

        スライドマスターの ``p:txStyles/p:bodyStyle/a:lvl{n}pPr/a:defRPr@sz``
        を lvl1 から順に読み取る．表・図が標準サイズで収まらないとき，下位レベルの
        小さいサイズへ段階的に切り替えるために用いる．取得できなければ既定 [18]．
        """
        if self._body_levels is not None:
            return self._body_levels
        levels = []
        try:
            master = self.prs.slide_masters[0]
            body = master.element.find(
                qn("p:txStyles") + "/" + qn("p:bodyStyle"))
            if body is not None:
                for lvl in range(1, 10):
                    el = body.find(
                        qn("a:lvl%dpPr" % lvl) + "/" + qn("a:defRPr"))
                    if el is not None and el.get("sz"):
                        levels.append(int(el.get("sz")) / 100.0)
        except Exception:
            pass
        if not levels:
            levels = [18.0]
        self._body_levels = levels
        return levels

    def _body_font_size(self):
        """本文プレースホルダの標準フォントサイズ（pt．lvl1）を返す．"""
        return self._body_font_levels()[0]

    @staticmethod
    def _text_width_pt(text, font_pt):
        """テキストの概算表示幅（pt）．全角は font_pt，半角は約 0.55×で見積もる．"""
        w = 0.0
        for ch in text or "":
            w += font_pt if ord(ch) > 0x2E80 else font_pt * 0.55
        return w

    def _fit_font(self, fits_at):
        """レベル別サイズを大きい順に試し，``fits_at(size)`` が真の最大サイズを返す．

        どのレベルでも収まらなければ最小レベル（最後）のサイズを返す（ベストエフォート）．
        """
        levels = self._body_font_levels()
        for sz in levels:
            if fits_at(sz):
                return sz
        return levels[-1]

    def content_slide(self, title, items):
        """タイトル＋箇条書きの基本スライドを 1 枚追加して返す．"""
        s = self.prs.slides.add_slide(self.L1)
        s.shapes.title.text = title
        body = self._body_placeholder(s)
        self.add_bullets(body.text_frame, items)
        self.fit_body(body.text_frame)
        self.add_slide_number(s)
        return s

    # --------------------------------------------------------- title slide
    def render_title_slide(self, ts: TitleSlide):
        """タイトルスライドを 1 枚目として追加する（番号は付けない）．

        base は 0 枚構成のため，タイトルレイアウト上に新規スライドを追加し，
        CENTER_TITLE（idx==0）に title を，SUBTITLE（idx==1）に
        subtitle＋author＋affiliation を流し込む．
        """
        s = self.prs.slides.add_slide(self.title_layout)

        title_ph = self._find_placeholder(s, 0)
        if title_ph is not None and ts.title:
            tf = title_ph.text_frame
            lines = ts.title.split("\n")
            tf.paragraphs[0].text = lines[0]
            for ln in lines[1:]:
                tf.add_paragraph().text = ln

        sub_ph = self._find_placeholder(s, 1)
        if sub_ph is not None:
            tf = sub_ph.text_frame
            sub_lines = []
            if ts.subtitle:
                sub_lines.append(ts.subtitle)
            if ts.author:
                sub_lines.append(ts.author)
            sub_lines.extend(ts.affiliation or [])
            if sub_lines:
                tf.paragraphs[0].text = sub_lines[0]
                for ln in sub_lines[1:]:
                    tf.add_paragraph().text = ln
        return s

    def _find_placeholder(self, slide, idx):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    # ----------------------------------------------------------- 図形（flow）
    def box(self, slide, l, t, w, h, text, tc, sub=None, fsize=None, ssize=None):
        """角丸四角ノードを描く（塗りはテーマ色 tc，文字は背景色 BG）．

        fsize / ssize（pt）を省略するとテーマ既定のフォントサイズを継承する．
        """
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        shp.fill.solid()
        shp.fill.fore_color.theme_color = tc
        shp.line.color.theme_color = self.TX
        shp.line.width = Pt(0.5)
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
        pa = tf.paragraphs[0]
        pa.alignment = PP_ALIGN.CENTER
        pa.text = text
        for r in pa.runs:
            r.font.color.theme_color = self.BG
            if fsize is not None:
                r.font.size = Pt(fsize)
            r.font.bold = True
        if sub:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.text = sub
            for r in p2.runs:
                r.font.color.theme_color = self.BG
                if ssize is not None:
                    r.font.size = Pt(ssize)
        return shp

    def arrow(self, slide, x1, y1, x2, y2, w=2.0):
        """矢印（直線コネクタ＋三角の矢じり）を描く．"""
        cn = slide.shapes.add_connector(2, x1, y1, x2, y2)
        cn.line.color.theme_color = self.TX
        cn.line.width = Pt(w)
        le = cn.line._get_or_add_ln()
        le.append(le.makeelement(
            qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
        return cn

    def note(self, slide, l, t, w, h, text, size, tc=None, bold=False,
             align=PP_ALIGN.LEFT):
        """注記用テキストボックスを描く（キャプション・矢印ラベル・省略記号）．"""
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        pa = tf.paragraphs[0]
        pa.alignment = align
        pa.text = text
        for r in pa.runs:
            if size is not None:
                r.font.size = Pt(size)
            if tc is not None:
                r.font.color.theme_color = tc
            if bold:
                r.font.bold = True
        return tb

    def _theme_color(self, name):
        """テーマ色名を MSO_THEME_COLOR へ解決する（未知なら None）．"""
        if not name:
            return None
        return self._theme_map.get(str(name).lower())

    def _box_fits(self, node, bw, bh, font_pt):
        """box（主ラベル＋副ラベル）が指定フォントサイズで収まるか概算判定する．"""
        line_h = font_pt * 1.2
        inner_w = max(1.0, bw / 12700.0 - 8)   # 左右マージン約 Pt(4)×2
        inner_h = bh / 12700.0 - 4             # 上下マージン約
        lines = max(1, math.ceil(self._text_width_pt(node.label, font_pt) / inner_w))
        if node.sublabel:
            lines += max(1, math.ceil(
                self._text_width_pt(node.sublabel, font_pt) / inner_w))
        return lines * line_h <= inner_h

    def render_flow(self, slide, flow, left, top, width, height):
        """Flow ブロックを矩形領域へ描画する（flow.plan_flow の座標プランを使用）．

        図中の文字サイズは本文プレースホルダの標準サイズに揃える．
        """
        plan = plan_flow(flow, left, top, width, height)
        # box が標準サイズで収まらなければ，全 box 一律で下位レベルへ切り替える．
        boxes = plan["boxes"]
        if boxes:
            bsz = self._fit_font(
                lambda sz: all(self._box_fits(node, bw, bh, sz)
                               for node, _, _, bw, bh in boxes))
        else:
            bsz = self._body_font_size()
        bi = 0
        for node, bl, bt, bw, bh in plan["boxes"]:
            tc = self._theme_color(node.color) or \
                self._box_palette[bi % len(self._box_palette)]
            bi += 1
            self.box(slide, bl, bt, bw, bh, node.label, tc,
                     sub=node.sublabel or None, fsize=bsz, ssize=bsz)
        for label, bl, bt, bw, bh in plan["ellipses"]:
            self.note(slide, bl, bt, bw, bh, label, bsz, tc=self.T2, bold=True,
                      align=PP_ALIGN.CENTER)
        for x1, y1, x2, y2 in plan["arrows"]:
            self.arrow(slide, x1, y1, x2, y2)
        for text, bl, bt, bw, bh in plan["labels"]:
            self.note(slide, bl, bt, bw, bh, text, bsz, tc=self.T2, bold=True,
                      align=PP_ALIGN.CENTER)
        for text, bl, bt, bw, bh, role in plan["captions"]:
            # caption のみ図に付随（note_top / note_bottom はプレースホルダ側で描く）．
            if role == "caption":
                self.note(slide, bl, bt, bw, bh, text, bsz, tc=self.T2,
                          align=PP_ALIGN.CENTER)

    def _table_col_widths(self, ncols, width, col_ratios):
        """表の列幅（EMU）リストを返す（均等 or 比率指定）．"""
        if col_ratios and len(col_ratios) == ncols and sum(col_ratios) > 0:
            tot = float(sum(col_ratios))
            return [int(width * r / tot) for r in col_ratios]
        cw = int(width / ncols)
        return [cw] * ncols

    def _table_height_emu(self, data, col_w, font_pt):
        """指定フォントサイズでの表の概算総高（EMU）を見積もる（折り返し考慮）．

        実際の PowerPoint レンダリングは行間・最小行高などで見積りより伸びがち
        なため，安全係数を掛けて保守的（やや大きめ）に見積もる．
        """
        line_h = font_pt * 1.32          # 行間込みの行高
        cell_pad_pt = 6                  # セル上下マージン＋最小余白（約）
        side_pad_pt = 18                 # セル左右マージン合計（約．Pt(10)+Pt(6)＋余裕）
        safety = 1.15                    # 折り返し・最小行高ぶんの安全係数
        total_pt = 0.0
        for row in data:
            row_h = line_h + cell_pad_pt
            for ci, cw in enumerate(col_w):
                text = row[ci] if ci < len(row) else ""
                inner_pt = max(1.0, cw / 12700.0 - side_pad_pt)
                lines = max(1, math.ceil(self._text_width_pt(text, font_pt) / inner_pt))
                row_h = max(row_h, lines * line_h + cell_pad_pt)
            total_pt += row_h
        return int(total_pt * safety * 12700)

    def render_table(self, slide, table, left, top, width, height, col_ratios=None):
        """Table ブロックを座標指定で 1 つ描画する（ヘッダ行をアクセント色で着色）．

        ``参照スクリプト`` の表描画を移植・一般化したもの．列幅は既定で均等，
        ``col_ratios`` を与えると比率配分する．配色はテーマ任せ（ヘッダのみ
        アクセント色 A2＋背景色 BG の文字）．
        """
        nrows = len(table.rows) + (1 if table.header else 0)
        ncols = max(
            len(table.header) if table.header else 0,
            max((len(r) for r in table.rows), default=0),
        )
        if nrows == 0 or ncols == 0:
            return None

        gf = slide.shapes.add_table(nrows, ncols, left, top, width, height)
        tbl = gf.table

        col_w = self._table_col_widths(ncols, width, col_ratios)
        for ci, cw in enumerate(col_w):
            tbl.columns[ci].width = cw

        data = ([table.header] if table.header else []) + list(table.rows)
        # フォントは本文標準（lvl1）を基本に，収まらなければ下位レベルへ切り替える．
        fsize = self._fit_font(
            lambda sz: self._table_height_emu(data, col_w, sz) <= height)

        for ri, row in enumerate(data):
            is_header = bool(table.header) and ri == 0
            for ci in range(ncols):
                cell = tbl.cell(ri, ci)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = Pt(10)
                cell.margin_right = Pt(6)
                cell.margin_top = Pt(2)
                cell.margin_bottom = Pt(2)
                pa = cell.text_frame.paragraphs[0]
                pa.text = row[ci] if ci < len(row) else ""
                for run in pa.runs:
                    run.font.size = Pt(fsize)
                    if is_header:
                        run.font.bold = True
                        run.font.color.theme_color = self.BG
                if is_header:
                    cell.fill.solid()
                    cell.fill.fore_color.theme_color = self.A2
        return gf

    # ------------------------------------------------------- content slide
    def render_slide(self, slide: Slide, slide_number=True, default_autofit=True):
        """コンテンツスライドを 1 枚追加して返す．

        ``slide.blocks`` を出現順に処理する．表を含まないスライドは Phase 1 と
        同じく本文プレースホルダへ Line を流し込み，表を含むスライドは座標スタック
        配置（テキスト→表→テキスト …）で描画する．Flow は Phase 3 でスキップ．
        """
        directives = slide.directives or {}
        layout_idx = directives.get("layout", slide.layout)
        try:
            layout = self.layouts[layout_idx]
        except (IndexError, TypeError):
            layout = self.L1

        s = self.prs.slides.add_slide(layout)
        if slide.title is not None and s.shapes.title is not None:
            s.shapes.title.text = slide.title

        # スライド既定の採番色（@autonum-color）．Line.num_color が優先．
        default_num_color = directives.get("autonum_color")
        scale = self._autofit_scale(directives)
        blocks = slide.blocks or []

        if any(isinstance(b, (Table, Flow)) for b in blocks):
            self._render_stacked(s, blocks, default_num_color, scale, default_autofit,
                                 self._col_ratios(directives))
        else:
            line_blocks = [b for b in blocks if isinstance(b, Line)]
            body = self._body_placeholder(s)
            if body is not None and line_blocks:
                tf = body.text_frame
                self._fill_lines(tf, line_blocks, default_num_color)
                self._apply_autofit(tf, scale, default_autofit)

        if slide_number:
            self.add_slide_number(s)
        return s

    # ----------------------------------------------------- 描画ユーティリティ
    def _append_lines(self, tf, line_blocks, first, default_num_color):
        """Line 列を text_frame に段落として追記する（採番／no_bullet を適用）．

        first=True なら最初の 1 行は既存の paragraphs[0] を使う．残りの行を
        追記しても良いよう，処理後の first 状態を返す．
        """
        for blk in line_blocks:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = blk.level
            p.text = blk.text
            if blk.kind == "autonum":
                fmt = blk.num_style or "arabicPeriod"
                color = blk.num_color or default_num_color
                self.set_autonum(p, fmt, color=color)
            elif blk.kind == "plain":
                self.no_bullet(p)
            # kind == "bullet" はテーマ既定のまま
        return first

    def _fill_lines(self, tf, line_blocks, default_num_color):
        """Line 列を text_frame の段落として流し込む（先頭から）．"""
        self._append_lines(tf, line_blocks, True, default_num_color)

    def _autofit_scale(self, directives):
        """@autofit ディレクティブを縮小率へ解釈する（非数値は警告して None）．"""
        autofit = directives.get("autofit")
        if autofit is None:
            return None
        try:
            return float(autofit)
        except (TypeError, ValueError):
            sys.stderr.write(
                f"md2pptx: warning: ignoring non-numeric @autofit value "
                f"{autofit!r}\n"
            )
            return None

    def _apply_autofit(self, tf, scale, default_autofit):
        """縮小率指定があれば焼き込み，無ければ既定の自動調整を設定する．"""
        if scale is not None:
            self.fit_body(tf, scale=scale)
        elif default_autofit:
            self.fit_body(tf)

    def _col_ratios(self, directives):
        """@col-widths ディレクティブ（"45,55" 等）を比率リストへ解釈する．"""
        v = directives.get("col_widths")
        if not v:
            return None
        try:
            return [float(x) for x in str(v).replace("，", ",").split(",")]
        except ValueError:
            return None

    def _content_rect(self, slide):
        """本文領域の矩形 (left, top, width, height) を返す（座標配置の基準）．"""
        ph = self._body_placeholder(slide)
        if ph is not None and None not in (ph.left, ph.top, ph.width, ph.height):
            return (ph.left, ph.top, ph.width, ph.height)
        try:
            for lph in slide.slide_layout.placeholders:
                if lph.placeholder_format.idx == 1 and None not in (
                    lph.left, lph.top, lph.width, lph.height
                ):
                    return (lph.left, lph.top, lph.width, lph.height)
        except Exception:
            pass
        # 既定：タイトル下の本文相当領域．
        return (Inches(0.6), Inches(1.7), self.SW - Inches(1.2),
                self.SH - Inches(2.3))

    def _note_to_line(self, text):
        """Flow の note 文字列を本文プレースホルダ用の Line へ変換する．

        行頭マーカー（→ は no_bullet）の最小限の解釈だけ行う．
        """
        t = (text or "").strip()
        if t.startswith("→"):
            return Line(text=t[len("→"):].lstrip(), kind="plain")
        return Line(text=t, kind="bullet")

    def _obj_weight(self, obj):
        """オブジェクト（Table / Flow）の縦方向の重み（高さ配分用）．"""
        if isinstance(obj, Flow):
            return max(4, len(obj.nodes) + 2)
        return max(2, len(obj.rows) + (1 if obj.header else 0))

    def _stack_objects(self, slide, objects, left, top, width, height, col_ratios):
        """Table / Flow を矩形領域内に重みづけで縦に積んで座標配置する．"""
        weights = [self._obj_weight(o) for o in objects]
        total = float(sum(weights)) or 1.0
        gap = Pt(6)
        avail = height - gap * (len(objects) - 1)
        y = top
        for obj, w in zip(objects, weights):
            seg_h = int(avail * w / total)
            if isinstance(obj, Flow):
                self.render_flow(slide, obj, left, y, width, seg_h)
            else:
                self.render_table(slide, obj, left, y, width, seg_h, col_ratios)
            y += seg_h + gap

    def _render_stacked(self, slide, blocks, default_num_color, scale,
                        default_autofit, col_ratios):
        """表／図を含むスライドを描画する．

        地の文（Line）は **標準の本文プレースホルダ**へ流し込み，表・図だけを
        座標配置する．プレースホルダには「導入文＋空行スペーサ＋結論文」を入れ，
        確保した中央帯に表・図を重ねる（``参照スクリプト`` の図スライドと同方式）．
        地の文を自由位置のテキストボックスには置かない．
        """
        left, top, width, height = self._content_rect(slide)
        body = self._body_placeholder(slide)

        # 地の文（前後）とオブジェクト（表・図）に分ける．
        # Flow の note(top)/note(bottom) も地の文としてプレースホルダへ回す．
        prose_before, objects, prose_after = [], [], []
        seen_obj = False
        for b in blocks:
            if isinstance(b, (Table, Flow)):
                if isinstance(b, Flow) and b.note_top:
                    bucket = prose_after if seen_obj else prose_before
                    bucket.append(self._note_to_line(b.note_top))
                objects.append(b)
                seen_obj = True
                if isinstance(b, Flow) and b.note_bottom:
                    prose_after.append(self._note_to_line(b.note_bottom))
            elif isinstance(b, Line):
                (prose_after if seen_obj else prose_before).append(b)
        if not objects:
            return

        # 地の文が無ければプレースホルダは使わず，領域全体にオブジェクトを置く．
        if not prose_before and not prose_after:
            if body is not None:
                body._element.getparent().remove(body._element)
            self._stack_objects(slide, objects, left, top, width, height, col_ratios)
            return

        # 地の文あり：プレースホルダに導入文＋空行＋結論文を流して中央帯を確保．
        # 帯と空行数はプレースホルダ矩形から逆算し，地の文＋空行＋結論文が
        # プレースホルダ高を超えないようにする（結論文がスライド外へ出ない）．
        bsz = self._body_font_size()
        line_h = int(Pt(bsz) * 1.32)        # 行間込みの保守的な行高
        nb, na = len(prose_before), len(prose_after)
        inset = Pt(4)
        band_h = height - (nb + na) * line_h - 2 * inset
        if band_h < Inches(0.8):
            band_h = Inches(0.8)
        band_top = top + nb * line_h + inset
        blanks = max(1, int(band_h / line_h))   # 帯を埋める空行数（超過しない）

        if body is not None:
            tf = body.text_frame
            first = self._append_lines(tf, prose_before, True, default_num_color)
            for _ in range(blanks):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
            self._append_lines(tf, prose_after, first, default_num_color)
            self._apply_autofit(tf, scale, default_autofit)

        # 結論文との重なりを避けるため帯を少しだけ詰めてオブジェクトを置く．
        obj_h = max(Inches(0.8), band_h - Pt(8))
        self._stack_objects(slide, objects, left, band_top, width, obj_h, col_ratios)

    # ------------------------------------------------------------- deck
    def render(self, deck: Deck):
        """Deck 全体を描画し，Presentation を返す．"""
        meta = deck.meta or {}
        slide_number = meta.get("slide_number", True)
        default_autofit = meta.get("default_autofit", True)

        if deck.title_slide is not None:
            self.render_title_slide(deck.title_slide)

        for sl in deck.slides:
            self.render_slide(
                sl,
                slide_number=slide_number,
                default_autofit=default_autofit,
            )
        return self.prs

    def save(self, path):
        """現在の Presentation を保存する．"""
        self.prs.save(path)
        return path


def build(deck, base_pptx_path, out_path):
    """Deck を base pptx 上に描画して out_path に保存する（CLI 用エントリ）．

    Returns:
        out_path（保存先パス）．
    """
    r = Renderer(base_pptx_path)
    r.render(deck)
    r.save(out_path)
    return out_path
