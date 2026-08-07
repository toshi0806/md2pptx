#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""``` ```arrow ``` フェンスを固定するテスト（Issue #137）．

#134 では「``↓`` だけの行」を矢印にしていたが、それだと**本文に ``↓`` の1文字を
書けなくなる**。他のブロック（flow / seq / image / note）はすべてフェンスなので、
矢印もフェンスに揃える。

`<!-- @col: arrow -->`（カラム間の右矢印）は別物として残る——あちらは
「区切りそのものを図形として描く」指定で、本文に置くブロックではない。
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

from md2pptx import render
from md2pptx.ir import Arrow
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"


def _blocks(body):
    slide, = parse(_FM + "### x\n\n" + body + "\n").slides
    return slide.blocks


def _build(tmp_path, src):
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _shapes(slide, kind):
    return [sh for sh in slide.shapes
            if not sh.is_placeholder
            and getattr(sh, "auto_shape_type", None) == kind]


def _fence(direction):
    return "```arrow\ndirection: %s\n```" % direction


# ---------------------------------------------------------------- 記法

@pytest.mark.parametrize(
    "name", ["down", "up", "right", "left", "updown", "leftright"])
def test_every_direction_parses(name):
    arrow, = [b for b in _blocks(_fence(name)) if isinstance(b, Arrow)]
    assert arrow.direction == name


def test_the_direction_is_required():
    """向きの無い矢印は描きようが無いのでエラー．"""
    with pytest.raises(ValueError, match="direction"):
        _blocks("```arrow\n```")


def test_an_unknown_direction_stops():
    with pytest.raises(ValueError, match="direction"):
        _blocks(_fence("sideways"))


def test_an_unknown_key_stops():
    """知らないキーはタイポとみなして止める（他のフェンスと同じ）．"""
    with pytest.raises(ValueError, match="arrow"):
        _blocks("```arrow\ndirection: down\ncolour: red\n```")


def test_a_lone_arrow_glyph_is_just_text():
    """``↓`` だけの行は**ふつうの本文**に戻った（記号を予約しない）．"""
    blocks = _blocks("↓")
    assert not any(isinstance(b, Arrow) for b in blocks)
    assert blocks[0].text == "↓"


def test_an_arrow_glyph_in_a_sentence_is_untouched():
    blocks = _blocks("- 上から ↓ へ進む")
    assert not any(isinstance(b, Arrow) for b in blocks)
    assert blocks[0].text == "上から ↓ へ進む"


# ---------------------------------------------------------------- 描画

@pytest.mark.parametrize("name,shape", [
    ("down", MSO_SHAPE.DOWN_ARROW),
    ("up", MSO_SHAPE.UP_ARROW),
    ("right", MSO_SHAPE.RIGHT_ARROW),
    ("left", MSO_SHAPE.LEFT_ARROW),
    ("updown", MSO_SHAPE.UP_DOWN_ARROW),
    ("leftright", MSO_SHAPE.LEFT_RIGHT_ARROW),
])
def test_the_shape_matches_the_direction(tmp_path, name, shape):
    prs = _build(tmp_path / name, _FM + "### x\n\n- 上\n\n" + _fence(name)
                 + "\n\n→ 下\n")
    assert len(_shapes(prs.slides[-1], shape)) == 1


def test_the_arrow_sits_between_the_prose(tmp_path):
    prs = _build(tmp_path, _FM + "### x\n\n- 上\n\n" + _fence("down") + "\n\n→ 下\n")
    slide = prs.slides[-1]
    arrow, = _shapes(slide, MSO_SHAPE.DOWN_ARROW)
    body, = [sh for sh in slide.shapes
             if sh.is_placeholder and sh.has_text_frame
             and sh != slide.shapes.title and "上" in sh.text_frame.text]
    assert body.top < arrow.top
    assert arrow.top + arrow.height <= body.top + body.height


def test_two_arrows_stack(tmp_path):
    """1枚に2つ置ける（cn2025-01 s28 は上下矢印と下矢印を続けて使う）．"""
    src = (_FM + "### x\n\n- A\n\n" + _fence("updown") + "\n\n- B\n\n"
           + _fence("down") + "\n\n→ C\n")
    prs = _build(tmp_path, src)
    slide = prs.slides[-1]
    ud, = _shapes(slide, MSO_SHAPE.UP_DOWN_ARROW)
    dn, = _shapes(slide, MSO_SHAPE.DOWN_ARROW)
    assert ud.top < dn.top


def test_it_works_inside_a_column(tmp_path):
    src = (_FM + "### x\n\n- 左\n\n<!-- @col -->\n\n- 右\n\n"
           + _fence("down") + "\n\n→ 結論\n")
    prs = _build(tmp_path, src)
    slide = prs.slides[-1]
    arrow, = _shapes(slide, MSO_SHAPE.DOWN_ARROW)
    right, = [sh for sh in slide.shapes
              if sh.is_placeholder and sh.placeholder_format.idx == 2]
    assert right.left <= arrow.left < right.left + right.width


# ---------------------------------------------------------------- @col: arrow

def test_the_column_break_arrow_still_works(tmp_path):
    """``<!-- @col: arrow -->`` はそのまま（役割が違うので残す）．"""
    src = _FM + "### x\n\n- 左\n\n<!-- @col: arrow -->\n\n- 右\n"
    prs = _build(tmp_path, src)
    assert len(_shapes(prs.slides[-1], MSO_SHAPE.RIGHT_ARROW)) == 1


# ---------------------------------------------------------------- 型と図形の対

def test_every_direction_has_a_shape():
    """``ArrowDirection`` を増やしたら、図形の対応表も足す合図になる．

    ``ARROW_DIRECTIONS`` が空になっていないことも、ここで一緒に見ている
    （``get_args`` は実行時に評価するので、型注釈の書き方によっては空になりうる）。
    """
    from md2pptx.ir import ARROW_DIRECTIONS
    assert len(ARROW_DIRECTIONS) == 6
    assert set(ARROW_DIRECTIONS) == set(render.Renderer._ARROW_SHAPES)


def test_a_horizontal_arrow_is_wider_than_tall(tmp_path):
    """横向きは長手が横．帯の高さではなく幅から長さを取る．"""
    prs = _build(tmp_path, _FM + "### x\n\n- 上\n\n" + _fence("right") + "\n\n→ 下\n")
    arrow, = _shapes(prs.slides[-1], MSO_SHAPE.RIGHT_ARROW)
    assert arrow.width > arrow.height


def test_a_vertical_arrow_is_taller_than_wide(tmp_path):
    prs = _build(tmp_path, _FM + "### x\n\n- 上\n\n" + _fence("down") + "\n\n→ 下\n")
    arrow, = _shapes(prs.slides[-1], MSO_SHAPE.DOWN_ARROW)
    assert arrow.height > arrow.width


def test_a_repeated_key_takes_the_last_value():
    """同じキーを 2 回書いたら後勝ち（``` ```image ``` と同じ扱い）．"""
    arrow, = [b for b in _blocks(
        "```arrow\ndirection: down\ndirection: up\n```") if isinstance(b, Arrow)]
    assert arrow.direction == "up"


def test_two_arrows_stay_readable(tmp_path):
    """帯を分け合っても、矢印が正方形（＝向きの読めない形）にならない．

    #139 で長手の比率を落としたとき、2 つ置いた版が 0.89×0.89cm になり
    ``updown`` が菱形に見えていた（Issue #141）。
    """
    src = (_FM + "### x\n\n- A\n  - a\n\n" + _fence("updown") + "\n\n- B\n  - b\n\n"
           + _fence("down") + "\n\n→ C\n")
    prs = _build(tmp_path, src)
    slide = prs.slides[-1]
    # ここで置くのは**縦向きだけ**なので「高さ＞幅」で見る．横向きを足すなら
    # 向きごとに条件を分けること（長手がどちらかは render が決める）．
    for kind in (MSO_SHAPE.UP_DOWN_ARROW, MSO_SHAPE.DOWN_ARROW):
        arrow, = _shapes(slide, kind)
        assert arrow.height > arrow.width * 1.2, f"{kind} が正方形に近い"
