#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""図の note(top) / note(bottom) を固定するテスト（Issue #129）．

``note(top)`` / ``note(bottom)`` は図の一部ではなく**地の文**で、本文
プレースホルダへ流し込まれる（SYNTAX.md「フロー図」「シーケンス図」）。
地の文である以上、本文行と**同じ**解釈を受けなければならない——
``→`` は行頭記号なし、行内装飾は装飾として効く。

実際には2つ壊れていた。

1. ``_note_to_line`` が ``_parse_spans`` を通らず、``[語]{red}`` が生の文字で出た
2. ``seq`` の note は分岐が ``Flow`` 限定で、**丸ごと落ちていた**

図の結論文は「→ …ので[返事が返れない]{red}」のように色が意味を持つ使い方を
するので、記号が生で出ると読めない。ここで両方を固定する。
"""
from __future__ import annotations

import pytest
from pptx import Presentation

from md2pptx import render
from md2pptx.parser import parse


def _build(tmp_path, src):
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _paragraphs(slide):
    """スライドの**タイトル以外**のテキストフレームの段落を (text, runs) で返す．

    タイトルを外すのは、見出しと同じ文字列を note に書いたときに取り違えないため．
    """
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame or sh == slide.shapes.title:
            continue
        for p in sh.text_frame.paragraphs:
            if p.text.strip():
                out.append((p.text, p.runs))
    return out


def _texts(slide):
    return [t for t, _ in _paragraphs(slide)]


def _find(slide, needle):
    for t, runs in _paragraphs(slide):
        if needle in t:
            return t, runs
    return None, None


_FLOW = """---
theme: t.pptx
---

### f

```flow
[#a 送信元] -> [#b 宛先]
note(top): {top}
note(bottom): {bottom}
```
"""

_SEQ = """---
theme: t.pptx
---

### s

```seq
lifelines: A, B
note(top): {top}
A -> B: x
note(bottom): {bottom}
```
"""


# ------------------------------------------------------- 行内装飾が効く

@pytest.mark.parametrize("src", [_FLOW, _SEQ], ids=["flow", "seq"])
def test_note_bottom_applies_inline_color(tmp_path, src):
    """``[語]{red}`` は色の付いた run になり、記号は本文に残らない．"""
    prs = _build(tmp_path, src.format(
        top="上", bottom="→ ここが[危険]{red}です"))
    slide = prs.slides[-1]
    text, runs = _find(slide, "危険")
    assert text == "→ ここが危険です"
    assert [r.text for r in runs] == ["→ ここが", "危険", "です"]
    assert runs[1].font.color.rgb == render.RGBColor(0xFF, 0x00, 0x00)


@pytest.mark.parametrize("src", [_FLOW, _SEQ], ids=["flow", "seq"])
def test_note_top_applies_inline_bold(tmp_path, src):
    """``**強調**`` は太字の run になる．"""
    prs = _build(tmp_path, src.format(top="**要点**はここ", bottom="下"))
    slide = prs.slides[-1]
    text, runs = _find(slide, "要点")
    assert text == "要点はここ"
    assert [(r.text, r.font.bold) for r in runs] == [
        ("要点", True), ("はここ", None)]


# ------------------------------------------------------- seq でも出る

def test_seq_notes_reach_the_body(tmp_path):
    """seq の note(top)/note(bottom) が落ちない（Flow 限定の分岐だった）．"""
    prs = _build(tmp_path, _SEQ.format(top="上の注記", bottom="→ 下の注記"))
    texts = _texts(prs.slides[-1])
    assert "上の注記" in texts
    assert "→ 下の注記" in texts


def test_seq_note_order_around_the_figure(tmp_path):
    """note(top) は図より前、note(bottom) は図より後の地の文になる．"""
    src = """---
theme: t.pptx
---

### s

導入文

```seq
lifelines: A, B
note(top): 上の注記
A -> B: x
note(bottom): → 下の注記
```

結論文
"""
    prs = _build(tmp_path, src)
    texts = [t for t in _texts(prs.slides[-1])
             if t in ("導入文", "上の注記", "→ 下の注記", "結論文")]
    assert texts == ["導入文", "上の注記", "→ 下の注記", "結論文"]


# ------------------------------------------------------- 従来の挙動

@pytest.mark.parametrize("src", [_FLOW, _SEQ], ids=["flow", "seq"])
def test_arrow_note_has_no_bullet(tmp_path, src):
    """``→`` で始まる note は行頭記号なし（従来どおり）．"""
    prs = _build(tmp_path, src.format(top="上", bottom="→ 結論"))
    slide = prs.slides[-1]
    text, _ = _find(slide, "結論")
    assert text == "→ 結論"
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            if p.text == "→ 結論":
                assert p._pPr.find(render.qn("a:buNone")) is not None
                return
    pytest.fail("→ の note が見つからない")


@pytest.mark.parametrize("src", [_FLOW, _SEQ], ids=["flow", "seq"])
def test_plain_note_stays_one_run(tmp_path, src):
    """装飾の無い note は 1 run のまま（従来経路を変えない）．"""
    prs = _build(tmp_path, src.format(top="ただの注記", bottom="下"))
    slide = prs.slides[-1]
    _, runs = _find(slide, "ただの注記")
    assert [r.text for r in runs] == ["ただの注記"]
