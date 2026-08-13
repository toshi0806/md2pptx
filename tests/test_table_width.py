#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""表の総幅を帯より狭くできることを固定するテスト（Issue #172）．

表は帯の幅いっぱいに描かれる。列が少なく短いセルばかりの表——「通信の種類」の
ような 2 列の対応表——では横に間延びして、内容に対して幅が広すぎる。
元にした cn2025 の表はスライド幅の 40〜80% で中央寄せされていた。

``@table-width`` で総幅を決め、余った帯は**左右に等分**して中央に置く。
``@table-widths``（列幅の比）とは別物で、こちらは総幅そのもの。
"""
from __future__ import annotations

import pytest
from pptx import Presentation

from md2pptx import render
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"

# 短いセルばかりの 2 列表（cn2026-04「通信の種類」と同じ形）
_SHORT = """| 種類 | 相手 |
|:--|:--|
| ユニキャスト | 1 対 1 |
| ブロードキャスト | 1 対 全 |
| マルチキャスト | 1 対 多 |
"""

# 帯に収まりきらない長いセルの表
_LONG = """| 種類 | 説明 |
|:--|:--|
| ユニキャスト | 相手を1つ指定して送る通常の通信で、宛先は必ず単一のノードになる |
| ブロードキャスト | 同じネットワーク上の全員に届き、ルータは原則として先へ転送しない |
"""


def _theme_path(tmp_path):
    """このテスト群が使う素のテーマ（既定レイアウトのまま）．"""
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    if not theme.exists():
        Presentation().save(str(theme))
    return theme


def _build(tmp_path, src):
    theme = _theme_path(tmp_path)
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _frame(prs):
    """最後のスライドの表シェイプ（GraphicFrame）を返す．"""
    for sh in prs.slides[-1].shapes:
        if sh.has_table:
            return sh
    raise AssertionError("表が無い")


def _geometry(prs) -> tuple[int, int]:
    """表の (左端, 総幅)．総幅は列幅の合計で測る（枠の値と食い違わせない）．"""
    fr = _frame(prs)
    return fr.left, sum(c.width for c in fr.table.columns)


def _src(rows: str, directive: str = "") -> str:
    head = "### 通信の種類\n" + (directive + "\n" if directive else "") + "\n"
    return _FM + head + rows


# ---------------------------------------------------------------- 既定（回帰）

def test_without_the_directive_the_table_fills_the_band(tmp_path):
    """未指定のときは帯いっぱい——既存のデッキの見た目を変えない．"""
    band_left, band_w = _geometry(_build(tmp_path, _src(_SHORT)))
    assert band_w > 0
    # 帯いっぱいなので、幅を絞ったときより広く、左端はより左にある
    left, w = _geometry(_build(tmp_path / "p", _src(_SHORT, "<!-- @table-width: 60% -->")))
    assert w < band_w
    assert left > band_left


# ---------------------------------------------------------------- 割合指定

def test_a_percent_narrows_the_table(tmp_path):
    """``60%`` で総幅が帯幅の 60% になる．"""
    _, band_w = _geometry(_build(tmp_path, _src(_SHORT)))
    _, w = _geometry(_build(tmp_path / "p", _src(_SHORT, "<!-- @table-width: 60% -->")))
    assert w == pytest.approx(band_w * 0.6, rel=0.02)


def test_a_narrowed_table_is_centred_in_the_band(tmp_path):
    """絞ったぶんは左右へ等分する（帯の中で中央寄せ）．"""
    band_left, band_w = _geometry(_build(tmp_path, _src(_SHORT)))
    left, w = _geometry(_build(tmp_path / "p", _src(_SHORT, "<!-- @table-width: 60% -->")))
    assert left == pytest.approx(band_left + (band_w - w) / 2, rel=0.02)


def test_a_percent_over_100_is_capped(tmp_path):
    """帯より広くはしない（はみ出させない）．"""
    _, band_w = _geometry(_build(tmp_path, _src(_SHORT)))
    _, w = _geometry(_build(tmp_path / "p", _src(_SHORT, "<!-- @table-width: 150% -->")))
    assert w == band_w


# ---------------------------------------------------------------- auto

def test_auto_fits_short_cells(tmp_path):
    """``auto`` は短いセルばかりの表を内容なりの幅に収める．"""
    _, band_w = _geometry(_build(tmp_path, _src(_SHORT)))
    _, w = _geometry(_build(tmp_path / "p", _src(_SHORT, "<!-- @table-width: auto -->")))
    assert w < band_w


def test_auto_is_capped_by_the_band(tmp_path):
    """``auto`` が要る幅が帯を超えるときは帯幅で頭打ちにする．"""
    _, band_w = _geometry(_build(tmp_path, _src(_LONG)))
    _, w = _geometry(_build(tmp_path / "p", _src(_LONG, "<!-- @table-width: auto -->")))
    assert w == band_w


def test_auto_gives_each_column_its_own_width(tmp_path):
    """``auto`` は列ごとの必要幅で割る——総幅だけ合わせて等分しない．

    「ブロードキャスト」の列と「1 対 全」の列を同じ幅にしては、
    横の間延びを別の場所へ移しただけになる．
    """
    prs = _build(tmp_path, _src(_SHORT, "<!-- @table-width: auto -->"))
    cols = [c.width for c in _frame(prs).table.columns]
    assert cols[0] > cols[1]


def test_explicit_ratios_win_over_auto(tmp_path):
    """``@table-widths`` を書いてあれば、列の配分はそちらに従う．"""
    src = _src(_SHORT, "<!-- @table-width: auto -->\n<!-- @table-widths: 50,50 -->")
    cols = [c.width for c in _frame(_build(tmp_path, src)).table.columns]
    assert cols[0] == pytest.approx(cols[1], rel=0.02)


def test_auto_leaves_room_for_the_longest_cell(tmp_path):
    """``auto`` の幅は最長セルが 1 行に収まるだけある（潰さない）．"""
    r = render.Renderer(str(_theme_path(tmp_path)))
    _, w = _geometry(_build(tmp_path, _src(_SHORT, "<!-- @table-width: auto -->")))
    need = r._text_width_pt("ブロードキャスト", r._body_font_size())
    assert w / 12700.0 > need          # 列 1 つぶんより広い（左右の余白を別にしても）


# ---------------------------------------------------------------- 併用

def test_it_combines_with_table_widths(tmp_path):
    """``@table-widths`` の比は、絞った総幅の中で配分される．"""
    src = _src(_SHORT, "<!-- @table-width: 60% -->\n<!-- @table-widths: 30,70 -->")
    prs = _build(tmp_path, src)
    cols = [c.width for c in _frame(prs).table.columns]
    total = sum(cols)
    assert cols[0] == pytest.approx(total * 0.3, rel=0.02)
    _, band_w = _geometry(_build(tmp_path / "p", _src(_SHORT)))
    assert total == pytest.approx(band_w * 0.6, rel=0.02)


# ---------------------------------------------------------------- 誤り

@pytest.mark.parametrize("bad", ["ひろめ", "70", "-20%", "0%"])
def test_a_bad_value_stops_with_the_line_number(tmp_path, bad):
    """解釈できない値は黙って無視せず、行番号を添えてエラーにする．"""
    with pytest.raises(ValueError) as e:
        _build(tmp_path, _src(_SHORT, f"<!-- @table-width: {bad} -->"))
    assert "@table-width" in str(e.value)
    assert "line" in str(e.value)
