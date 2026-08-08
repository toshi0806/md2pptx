#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""折り返した ``{box}`` の高さを固定するテスト（Issue #150）．

2行に折り返す項目を囲むと、**枠が1行ぶん下へ伸びて次の項目に掛かっていた**。
枠の高さを「折り返し行数 × 段落の高さ」で出していて、その段落の高さには
#147 で入れた段落前アキ（``spcBef``）が入っていたため。

アキは**段落に1回**付くもので、折り返した行ごとには付かない。cn2026-theme は
全レベル 20% なので、2行の項目は 0.2 行ぶん余計に伸びていた。
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Pt

from md2pptx import render
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"


def _theme(tmp_path, spc_pct=20000, sizes=(3000, 2600, 2200)):
    """マスター本文スタイルに spcBef を持つテーマを作る．"""
    tmp_path.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    body = prs.slide_masters[0].element.find(
        qn("p:txStyles") + "/" + qn("p:bodyStyle"))
    for lvl, sz in enumerate(sizes, start=1):
        el = body.find(qn(f"a:lvl{lvl}pPr"))
        assert el is not None, f"既定テーマに lvl{lvl}pPr が無い（前提が崩れている）"
        for tag in ("a:spcBef", "a:defRPr"):
            old = el.find(qn(tag))
            if old is not None:
                el.remove(old)
        if spc_pct:
            spc = el.makeelement(qn("a:spcBef"), {})
            spc.append(spc.makeelement(qn("a:spcPct"), {"val": str(spc_pct)}))
            el.insert(0, spc)
        el.append(el.makeelement(qn("a:defRPr"), {"sz": str(sz)}))
    path = tmp_path / "theme.pptx"
    prs.save(str(path))
    return path


def _build(tmp_path, src, theme):
    tmp_path.mkdir(parents=True, exist_ok=True)
    r = render.Renderer(str(theme))
    r.render(parse(src))
    out = tmp_path / "out.pptx"
    r.save(str(out))
    return Presentation(str(out)), r


def _boxes(prs):
    return [sh for sh in prs.slides[-1].shapes
            if not sh.is_placeholder
            and getattr(sh, "auto_shape_type", None) == MSO_SHAPE.ROUNDED_RECTANGLE]


# 既定テンプレートの本文枠は 8.5in ≈ 21.6cm．lvl1 は 30pt（≈1.06cm/全角）なので
# 20 文字ほどで折り返す．24 文字ならちょうど 2 行になる。
_LONG = "とても長い項目名" * 3


def _wrapped(r, prs, text, level=0):
    """``text`` が何行に折り返るか——md2pptx 自身の見積もりで数える．

    テストが確かめたいのは**高さの計算**なので、折り返し行数のほうは前提として
    確認するだけにする。テーマや python-pptx が変わって 2 行でなくなったとき、
    「バグが出た」ではなく「テストの前提が崩れた」と分かるようにするため。
    """
    body, = [sh for sh in prs.slides[-1].shapes
             if sh.is_placeholder and sh.placeholder_format.idx == 1]
    tf = body.text_frame
    avail_pt = (body.width - tf.margin_left - tf.margin_right) / 12700.0
    sz = r._body_font_levels()[min(level, len(r._body_font_levels()) - 1)]
    import math
    return max(1, math.ceil(r._text_width_pt(text, sz) / avail_pt))


def test_a_wrapped_box_is_two_lines_tall_not_more(tmp_path):
    """2行に折り返す項目の枠は、1行の枠のちょうど2倍にはならない．

    アキは段落に1回きりなので、``2*行高 + アキ`` になる。
    ``2*(行高+アキ)`` だとアキ1つぶん背が高い。
    """
    theme = _theme(tmp_path)
    one, r = _build(tmp_path / "a", _FM + "### x\n\n- {box} 短い\n", theme)
    two, _ = _build(tmp_path / "b", _FM + "### x\n\n- {box} " + _LONG + "\n", theme)
    (a,), (b,) = _boxes(one), _boxes(two)
    assert _wrapped(r, two, _LONG) == 2, "前提が崩れた（_LONG が2行に折り返らない）"
    sz = r._body_font_levels()[0]
    line = r._line_height(sz)
    assert a.height == pytest.approx(line, abs=Pt(1))
    assert b.height == pytest.approx(2 * line, abs=Pt(1))


def test_the_next_box_is_not_pushed_down(tmp_path):
    """折り返す項目の**下**にある枠も、そのぶんずれない．"""
    theme = _theme(tmp_path)
    prs, r = _build(tmp_path, _FM + "### x\n\n- " + _LONG + "\n- {box} 次の項目\n",
                    theme)
    assert _wrapped(r, prs, _LONG) == 2, "前提が崩れた（_LONG が2行に折り返らない）"
    box, = _boxes(prs)
    sz = r._body_font_levels()[0]
    body, = [sh for sh in prs.slides[-1].shapes
             if sh.is_placeholder and sh.placeholder_format.idx == 1]
    top = (body.top + body.text_frame.margin_top
           + r._para_height(0, sz, 2)      # 折り返して2行になった上の項目
           + r._space_before(0, sz)        # 囲む項目そのものの段落前アキ
           - int(r._line_height(sz) * r._BOX_LIFT))   # 字に合わせた持ち上げ
    assert box.top == pytest.approx(top, abs=Pt(2))


def test_a_theme_without_space_before_is_unchanged(tmp_path):
    """アキの無いテーマでは、折り返した枠はちょうど2行ぶん（回帰させない）．"""
    theme = _theme(tmp_path, spc_pct=0)
    one, _ = _build(tmp_path / "a", _FM + "### x\n\n- {box} 短い\n", theme)
    two, _r = _build(tmp_path / "b", _FM + "### x\n\n- {box} " + _LONG + "\n", theme)
    (a,), (b,) = _boxes(one), _boxes(two)
    assert _wrapped(_r, two, _LONG) == 2, "前提が崩れた（_LONG が2行に折り返らない）"
    assert b.height == pytest.approx(2 * a.height, abs=Pt(1))



def test_the_pieces_add_up_to_the_paragraph_height(tmp_path):
    """``_para_height`` は「行の高さ ＋ アキ」に一致する（分けても意味は同じ）．"""
    theme = _theme(tmp_path)
    _, r = _build(tmp_path, _FM + "### x\n\n- {box} 短い\n", theme)
    for lvl, sz in ((0, 30.0), (1, 26.0), (2, 22.0)):
        assert r._para_height(lvl, sz) == pytest.approx(
            r._line_height(sz) + r._space_before(lvl, sz), abs=2)
