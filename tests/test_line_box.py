#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""行を枠で囲む ``{box}`` を固定するテスト（Issue #133）．

シラバスや講義日程で「その日にあたる項目だけを枠で囲む」——cn2025 の全14ファイルが
やっている使い方で、囲むのは **run ではなく段落**（2行に折り返しても枠は1つ）。

トークンの置き場は**相対サイズトークンと同じ**（行頭マーカーの直後）。
新しい構文の場所を増やさないため。
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

from md2pptx import render
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"


def _lines(body):
    slide, = parse(_FM + "### x\n\n" + body + "\n").slides
    return slide.blocks


def _build(tmp_path, src):
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _boxes(slide):
    """描かれた枠（塗りつぶし無しの角丸四角）を返す．"""
    return [sh for sh in slide.shapes
            if not sh.is_placeholder
            and getattr(sh, "auto_shape_type", None) == MSO_SHAPE.ROUNDED_RECTANGLE]


def _body(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh != slide.shapes.title and sh.is_placeholder:
            return sh
    return None


# ---------------------------------------------------------------- 記法

def test_bullet_can_be_boxed():
    ln, = _lines("- {box} 囲む行")
    assert (ln.text, ln.kind, ln.boxed) == ("囲む行", "bullet", True)


def test_autonum_can_be_boxed():
    ln, = _lines("1. {box} 囲む行")
    assert (ln.text, ln.kind, ln.boxed, ln.num_start) == ("囲む行", "autonum", True, 1)


def test_arrow_line_can_be_boxed():
    ln, = _lines("→ {box} 囲む行")
    assert (ln.text, ln.kind, ln.boxed) == ("→ 囲む行", "plain", True)


def test_a_plain_line_is_not_boxed():
    ln, = _lines("- ふつうの行")
    assert ln.boxed is False


@pytest.mark.parametrize("src", ["- {box}{-1} 囲む行", "- {-1}{box} 囲む行"],
                         ids=["box-first", "size-first"])
def test_box_and_size_token_combine_in_either_order(src):
    """サイズトークンと併記できる．順序は問わない．"""
    ln, = _lines(src)
    assert (ln.text, ln.boxed, ln.size_delta) == ("囲む行", True, -1)


def test_an_unknown_token_stays_text():
    """``{boxx}`` はトークンではないので本文に残る（``{+x}`` と同じ扱い）．"""
    ln, = _lines("- {boxx} 本文")
    assert ln.text == "{boxx} 本文"
    assert ln.boxed is False


def test_the_token_is_removed_from_the_text():
    """トークンは本文に残さない（枠は見た目であって内容ではない）．"""
    ln, = _lines("- {box} 囲む行")
    assert "{box}" not in ln.text


# ---------------------------------------------------------------- 描画

def test_a_boxed_line_draws_a_frame(tmp_path):
    prs = _build(tmp_path, _FM + "### x\n\n- {box} 囲む行\n- ふつうの行\n")
    box, = _boxes(prs.slides[-1])
    assert box.fill.type is None or box.fill.type == 5   # 塗りつぶし無し


def test_no_frame_without_the_token(tmp_path):
    """``{box}`` の無い原稿の出力は変わらない．"""
    prs = _build(tmp_path, _FM + "### x\n\n- ふつうの行\n- もう1行\n")
    assert _boxes(prs.slides[-1]) == []


def test_the_frame_sits_inside_the_body(tmp_path):
    """短い 1 行なら枠は本文の枠に収まる．

    見積もりなので、枠いっぱいまで折り返す行では下端を超えうる。ここで固定して
    いるのは「ふつうの行で枠が飛び出さない」ことで、上限そのものではない。
    """
    prs = _build(tmp_path, _FM + "### x\n\n- {box} 囲む行\n")
    slide = prs.slides[-1]
    box, = _boxes(slide)
    body = _body(slide)
    assert body.left <= box.left
    assert box.left + box.width <= body.left + body.width
    assert body.top <= box.top
    assert box.top + box.height <= body.top + body.height


def test_the_frame_follows_the_paragraph_down(tmp_path):
    """3 行目を囲めば、1 行目を囲んだときより下に出る．"""
    first = _build(tmp_path / "a", _FM + "### x\n\n- {box} A\n- B\n- C\n")
    third = _build(tmp_path / "b", _FM + "### x\n\n- A\n- B\n- {box} C\n")
    a, = _boxes(first.slides[-1])
    c, = _boxes(third.slides[-1])
    assert c.top > a.top


def test_the_frame_lands_in_the_right_column(tmp_path):
    """2カラムの右側の行を囲むと、枠も右カラムに出る．"""
    src = _FM + "### x\n\n- 左\n\n<!-- @col -->\n\n- {box} 右\n"
    prs = _build(tmp_path, src)
    slide = prs.slides[-1]
    box, = _boxes(slide)
    # 右カラムは idx==2 のプレースホルダ．座標の最大で選ぶと
    # スライド番号のプレースホルダを拾ってしまう．
    right, = [sh for sh in slide.shapes
              if sh.is_placeholder and sh.placeholder_format.idx == 2]
    assert right.left <= box.left < right.left + right.width


def test_the_frame_works_next_to_a_table(tmp_path):
    """表と同居するスライド（帯の上下に地の文が割れる）でも位置が合う．"""
    src = (_FM + "### x\n\n- {box} 導入文\n\n| a | b |\n|:--|:--|\n| 1 | 2 |\n\n"
           "→ 結論文\n")
    prs = _build(tmp_path, src)
    slide = prs.slides[-1]
    box, = _boxes(slide)
    body = _body(slide)
    # 導入文は枠の 1 行目なので、帯（表）より上にある
    tbl = [sh for sh in slide.shapes if sh.has_table][0]
    assert box.top < tbl.top


def test_a_wrapped_line_gets_a_taller_frame(tmp_path):
    """折り返す行は、枠も 2 行ぶんの高さになる（枠は 1 つのまま）．"""
    short = _build(tmp_path / "a", _FM + "### x\n\n- {box} 短い\n")
    long_ = _build(tmp_path / "b",
                   _FM + "### x\n\n- {box} " + "とても長い項目名" * 8 + "\n")
    s, = _boxes(short.slides[-1])
    l, = _boxes(long_.slides[-1])
    assert l.height > s.height


# ---------------------------------------------------------------- 枠の色

def test_the_frame_takes_a_colour():
    """色は**正規化して**持つ（Span.color と同じ扱い．"#f00" と "#F00" を分けない）．"""
    ln, = _lines("- {box:blue} 囲む行")
    assert (ln.boxed, ln.box_color, ln.text) == (True, "#0000FF", "囲む行")


def test_a_theme_colour_stays_a_name():
    """テーマ色名は RGB へ潰さない（テーマ差し替えに追従させるため）．"""
    ln, = _lines("- {box:accent1} 囲む行")
    assert ln.box_color == "accent1"


def test_hex_case_is_normalised():
    ln, = _lines("- {box:#f00} 囲む行")
    assert ln.box_color == "#FF0000"


def test_a_misspelt_colour_stops(tmp_path):
    """色名の綴り違いはエラーで止まる（行内装飾と同じ扱い）．"""
    with pytest.raises(Exception):
        _lines("- {box:bleu} 囲む行")


def test_the_colour_reaches_the_shape(tmp_path):
    prs = _build(tmp_path, _FM + "### x\n\n- {box:blue} 囲む行\n")
    box, = _boxes(prs.slides[-1])
    assert str(box.line.color.rgb) == "0000FF"


def test_without_a_colour_the_theme_decides(tmp_path):
    """色を書かなければテーマ色（差し替えに追従する）．"""
    prs = _build(tmp_path, _FM + "### x\n\n- {box} 囲む行\n")
    box, = _boxes(prs.slides[-1])
    assert box.line.color.type is not None
    assert box.line.color.theme_color is not None


# ---------------------------------------------------------------- 端の扱い

def test_the_frame_never_leaves_the_slide(tmp_path):
    """左へ余白を取っても、スライドの外（負の座標）へは出さない．"""
    prs = _build(tmp_path, _FM + "### x\n\n- {box} 囲む行\n")
    box, = _boxes(prs.slides[-1])
    assert box.left >= 0
    assert box.left + box.width <= prs.slide_width


def test_a_size_token_above_shifts_the_frame(tmp_path):
    """先行する行が大きいと、その下の枠もそのぶん下がる．

    先行ぶんを一律の本文標準サイズで数えていると、ここがずれる。
    """
    plain = _build(tmp_path / "a", _FM + "### x\n\n- 上の行\n- {box} 囲む行\n")
    big = _build(tmp_path / "b", _FM + "### x\n\n- {+3} 上の行\n- {box} 囲む行\n")
    a, = _boxes(plain.slides[-1])
    b, = _boxes(big.slides[-1])
    assert b.top > a.top


# ------------------------------------------------ 字に対する上下（Issue #160）

def test_the_frame_is_lifted_off_the_line_box(tmp_path):
    """枠は行の箱より少し上に出す．

    PowerPoint は行の箱の中で字を上寄りに置き、下に descent ぶんの空きを残す。
    行の箱にそのまま合わせると、字に対して枠が下がって見える——**枠線（3pt）を
    字と数えないように色で分けて**測ると、30pt の行で上 2.9pt / 下 8.6pt だった
    （Issue #164）。
    """
    prs = _build(tmp_path, _FM + "### x\n\n- {box} 囲む行\n")
    slide = prs.slides[-1]
    box, = _boxes(slide)
    body = _body(slide)
    r = render.Renderer(str(tmp_path / "theme.pptx"))
    sz = r._body_font_size()
    line_top = (body.top + body.text_frame.margin_top
                + r._space_before(0, sz))
    assert box.top < line_top                      # 持ち上がっている
    assert line_top - box.top < r._line_height(sz) // 4   # 上げすぎない
    # 字の高さ（行の箱の 7 割ほど）の半分は超えない——超えると今度は上へ寄る
    assert line_top - box.top < r._line_height(sz) * 0.35


def test_the_frame_stays_inside_the_body_after_the_lift(tmp_path):
    """持ち上げても本文の枠から上へ出ない．

    ``max(0, ...)`` のクランプは**防御**——持ち上げは行の 8.6% で、本文枠は
    タイトルの下にあるので、実際のテーマで負になる道は無い。ここで見ているのは
    「持ち上げが本文枠を越えるほど大きくない」ことのほう。
    """
    prs = _build(tmp_path, _FM + "### x\n\n- {box} 囲む行\n")
    slide = prs.slides[-1]
    box, = _boxes(slide)
    assert box.top >= 0
    assert box.top >= _body(slide).top
