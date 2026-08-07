#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""``` ```arrow ``` の大きさと色を固定するテスト（Issue #143）．

自動の大きさは帯から決まり 1.0in で頭打ちになる。cn2025 の矢印を実測すると
**0.6×1.1cm から 12.3×0.9cm まで**あり、`1.5×7.6cm`（層をまたぐ上下矢印）は
いまの上限では書けない。色も 88 個中 37 個がアクセント色ではなかった。

語彙は ``` ```image ``` と共通（``_parse_length`` / ``parse_color``）。
新しい書き方を増やさないため。
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm, Emu, Inches

from md2pptx import render
from md2pptx.ir import Arrow
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"


def _blocks(body):
    slide, = parse(_FM + "### x\n\n" + body + "\n").slides
    return slide.blocks


def _arrow(body):
    return [b for b in _blocks(body) if isinstance(b, Arrow)][0]


def _build(tmp_path, body):
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(_FM + "### x\n\n- 上\n\n" + body + "\n\n→ 下\n"))
    r.save(str(out))
    return Presentation(str(out))


def _shape(slide, kind=MSO_SHAPE.DOWN_ARROW):
    got = [sh for sh in slide.shapes
           if not sh.is_placeholder
           and getattr(sh, "auto_shape_type", None) == kind]
    assert len(got) == 1, f"{kind} が {len(got)} 個"
    return got[0]


def _fence(**kw):
    body = "\n".join(f"{k}: {v}" for k, v in kw.items())
    return "```arrow\n" + body + "\n```"


# ---------------------------------------------------------------- 記法

def test_a_bare_arrow_has_no_size():
    """書かなければ従来どおり（帯から自動）．"""
    a = _arrow(_fence(direction="down"))
    assert (a.width, a.height, a.color) == (None, None, None)


@pytest.mark.parametrize("text,emu", [
    ("2.5cm", Cm(2.5)), ("1in", Inches(1)), ('1"', Inches(1)),
    ("72pt", Inches(1)),
])
def test_absolute_lengths(text, emu):
    a = _arrow(_fence(direction="down", height=text))
    assert a.height.unit == "emu"
    assert round(a.height.value) == round(float(emu))


def test_a_percentage_stays_relative():
    """``%`` は帯に対する割合（``` ```image ``` と同じ）．"""
    a = _arrow(_fence(direction="down", width="40%"))
    assert (a.width.unit, a.width.value) == ("percent", 40.0)


def test_only_one_side_can_be_given():
    a = _arrow(_fence(direction="down", height="3cm"))
    assert a.width is None and a.height is not None


def test_a_bad_length_stops():
    with pytest.raises(ValueError):
        _arrow(_fence(direction="down", width="ふとい"))


def test_the_colour_is_normalised():
    """色は検証して正規化して持つ（``{box:…}`` と同じ扱い）．"""
    assert _arrow(_fence(direction="down", color="blue")).color == "#0000FF"
    assert _arrow(_fence(direction="down", color="#f00")).color == "#FF0000"
    assert _arrow(_fence(direction="down", color="accent2")).color == "accent2"


def test_a_misspelt_colour_stops():
    with pytest.raises(ValueError):
        _arrow(_fence(direction="down", color="bleu"))


def test_an_unknown_key_still_stops():
    with pytest.raises(ValueError, match="arrow"):
        _arrow("```arrow\ndirection: down\nthickness: 3\n```")


# ---------------------------------------------------------------- 描画

def test_an_absolute_size_is_used_as_written(tmp_path):
    prs = _build(tmp_path, _fence(direction="down", width="2.5cm", height="2.1cm"))
    arrow = _shape(prs.slides[-1])
    assert abs(arrow.width - Cm(2.5)) < Emu(10000)
    assert abs(arrow.height - Cm(2.1)) < Emu(10000)


def test_an_explicit_size_may_exceed_the_automatic_cap(tmp_path):
    """自動のときの上限（1.0in）を超えて書ける．

    cn2025-02 s8 の「層をまたぐ上下矢印」は 1.5×7.6cm で、
    上限に縛られていると**そもそも書けない**。
    """
    prs = _build(tmp_path, _fence(direction="updown", width="1.5cm", height="7.6cm"))
    arrow = _shape(prs.slides[-1], MSO_SHAPE.UP_DOWN_ARROW)
    assert arrow.height > Inches(1.0)
    assert abs(arrow.height - Cm(7.6)) < Emu(10000)


def test_one_side_given_leaves_the_other_automatic(tmp_path):
    auto = _shape(_build(tmp_path / "a", _fence(direction="down")).slides[-1])
    half = _shape(_build(tmp_path / "b",
                         _fence(direction="down", width="4cm")).slides[-1])
    assert abs(half.width - Cm(4)) < Emu(10000)
    assert half.height == auto.height


def test_the_arrow_stays_centred(tmp_path):
    """大きさを変えても帯の中央に置く．"""
    prs = _build(tmp_path, _fence(direction="down", width="4cm"))
    slide = prs.slides[-1]
    arrow = _shape(slide)
    body, = [sh for sh in slide.shapes
             if sh.is_placeholder and sh.placeholder_format.idx == 1]
    centre = body.left + body.width // 2
    assert abs((arrow.left + arrow.width // 2) - centre) < Emu(20000)


def test_the_colour_reaches_the_shape(tmp_path):
    prs = _build(tmp_path, _fence(direction="down", color="blue"))
    assert str(_shape(prs.slides[-1]).fill.fore_color.rgb) == "0000FF"


def test_a_theme_colour_stays_a_theme_colour(tmp_path):
    prs = _build(tmp_path, _fence(direction="down", color="accent2"))
    assert _shape(prs.slides[-1]).fill.fore_color.theme_color is not None


def test_without_a_colour_the_theme_decides(tmp_path):
    prs = _build(tmp_path, _fence(direction="down"))
    assert _shape(prs.slides[-1]).fill.fore_color.theme_color is not None


# ---------------------------------------------------------------- 形の比率

def test_the_shape_carries_its_proportions(tmp_path):
    """矢じりと軸の比率を書き込む．

    書かないと PowerPoint の既定に落ち、両端に矢じりのある形は
    **箱いっぱいの菱形**になって矢印に見えない。
    """
    from pptx.oxml.ns import qn
    prs = _build(tmp_path, _fence(direction="updown", width="3.4cm",
                                  height="2.5cm"))
    shp = _shape(prs.slides[-1], MSO_SHAPE.UP_DOWN_ARROW)
    gds = shp._element.spPr.find(qn("a:prstGeom")).find(qn("a:avLst"))
    got = {gd.get("name"): gd.get("fmla") for gd in gds}
    assert got == {"adj1": "val 50000", "adj2": "val 20000"}


def test_a_single_headed_arrow_uses_a_bigger_head(tmp_path):
    """片側だけの矢じりは 25%（両側の 20% より大きく取れる）．"""
    from pptx.oxml.ns import qn
    prs = _build(tmp_path, _fence(direction="down"))
    shp = _shape(prs.slides[-1])
    gds = shp._element.spPr.find(qn("a:prstGeom")).find(qn("a:avLst"))
    got = {gd.get("name"): gd.get("fmla") for gd in gds}
    assert got == {"adj1": "val 50000", "adj2": "val 25000"}
