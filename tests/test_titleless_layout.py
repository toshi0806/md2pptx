#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""タイトルの無いレイアウトの置き場を固定するテスト（Issue #138）．

「白紙」レイアウト（``<!-- @layout: 6 -->``）で図だけのスライドを作ると、
**上にタイトルぶんの空きが残っていた**。`_content_rect` の既定値が
「タイトル下の本文相当領域」で、タイトルが無いレイアウトでもその 1.7in を空けていた。

cn2025 には「タイトルなしで図だけ」のスライドが十数枚あり、そこは図が
スライドいっぱいに置いてある。
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from md2pptx import render
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"


def _build(tmp_path, src):
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme), base_dir=str(tmp_path))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _fig(tmp_path, name="fig.png"):
    """1x1 の PNG を置いてパスを返す（中身は問わない）．"""
    from PIL import Image
    tmp_path.mkdir(parents=True, exist_ok=True)
    p = tmp_path / name
    Image.new("RGB", (400, 300), "white").save(p)
    return p


def _picture(slide):
    pics = [sh for sh in slide.shapes if not sh.is_placeholder]
    assert pics, "図が置かれていない（画像パスの解決に失敗した？）"
    return pics[0]


def _layouts(prs):
    """(タイトルも本文も無いレイアウト, タイトルだけあるレイアウト) の番号を返す．

    **名前では探さない**——``"Blank"`` / ``"Title Only"`` は python-pptx の既定
    テンプレートの英語名で、テーマや版が変われば変わる。ここで要るのは
    「タイトル枠があるか」「本文枠があるか」という**構造**そのものなので、
    プレースホルダを数えて選ぶ。
    """
    blank = titled = None
    for i, lay in enumerate(prs.slide_masters[0].slide_layouts):
        idxs = {ph.placeholder_format.idx for ph in lay.placeholders}
        has_title, has_body = 0 in idxs, 1 in idxs
        if not has_title and not has_body and blank is None:
            blank = i
        if has_title and not has_body and titled is None:
            titled = i
    assert blank is not None, "タイトルも本文も無いレイアウトが見つからない"
    assert titled is not None, "タイトルだけのレイアウトが見つからない"
    return blank, titled


def _src(layout, img):
    return (_FM + "### x\n<!-- @layout: %d -->\n\n"
            "```image\nsrc: %s\nwidth: 100%%\n```\n" % (layout, img.name))


def test_a_blank_layout_gives_the_figure_the_top(tmp_path):
    """「白紙」ではタイトルぶんを空けない．"""
    img = _fig(tmp_path)
    blank, _ = _layouts(Presentation())
    out = _build(tmp_path, _src(blank, img))
    pic = _picture(out.slides[-1])
    assert pic.top < Inches(1.0)


def test_a_titled_layout_still_reserves_the_top(tmp_path):
    """タイトルのあるレイアウトでは従来どおり空ける（回帰させない）．"""
    img = _fig(tmp_path / "b")
    _, only = _layouts(Presentation())
    out = _build(tmp_path / "b", _src(only, img))
    pic = _picture(out.slides[-1])
    assert pic.top >= Inches(1.0)


def test_the_figure_grows_on_a_blank_layout(tmp_path):
    """空きが減ったぶん、図は大きくなる．"""
    blank, only = _layouts(Presentation())

    img_a = _fig(tmp_path / "a")
    img_b = _fig(tmp_path / "b")
    on_blank = _picture(_build(tmp_path / "a", _src(blank, img_a)).slides[-1])
    on_title = _picture(_build(tmp_path / "b", _src(only, img_b)).slides[-1])
    assert on_blank.height > on_title.height


def test_the_figure_stays_on_the_slide(tmp_path):
    img = _fig(tmp_path)
    blank, _ = _layouts(Presentation())
    out = _build(tmp_path, _src(blank, img))
    slide = out.slides[-1]
    pic = _picture(slide)
    assert pic.top >= 0
    assert pic.left >= 0
    assert pic.top + pic.height <= out.slide_height
    assert pic.left + pic.width <= out.slide_width


def test_a_layout_with_a_body_is_unchanged(tmp_path):
    """本文プレースホルダのあるレイアウトはそもそもこの分岐へ来ない．

    図は本文の枠の中に置かれる——`_content_rect` はその枠を返し、
    既定値（この修正で触ったところ）まで落ちない。

    地の文を 1 行入れてあるのは、**図だけのスライドでは本文プレースホルダが
    消される**ため（地の文が無ければ枠ごと外して図を広く使う既存の動き）。
    ここで確かめたいのは枠がある場合なので、消えない形にしておく。
    """
    img = _fig(tmp_path)
    src = (_FM + "### x\n\n- 導入文\n\n```image\nsrc: %s\nwidth: 100%%\n```\n"
           % img.name)
    out = _build(tmp_path, src)
    slide = out.slides[-1]
    pic = _picture(slide)
    bodies = [sh for sh in slide.shapes
              if sh.is_placeholder and sh.placeholder_format.idx == 1]
    assert bodies, "既定レイアウトに本文プレースホルダが無い（前提が崩れている）"
    assert pic.top >= bodies[0].top


def test_the_bottom_margin_is_unchanged(tmp_path):
    """下の余白は 0.6in のまま（上を詰めた影響を下へ波及させない）．"""
    img = _fig(tmp_path)
    blank, _ = _layouts(Presentation())
    out = _build(tmp_path, _src(blank, img))
    pic = _picture(out.slides[-1])
    assert out.slide_height - (pic.top + pic.height) >= 0
    # 帯の下端は SH - 0.6in．図はその中に収まる．
    assert pic.top + pic.height <= out.slide_height - Inches(0.6) + 1
