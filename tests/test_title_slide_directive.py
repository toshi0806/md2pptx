#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""``@title-slide`` を固定するテスト（Issue #94）．

表紙はデッキに必ず 1 枚あるのに、指定に使う ``@layout`` は DESIGN.md §5.6 で
「**必要時のみ補う**」逃げ道として定義されている．しかも値がレイアウト番号で、
どのテーマでも 0 が表紙という並びの慣行を原稿が知っている必要があった．

``@title-slide`` は値を取らない（``@col`` と同じ形）．**``@layout`` との併記は
結果が同じ ``@layout: 0`` も含めてエラー**にする——「同じ結果なら許す」を入れると、
矛盾する組み合わせをどちらで解決するかを別に決めることになる．
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from md2pptx import render
from md2pptx.ir import TITLE_LAYOUT
from md2pptx.parser import parse

# @title-slide は syntax 0 専用（1 では "#" 自体が表紙）．
_FM = "---\ntheme: t.pptx\nsyntax: 0\n---\n\n"


def _theme(tmp_path):
    path = tmp_path / "theme.pptx"
    Presentation().save(str(path))
    return str(path)


def _build(tmp_path, src):
    out = tmp_path / "out.pptx"
    r = render.Renderer(_theme(tmp_path))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _numbered(slide):
    return any(ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER
               for ph in slide.placeholders)


def test_it_selects_the_title_layout():
    """``@title-slide`` は表紙レイアウトを選ぶ．"""
    slide = parse(_FM + "# 主題\n<!-- @title-slide -->\n\n- 著者\n").slides[0]
    assert slide.layout == TITLE_LAYOUT
    assert slide.directives == {"title_slide": True}


def test_the_title_slide_gets_no_number(tmp_path):
    """表紙には番号が付かず、扉と本文には付く．"""
    prs = _build(tmp_path, _FM + """# 表紙
<!-- @title-slide -->

- 著者

## 本文

- x

# 章扉
""")
    assert [(s.slide_layout.name, _numbered(s)) for s in prs.slides] == [
        ("Title Slide", False),
        ("Title and Content", True),
        ("Section Header", True),
    ]


@pytest.mark.parametrize("body", [
    "# 主題\n<!-- @title-slide -->\n<!-- @layout: 0 -->\n",   # 結果は同じでも
    "# 主題\n<!-- @layout: 0 -->\n<!-- @title-slide -->\n",   # 順序が逆でも
    "# 主題\n<!-- @title-slide -->\n<!-- @layout: 5 -->\n",   # 矛盾していても
])
def test_combining_with_layout_is_an_error(body):
    """``@layout`` との併記はエラー（結果が同じ組み合わせも含む）．"""
    with pytest.raises(ValueError, match="conflicts"):
        parse(_FM + body)


def test_it_takes_no_value():
    """値付きはエラー（``@col`` と同じ扱い）．"""
    with pytest.raises(ValueError, match="takes no value"):
        parse(_FM + "# 主題\n<!-- @title-slide: 0 -->\n")


def test_layout_still_works_on_its_own(tmp_path):
    """``@layout`` は従来どおり動く（表紙の書き方として案内しないだけ）．"""
    prs = _build(tmp_path, _FM + "# 主題\n<!-- @layout: 0 -->\n")
    assert prs.slides[0].slide_layout.name == "Title Slide"
    assert _numbered(prs.slides[0]) is False


def test_the_unknown_directive_message_lists_it():
    """未知のディレクティブの案内に ``@title-slide`` が並ぶ．

    タイポしたときに一覧から探せることが、このエラーの役目．
    """
    with pytest.raises(ValueError, match="@title-slide"):
        parse(_FM + "## 見出し\n<!-- @typo: 1 -->\n")
