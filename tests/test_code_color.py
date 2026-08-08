#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""等幅ブロックの中の色付けを固定するテスト（Issue #162）．

フェンスの中は**書いたまま**出すのが既定。``[`` や ``{`` はコードにふつうに
現れるので、無条件に解釈すると壊れる。そこで **info string でのオプトイン**にする。

    ```text color
     10000101  [10010100]{blue}
    ```

講義スライドでは IP アドレスのビット列を「ネットワーク部は青／ホスト部は赤」で
塗り分けている（cn2025-03）。桁を揃えるので等幅ブロックで書くしかない。
"""
from __future__ import annotations

import pytest
from pptx import Presentation

from md2pptx import render
from md2pptx.ir import Line
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"


def _lines(body):
    slide, = parse(_FM + "### x\n\n" + body + "\n").slides
    return [b for b in slide.blocks if isinstance(b, Line)]


def _fence(info, *body):
    return "```" + info + "\n" + "\n".join(body) + "\n```"


def _build(tmp_path, src):
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    r = render.Renderer(str(theme))
    r.render(parse(src))
    out = tmp_path / "out.pptx"
    r.save(str(out))
    return Presentation(str(out))


# ---------------------------------------------------------------- オプトイン

def test_a_plain_fence_keeps_the_markup_as_text():
    """``color`` を書かないフェンスは従来どおり——記号もそのまま出る．"""
    ln, = _lines(_fence("text", "[10010100]{blue}"))
    assert ln.text == "[10010100]{blue}"
    assert ln.spans == []


def test_the_color_info_turns_it_on():
    ln, = _lines(_fence("text color", "[10010100]{blue}"))
    assert ln.text == "10010100"
    assert [s.color for s in ln.spans] == ["#0000FF"]


def test_the_line_is_still_code():
    """色を付けても等幅の行のまま（行頭記号は消え、桁は揃う）．"""
    ln, = _lines(_fence("text color", "[10010100]{blue}"))
    assert ln.kind == "code"


def test_the_info_word_can_stand_alone():
    ln, = _lines(_fence("color", "[a]{red}"))
    assert ln.text == "a"


# ------------------------------------------------------------------ 桁揃え

def test_the_markup_does_not_take_columns():
    """記号は桁に数えない——等幅の桁揃えが崩れない．"""
    a, b = _lines(_fence("text color",
                         " 10000101  10010100  00000001",
                         " 10000101  [10010100]{blue}  [00000001]{red}"))
    assert a.text == b.text


def test_several_colors_on_one_line():
    ln, = _lines(_fence("text color", "[aa]{blue}bb[cc]{red}"))
    assert ln.text == "aabbcc"
    assert [(s.text, s.color) for s in ln.spans] == [
        ("aa", "#0000FF"), ("bb", None), ("cc", "#FF0000")]


# ------------------------------------------------- 色でない ``{…}`` は壊さない

def test_a_bracket_that_is_not_a_colour_stays_text():
    """``[133.69.130.4]{n}`` のようなふつうの角括弧は文字のまま．"""
    ln, = _lines(_fence("text color", "from mx [133.69.130.4]{n}"))
    assert ln.text == "from mx [133.69.130.4]{n}"
    assert ln.spans == []


def test_a_bare_bracket_stays_text():
    ln, = _lines(_fence("text color", "for <example@smkwlab.net>; [...]"))
    assert ln.text == "for <example@smkwlab.net>; [...]"


def test_bold_and_backticks_are_left_alone():
    """効くのは色だけ——もう等幅なので他の装飾は解釈しない．"""
    ln, = _lines(_fence("text color", "**not bold** `not code`"))
    assert ln.text == "**not bold** `not code`"


# ---------------------------------------------------------------- 描画

def test_the_colour_reaches_the_run(tmp_path):
    prs = _build(tmp_path, _FM + "### x\n\n"
                 + _fence("text color", "aa[bb]{blue}") + "\n")
    body, = [sh for sh in prs.slides[-1].shapes
             if sh.is_placeholder and sh.placeholder_format.idx == 1]
    runs = body.text_frame.paragraphs[0].runs
    assert [r.text for r in runs] == ["aa", "bb"]
    assert str(runs[1].font.color.rgb) == "0000FF"


def test_every_run_stays_monospace(tmp_path):
    """色を付けた run も等幅のまま——ここが崩れると桁が揃わない．"""
    prs = _build(tmp_path, _FM + "### x\n\n"
                 + _fence("text color", "aa[bb]{blue}") + "\n")
    body, = [sh for sh in prs.slides[-1].shapes
             if sh.is_placeholder and sh.placeholder_format.idx == 1]
    names = {r.font.name for r in body.text_frame.paragraphs[0].runs}
    assert len(names) == 1 and None not in names


# ------------------------------------ 色でない ``{…}`` と色が混ざる並び（#163）

@pytest.mark.parametrize("src,plain,colors", [
    # 色でないものが後ろ
    ("[a]{blue}[b]{n}", "a[b]{n}", ["#0000FF", None]),
    # 色でないものが前
    ("[b]{n}[a]{blue}", "[b]{n}a", [None, "#0000FF"]),
    # 色でないものが間に挟まる
    ("[a]{blue}[b]{n}[c]{red}", "a[b]{n}c", ["#0000FF", None, "#FF0000"]),
    # 色でないものが両端
    ("[x]{n}[a]{blue}[y]{m}", "[x]{n}a[y]{m}", [None, "#0000FF", None]),
])
def test_a_skipped_match_keeps_its_text(src, plain, colors):
    """色名でない ``{…}`` は**1文字も落とさず**に残る．

    飛ばしたぶんは「次のマッチまでの地の文」または末尾で拾われる。並びによって
    拾われ方が変わるので、代表的な4通りを固定しておく。
    """
    ln, = _lines(_fence("text color", src))
    assert ln.text == plain
    assert [s.color for s in ln.spans] == colors


def test_nothing_is_lost_whatever_the_order():
    """記号を除いた文字列は、色の付いた語と地の文の連結に必ず一致する．"""
    src = "pre[a]{blue}mid[b]{n}post[c]{red}end"
    ln, = _lines(_fence("text color", src))
    assert "".join(s.text for s in ln.spans) == ln.text
    assert ln.text == "prea" + "mid[b]{n}post" + "c" + "end"
