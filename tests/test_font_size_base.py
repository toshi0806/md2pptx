#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""相対フォントサイズ（``{±n}``）の基点を固定するテスト（Issue #83）．

``{±n}`` は「テーマ既定サイズからの相対段数」なので，**基点がその文字の実際に
出るサイズでなければ意味を持たない**．PowerPoint の継承は
スライド → レイアウト → マスターの順で，テーマはレイアウトのプレースホルダに
``a:lstStyle`` を置いてサイズを上書きする（既定テンプレートでも「セクション見出し」
「2 つのコンテンツ」がそうしている）．マスターの ``txStyles`` だけを見ていると，
**画面に出ているサイズとは違う値を基点にして縮小・拡大する**ことになる．

Issue #83 で実測した症状は 2 つで，ここで固定するのは①．

①同じ ``{-2}`` が，同じ枠なのに経路で違うサイズになっていた（front matter 由来の
  著者・所属は 22pt，本文行は 24pt）．基点が前者はレイアウト，後者はマスターだった．
②front matter の ``subtitle:`` の基点がマスターの表題サイズ×0.8 で，実際に描かれる
  タイトルサイズを見ていない．こちらは**直さず凍結**した——Issue #82 で
  front matter の表紙記述ごと本文記法へ移して非推奨にするため，消える直前に
  見た目を変える益がない（判断の経緯は ``render_title_slide`` のコメント）．

テーマは python-pptx 同梱の既定テンプレートから作るので，リポジトリの
``OfficeTheme.pptx`` にも実 PowerPoint にも依存しない．既定テンプレートで idx 1 の枠が
持つ実効 lvl1 はこうなっている（実測）．

- レイアウト 0「Title Slide」… 上書きなし → マスターの 32pt
- レイアウト 1「Title and Content」… 上書きなし → マスターの 32pt
- レイアウト 2「Section Header」… ``{1:20, 2:18, 3:16}`` → 20pt
- レイアウト 3「Two Content」… ``{1:28, 2:24, 3:20}`` → 28pt

**レイアウト 0 に sz の上書きが無いこと**が
``test_same_frame_same_size_across_paths`` の期待値の前提．``lstStyle`` 要素自体は
9 レベル分あるが，持っているのは配置とビュレットだけで ``sz`` を含まない．
"""
from __future__ import annotations

from pptx import Presentation

from md2pptx import render
from md2pptx.parser import parse

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# 既定テンプレートのマスター本文 lvl1（このテストの「上書きが無いときの基点」）．
_MASTER_BODY_LVL1 = 32.0
_STEP = render.Renderer._SIZE_STEP_RATIO


def _theme(tmp_path):
    path = tmp_path / "theme.pptx"
    Presentation().save(str(path))
    return str(path)


def _build(tmp_path, src, name="out.pptx"):
    """Markdown を描画して保存し，読み直した Presentation を返す．

    ``r.prs`` ではなく保存した pptx を読み直すのは，利用者が開くのがファイルの
    ほうだから（test_text_language.py と同じ理由）．
    """
    out = tmp_path / name
    deck = parse(src)
    r = render.Renderer(_theme(tmp_path))
    r.render(deck)
    r.save(str(out))
    return Presentation(str(out))


def _para_sizes(prs, slide_idx=0, ph_idx=1):
    """指定スライドの指定プレースホルダの (本文, 段落サイズ pt) を並べて返す．

    サイズは段落の既定文字書式（``a:pPr/a:defRPr@sz``）から読む．``_apply_size_delta``
    が run ではなくここへ書くのは，run の無い空段落や採番記号にも効かせるため．
    未指定（テーマ継承のまま）は None．

    ``ph_idx=1`` は本文（タイトルレイアウトでは著者・所属）の枠．既定テンプレートは
    このテストが使うどのレイアウトにも idx 1 を持つ．見つからなければ ``assert`` で
    止まる——**その枠が無ければテストの前提が崩れている**ので，静かに読み飛ばす
    のではなく落ちるのが正しい．
    """
    slide = prs.slides[slide_idx]
    ph = next((s for s in slide.placeholders
               if s.placeholder_format.idx == ph_idx), None)
    # 失敗時はレイアウト名まで出す．このテストの期待値はレイアウトごとの上書きに
    # 依存するので，どのレイアウトを見に行けばよいかが分かるほうが早い．
    assert ph is not None, (
        f"placeholder idx={ph_idx} not found on slide {slide_idx} "
        f"(layout {slide.slide_layout.name!r})")
    out = []
    for para in ph.text_frame._txBody.findall(f"{_A}p"):
        text = "".join(t.text or "" for t in para.iter(f"{_A}t"))
        el = para.find(f"{_A}pPr/{_A}defRPr")
        sz = int(el.get("sz")) / 100.0 if el is not None and el.get("sz") else None
        out.append((text, sz))
    return out


def test_layout_override_is_the_base(tmp_path):
    """レイアウトが上書きしたサイズが基点になる（マスターより優先）．

    レイアウト 2「Section Header」は本文 lvl1 を 20pt に上書きしている．
    修正前はマスターの 32pt を基点にしていたので 28pt になっていた．
    """
    prs = _build(tmp_path, """---
theme: t.pptx
---

### 章扉
<!-- @layout: 2 -->

- {-1} 添え書き
""")
    assert _para_sizes(prs) == [("添え書き", round(20.0 / _STEP))]


def test_master_style_is_the_base_without_override(tmp_path):
    """レイアウトが上書きしない枠ではマスター本文サイズが基点のまま．

    通常のコンテンツスライド（レイアウト 1）はどのテーマでも上書きが無いことが
    多く，**既存原稿の本文サイズが動かない**ことがこの修正の前提だった．
    """
    prs = _build(tmp_path, """---
theme: t.pptx
---

### 見出し

- {-1} 本文
""")
    assert _para_sizes(prs) == [("本文", round(_MASTER_BODY_LVL1 / _STEP))]


def test_override_is_merged_per_level(tmp_path):
    """上書きはレベル単位で重なる（インデントした行はその段の値が基点）．

    レイアウト 3「Two Content」は lvl1=28 / lvl2=24．テーマが一部のレベルだけ
    上書きする場合に備え，欠けたレベルはマスター側で埋める設計になっている．
    """
    prs = _build(tmp_path, """---
theme: t.pptx
---

### 比較

- {+1} 親
  - {+1} 子

<!-- @col -->

- 右
""")
    assert _para_sizes(prs) == [("親", round(28.0 * _STEP)),
                                ("子", round(24.0 * _STEP))]


def test_same_frame_same_size_across_paths(tmp_path):
    """同じ枠に出る同じ記法は，front matter 経由でも本文経由でも同じサイズ．

    これが Issue #83 の症状①そのもの．どちらもタイトルレイアウトの idx 1
    （著者・所属の枠）へ入るのに，基点が違うため 2pt ずれていた．
    """
    front = _build(tmp_path, """---
theme: t.pptx
title: 表題
author: 著者
affiliation:
  - "{-1} 所属"
---
""", name="front.pptx")
    body = _build(tmp_path, """---
theme: t.pptx
---

### 表題
<!-- @layout: 0 -->

- 著者
- {-1} 所属
""", name="body.pptx")
    assert _para_sizes(front) == _para_sizes(body)
    # 値そのものも固定する（両方が同じだけ間違っていても気づけるように）．
    assert _para_sizes(body) == [("著者", None),
                                 ("所属", round(_MASTER_BODY_LVL1 / _STEP))]
