#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""表のセルを個別に塗れることを固定するテスト（Issue #112）．

講義スライドの表14個のうち**7個がセルを個別に着色**していた。NAT 変換表で
「変換前」と「変換後」を色分けするなど、**着色そのものが説明の一部**になっている。
これまではヘッダ行しか塗れなかった。

記法は #105 で入れた行内装飾の色指定と同じ ``{色}`` を使う。セルの中身の末尾に
書くと**そのセルの背景**になる:

    | 方向 | 変換前 | 変換後 |
    |:--|:--|:--|
    | 送信 | 192.168.0.2 {bg2} | 203.0.113.1 {accent2} |

**同じものを2通りで書けるようにしない**——色名の語彙も ``colors.parse_color`` の
ままで、テーマ色名／CSS の色名／16進がそのまま使える。
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR

from md2pptx import render
from md2pptx.parser import parse

_FM = "---\ntheme: t.pptx\n---\n\n"


def _table(src_body):
    slide, = parse(_FM + "### x\n\n" + src_body + "\n").slides
    return [b for b in slide.blocks if hasattr(b, "rows")][0]


def _build(tmp_path, src):
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _grid(slide):
    for sh in slide.shapes:
        if sh.has_table:
            return sh.table
    raise AssertionError("表が無い")


_NAT = """| 方向 | 変換前 | 変換後 |
|:--|:--|:--|
| 送信 | 192.168.0.2 {bg2} | 203.0.113.1 {accent2} |
| 受信 | 203.0.113.1 {accent2} | 192.168.0.2 {bg2} |"""


# ---------------------------------------------------------------- パース

def test_the_colour_token_is_taken_out_of_the_text():
    """``{色}`` はセルの指定であって中身ではない（表示テキストから外す）．"""
    t = _table(_NAT)
    assert t.rows[0] == ["送信", "192.168.0.2", "203.0.113.1"]


def test_the_colour_is_kept_per_cell():
    """セルごとの色を保持する（指定の無いセルは None）．"""
    t = _table(_NAT)
    assert t.fills[0] == [None, "bg2", "accent2"]
    assert t.fills[1] == [None, "accent2", "bg2"]


def test_a_table_without_colours_has_no_fills():
    """色を書かない表は従来どおり（``fills`` は空）．"""
    t = _table("| a | b |\n|---|---|\n| 1 | 2 |")
    assert t.fills == []


def test_a_brace_that_is_not_a_colour_stays_as_text():
    """色名でない ``{…}`` は文字のまま残す（式や記号を壊さない）．"""
    t = _table("| a | b |\n|---|---|\n| {n} 個 | x |")
    assert t.rows[0] == ["{n} 個", "x"] and t.fills == []


def test_an_unknown_colour_name_stops():
    """色として書いたつもりの綴り違いは止める．

    ``{akairo}`` のような**色名らしい語**だけを対象にする——``{n}`` のような
    式まで拾うと、書けるものを勝手に減らすことになる．
    """
    with pytest.raises(ValueError, match="unknown color"):
        parse(_FM + "### x\n\n| a |\n|---|\n| v {redish} |\n")


def test_the_header_can_be_coloured_too():
    """ヘッダのセルにも書ける（既定のアクセント色を上書きする）．"""
    t = _table("| 左 {tx2} | 右 |\n|---|---|\n| 1 | 2 |")
    assert t.header == ["左", "右"]
    assert t.header_fills == ["tx2", None]


# ---------------------------------------------------------------- 描画

def test_the_fill_reaches_the_cell(tmp_path):
    """テーマ色はテーマ色として塗る（RGB へ潰さない）．"""
    slide, = _build(tmp_path, _FM + "### x\n\n" + _NAT + "\n").slides
    grid = _grid(slide)
    assert grid.cell(1, 2).fill.fore_color.theme_color == MSO_THEME_COLOR.ACCENT_2


def test_a_concrete_colour_is_written_as_rgb(tmp_path):
    """具体的な色名は RGB として塗る．"""
    src = _FM + "### x\n\n| a |\n|---|\n| v {red} |\n"
    slide, = _build(tmp_path, src).slides
    assert str(_grid(slide).cell(1, 0).fill.fore_color.rgb) == "FF0000"


def test_an_uncoloured_cell_is_left_alone(tmp_path):
    """指定の無いセルはテーマ任せのまま（塗りを足さない）．"""
    src = _FM + "### x\n\n| a | b |\n|---|---|\n| 1 | 2 {red} |\n"
    slide, = _build(tmp_path, src).slides
    xml = _grid(slide).cell(1, 0)._tc.xml
    assert "srgbClr" not in xml
