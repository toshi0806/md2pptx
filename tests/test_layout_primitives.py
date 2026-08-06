#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""座標計算の共通プリミティブと線ヘルパーを固定するテスト（Issue #108）．

図 DSL を足すたびに ``Rect`` や ``Placed*`` を各モジュールで書き直すのは無駄だし、
**``ObjectBlock`` を増やしたときの取りこぼし**が起きる——``(Table, Flow, Image)``
というタプルが render の 5 か所に複製されていた。

これは**機能を変えないリファクタ**なので、いちばん大事なのは「出力が変わらない」
こと。ここでは移した先が使えることと、``is_object_block`` が
``ObjectBlock`` の定義に自動で追従することを見る。

``line`` は ``block_arrow`` と役割が違う。``block_arrow`` はノード間の**すき間に
収まる塗り矢印**で box に食い込ませないための道具、``line`` は任意の 2 点を結ぶ
細い線で、シーケンス図のライフラインのように**図の骨格**を描くためのもの。
"""
from __future__ import annotations

from typing import get_args

from pptx import Presentation

from md2pptx import ir, render
from md2pptx.flow import parse_flow, plan_flow
from md2pptx.layout import EMU, PlacedLine, Rect, emu


def test_rect_reads_its_edges_by_name():
    """端と中心は名前で読む（足し算で出すと取り違えても気づけない）．"""
    r = Rect(left=10, top=20, width=100, height=50)
    assert (r.right, r.bottom) == (110, 70)
    assert (r.center_x, r.center_y) == (60, 45)


def test_emu_converts_inches():
    assert emu(1) == EMU == 914400
    assert emu(0.5) == 457200


def test_flow_still_plans_through_the_shared_primitives():
    """flow は移した先の ``Rect`` を使って従来どおり配置する．

    移設で座標が変わっていないことが、このリファクタの受け入れ条件。
    """
    plan = plan_flow(parse_flow("[a] -> [b]"), 0, 0, emu(10), emu(5))
    assert len(plan.boxes) == 2 and len(plan.arrows) == 1
    first, second = (b.rect for b in plan.boxes)
    assert isinstance(first, Rect)
    assert first.left < second.left           # lr は左から右へ
    assert first.center_y == second.center_y  # 高さは揃う


def test_is_object_block_follows_the_type_alias():
    """``is_object_block`` は ``ObjectBlock`` の定義に自動で追従する．

    型注釈と別に ``(Table, Flow, Image)`` を書き並べると、種類を増やしたときに
    必ずどこかが漏れる。**ここが漏れなければ render の 5 か所も漏れない**。
    """
    samples = {ir.Table: lambda: ir.Table(),
               ir.Flow: lambda: ir.Flow(),
               ir.Image: lambda: ir.Image(src="fig.png")}
    kinds = get_args(ir.ObjectBlock)
    # 種類が増えたらこのテストが先に落ちる（作り方を書き足せ、という合図）．
    assert set(kinds) == set(samples), "ObjectBlock が増えた——見本を足すこと"
    for cls in kinds:
        assert ir.is_object_block(samples[cls]()), cls
    assert not ir.is_object_block(ir.Line(text="ふつうの行"))


def test_placed_line_carries_the_dashed_flag():
    """``PlacedLine`` は破線かどうかを持つ（矢尻の有無は ``PlacedArrow`` と分ける）．"""
    assert PlacedLine(0, 0, 10, 10).dashed is False
    assert PlacedLine(0, 0, 10, 10, dashed=True).dashed is True


def _slide(tmp_path):
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    r = render.Renderer(str(theme))
    return r, r.prs.slides.add_slide(r.L1)


def test_line_draws_a_connector(tmp_path):
    """``line`` は 2 点を結ぶ線を置く（塗り矢印ではない）．"""
    r, slide = _slide(tmp_path)
    before = len(slide.shapes)
    r.line(slide, emu(1), emu(1), emu(3), emu(2))
    assert len(slide.shapes) == before + 1


def test_line_can_be_dashed_and_arrowed(tmp_path):
    """破線と矢じりを指定できる（ライフラインとメッセージの描き分けに要る）．"""
    r, slide = _slide(tmp_path)
    conn = r.line(slide, 0, 0, emu(2), 0, dashed=True, arrow=True)
    xml = conn._element.xml
    assert "dash" in xml
    assert "tailEnd" in xml


def test_a_plain_line_has_neither(tmp_path):
    """既定は実線・矢尻なし（ライフラインはこちら）．"""
    r, slide = _slide(tmp_path)
    xml = r.line(slide, 0, 0, emu(2), 0)._element.xml
    assert "tailEnd" not in xml
