#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""行内装飾を固定するテスト（Issue #105）．

md2pptx には行内装飾が**一つも無かった**（``**bold**`` すら文字として出た）。
講義スライドの実測では色付きの run が 1,122・太字 910・上付下付 163 あり、
「強調したいアドレス部分だけ色を変える」「変換前と変換後を色分けする」が
説明の要になっている。

**色は3系統を受ける**——テーマ色名／CSS の色名／16進。テーマ色名を既定に
薦めるのはテーマ差し替えに追従するからで、具体的な色名を入れるのは
「危険＝赤」のように**意味が色そのものに宿る**場面があるから。

装飾の無い行は ``Line.spans`` が空のまま返り、描画も従来の1 run 経路を通る。
装飾を使わない原稿の出力が変わらないことを、ここで固定しておく。
"""
from __future__ import annotations

import pytest
from pptx import Presentation

from md2pptx import render
from md2pptx.parser import parse

_FM = "---\ntheme: t.pptx\n---\n\n"


def _line(src_body):
    slide, = parse(_FM + "### x\n\n" + src_body + "\n").slides
    return slide.blocks[0]


def _spans(src_body):
    return [(s.text, s.bold, s.mono, s.color, s.link, s.script)
            for s in _line(src_body).spans]


def _build(tmp_path, src):
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _runs(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh != slide.shapes.title:
            return sh.text_frame.paragraphs[0].runs
    return []


# ---------------------------------------------------------------- 記法

def test_bold():
    assert _spans("- 前 **強調** 後") == [
        ("前 ", False, False, None, None, None),
        ("強調", True, False, None, None, None),
        (" 後", False, False, None, None, None),
    ]


def test_inline_code():
    assert _spans("- `dig` を使う")[0] == ("dig", False, True, None, None, None)


def test_superscript_and_subscript():
    assert _spans("- 2^32^ 個")[1] == ("32", False, False, None, None, "sup")
    assert _spans("- H~2~O")[1] == ("2", False, False, None, None, "sub")


def test_link():
    assert _spans("- [JPNIC](https://www.nic.ad.jp/)") == [
        ("JPNIC", False, False, None, "https://www.nic.ad.jp/", None)]


@pytest.mark.parametrize("name", ["accent2", "red", "#ff0000", "#f00", "RED"])
def test_color_accepts_three_families(name):
    """テーマ色名・CSS の色名・16進のいずれも受ける（大小文字は問わない）．"""
    assert _spans(f"- [色]{{{name}}} の行")[0][3] == name


def test_an_unknown_color_stops():
    """綴りを間違えた色は黙って既定色にせず止める．

    「効かない」で済ませると、色が付いていないことに気づかないまま配る。
    """
    with pytest.raises(ValueError, match="unknown color"):
        parse(_FM + "### x\n\n- [これ]{akairo} は色か\n")


def test_decorations_nest():
    """``[**赤い強調**]{red}`` のように重ねられる．

    重ねられないと「色を付けたら太字にできない」という例外を説明することになる。
    """
    assert _spans("- [**赤い強調**]{red}") == [
        ("赤い強調", True, False, "red", None, None)]


def test_a_plain_line_has_no_spans():
    """装飾の無い行は spans が空——従来の1 run 経路をそのまま通る．"""
    line = _line("- ふつうの行")
    assert line.spans == [] and line.text == "ふつうの行"


def test_the_text_keeps_the_plain_string():
    """``Line.text`` は装飾記号を除いた素のテキストのまま．

    段落の折り返しや帯の高さの見積もりが text の長さを見るので、
    記号が残っていると勘定が狂う。
    """
    assert _line("- **強調**と`等幅`").text == "強調と等幅"


def test_markup_works_on_every_line_kind():
    """箇条書き・採番・結論行のどれでも効く．"""
    assert _spans("1. **採番**")[0][1] is True
    assert _spans("→ **結論**")[1][1] is True


# ---------------------------------------------------------------- 描画

def test_bold_and_color_reach_the_runs(tmp_path):
    """太字とテーマ色が run に届く（テーマ色は RGB へ潰さない）．

    潰すとテーマを差し替えたときに色が取り残される。
    """
    from pptx.enum.dml import MSO_THEME_COLOR
    slide, = _build(tmp_path, _FM + "### x\n\n- **太**と[色]{accent2}\n").slides
    runs = _runs(slide)
    assert [r.text for r in runs] == ["太", "と", "色"]
    assert runs[0].font.bold is True
    assert runs[2].font.color.theme_color == MSO_THEME_COLOR.ACCENT_2


def test_a_concrete_colour_reaches_the_runs_as_rgb(tmp_path):
    """具体的な色名は RGB として書かれる（テーマに無い色を指定できる）．"""
    slide, = _build(tmp_path, _FM + "### x\n\n- [危険]{red}\n").slides
    assert str(_runs(slide)[0].font.color.rgb) == "FF0000"


def test_superscript_reaches_the_runs(tmp_path):
    """上付きは ``a:rPr/@baseline`` として書かれる（python-pptx に API が無い）．"""
    slide, = _build(tmp_path, _FM + "### x\n\n- 2^32^\n").slides
    runs = _runs(slide)
    assert runs[1].text == "32"
    assert runs[1].font._rPr.get("baseline") == "30000"


def test_a_link_reaches_the_runs(tmp_path):
    """ハイパーリンクが run に付く．"""
    slide, = _build(tmp_path, _FM + "### x\n\n- [JPNIC](https://www.nic.ad.jp/)\n").slides
    assert _runs(slide)[0].hyperlink.address == "https://www.nic.ad.jp/"


def test_segment_sizes_still_land_on_the_right_segment(tmp_path):
    """``<br>`` のセグメントごとの相対サイズが、装飾で run が割れてもずれない．

    run とセグメントを**位置で**対応させていると、1 セグメントが複数 run に
    割れた時点で別のセグメントへサイズが付く。1 つずれたサイズは正しく見えて
    しまうので、崩れより気づきにくい。
    """
    slide, = _build(
        tmp_path, _FM + "### x\n\n- **太**い前半<br>{-2} 小さい後半\n").slides
    runs = _runs(slide)
    assert [r.text for r in runs] == ["太", "い前半", "小さい後半"]
    assert runs[0].font.size is None and runs[1].font.size is None
    assert runs[2].font.size is not None      # 縮んでいるのは後半だけ
