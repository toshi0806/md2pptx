#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""表紙を本文記法で書けることを固定するテスト（Issue #82）．

フロントマターの ``title`` / ``subtitle`` / ``author`` / ``affiliation`` は
「文書のメタデータ」の名前を持ちながら，実体は表紙スライドの描画記述だった
（pptx のコアプロパティへは一度も書いていない）．同じ絵を出す記法が 2 つある状態が
続き，``<br>`` を片方だけ通し忘れる取りこぼし（Issue #79）まで起きた．

表紙は ``# 見出し`` に ``<!-- @layout: 0 -->`` を添えて書く．扉（``#`` ＝ レイアウト 2）
との違いはレイアウト番号だけで，記法も描画経路も共通になる．**「先頭スライドなら
自動的に表紙」という位置依存の判定は採らない**——front matter を排した理由が
「書いたものがそのまま出る」ことなので，暗黙の位置ルールを持ち込むと同じ問題を
別の形で作る．ここでは 2 枚目でも表紙になれることを固定して，その設計を守る．

番号は**レイアウト 0 のスライドに付けない**．そのレイアウトを選ぶこと自体が
「これは表紙」の宣言なので，番号の有無をそこに紐づける（フロントマター由来の
表紙を描く ``render_title_slide`` が番号を付けないのと揃う）．

非推奨のフロントマターは**受理をやめない**．警告だけ出して従来どおり描く——
非推奨の間は既存原稿が同じ見た目で動き続けるほうが価値がある（Issue #83 で
副題の基点を凍結したのと同じ理由）．
"""
from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from md2pptx import render
from md2pptx.parser import parse


def _theme(tmp_path):
    path = tmp_path / "theme.pptx"
    Presentation().save(str(path))
    return str(path)


def _build(tmp_path, src, name="out.pptx"):
    out = tmp_path / name
    r = render.Renderer(_theme(tmp_path))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _numbered(slide):
    """スライドに番号プレースホルダが入っているか．

    **種別で判定する**（``add_slide_number`` は ``idx == 12`` で探すが，それは
    実装側の前提）．テストまで同じ添字を見ると，添字の取り違えを検出できない．
    """
    return any(ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER
               for ph in slide.placeholders)


def test_layout_zero_slides_get_no_number(tmp_path):
    """レイアウト 0 だけ番号が付かない（他のレイアウトは従来どおり）．"""
    prs = _build(tmp_path, """---
theme: t.pptx
slide_number: true
---

# 表紙
<!-- @layout: 0 -->

- 著者

## 本文スライド

- A

# 章扉
""")
    got = [(s.slide_layout.name, _numbered(s)) for s in prs.slides]
    assert got == [("Title Slide", False),
                   ("Title and Content", True),
                   ("Section Header", True)]


def test_a_title_slide_can_sit_anywhere(tmp_path):
    """表紙は位置ではなくレイアウト指定で決まる（暗黙の位置ルールを持たない）．

    2 枚目に置いても表紙になり，1 枚目の扉は番号が付いたままになる．
    """
    prs = _build(tmp_path, """---
theme: t.pptx
slide_number: true
---

# 章扉

# 表紙
<!-- @layout: 0 -->
""")
    assert [_numbered(s) for s in prs.slides] == [True, False]


def test_slide_number_false_still_wins(tmp_path):
    """``slide_number: false`` は従来どおり全体に効く．"""
    prs = _build(tmp_path, """---
theme: t.pptx
slide_number: false
---

## 本文スライド
""")
    assert _numbered(prs.slides[0]) is False


def test_the_title_frame_and_the_body_frame_are_filled(tmp_path):
    """主題・副題はタイトル枠へ，著者・所属は本文枠へ入る．

    副題をタイトル枠に置くのはテーマがそう作られているため——多くのテーマは
    タイトル枠と著者欄の間に罫線を引いており，副題を下の枠へ移すと罫線の下に
    落ちて著者情報の塊に混ざる（実 PowerPoint で確認．Issue #82）．
    """
    prs = _build(tmp_path, """---
theme: t.pptx
---

# 主題<br>{-5} ― 副題 ―
<!-- @layout: 0 -->

- 著者
- 所属
""")
    slide = prs.slides[0]
    frames = {ph.placeholder_format.idx: ph.text_frame.text
              for ph in slide.placeholders
              if ph.placeholder_format.idx in (0, 1)}
    # 枠が無ければ「テストの前提が崩れた」ので落ちるのが正しいが，KeyError では
    # 何が起きたか読み取れない．どちらの枠が欠けたかを言う．
    assert set(frames) == {0, 1}, (
        f"title/body frames not found on {slide.slide_layout.name!r}: "
        f"got idx {sorted(frames)}")
    # 区切り文字の違いは枠ではなく**段落構造**の違い．python-pptx の
    # text_frame.text は段落の境目を "\n"，段落内の改行（a:br）を "\v" で返す．
    # 主題と副題は <br> でつないだ 1 段落なので "\v"，著者と所属は別々の行として
    # 書いたので別段落＝"\n" になる．
    assert frames[0] == "主題\v― 副題 ―"
    assert frames[1] == "著者\n所属"


def test_deprecated_front_matter_still_renders(tmp_path, capsys):
    """非推奨のフロントマターは警告を出しつつ従来どおり表紙を描く．

    受理をやめると動く原稿が黙って壊れる．移行期間の価値は「同じ見た目で
    動き続けること」なので，出力は変えない．
    """
    prs = _build(tmp_path, """---
theme: t.pptx
title: 主題
author: 著者
---

## 本文スライド
""")
    assert len(prs.slides) == 2                     # 表紙＋本文
    assert _numbered(prs.slides[0]) is False        # 表紙に番号は付かない
    err = capsys.readouterr().err
    assert "deprecated" in err
    # 移行先を示すこと．警告だけ出して書き換え方が分からないと動きようがない．
    assert "@layout: 0" in err
    assert "title" in err and "author" in err


def test_no_warning_without_the_deprecated_keys(tmp_path, capsys):
    """表紙を本文記法で書いた原稿には警告を出さない．"""
    _build(tmp_path, """---
theme: t.pptx
---

# 主題
<!-- @layout: 0 -->
""")
    assert "deprecated" not in capsys.readouterr().err
