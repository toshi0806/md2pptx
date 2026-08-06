#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""採番の開始番号を固定するテスト（Issue #107）．

PowerPoint の自動採番は**プレースホルダごとに 1 から数え直す**。そのため
``<!-- @col -->`` で 2 カラムにした採番リストは、右カラムが ``1.`` に戻っていた。
シラバスや講義日程のように「番号付きリストを 2 カラムに割る」形は毎年書くので、
これは実害が大きい。

原稿に書いた番号を捨てず、**リストの先頭の番号**を種にして md2pptx 側で数える。
先頭だけを種にするのは CommonMark の規則（「順序付きリストの開始番号は最初の項目の
番号で、以降の番号は無視される」）と同じで、``1. 1. 1.`` と書けば 1・2・3 になる
従来の書き方もそのまま動く。

**番号は全ての採番段落へ明示的に書く。** 先頭にだけ ``startAt`` を付ける実装を
最初に試したが、PowerPoint は ``startAt`` の付いた段落の**次から数え直す**ため、
右カラムが「8. 1. 2. 3. …」になった（実 PowerPoint で確認）。数えるのは
md2pptx 側の仕事にする。

観測するのは pptx に出た ``a:buAutoNum`` の ``startAt``——これは出力そのもので、
実装の内部ではない。python-pptx に読み取り API が無いので XML を直接見る。
"""
from __future__ import annotations

from pptx import Presentation

from md2pptx import render
from md2pptx.parser import parse

_FM = "---\ntheme: t.pptx\n---\n\n"


def _build(tmp_path, src):
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _numbers(slide):
    """スライド内の採番段落を (テキスト, 番号) で返す（出現順）．

    番号は ``a:buAutoNum`` の ``startAt``。全ての採番段落に付いているはずで、
    付いていない段落は None として現れる（＝数え直しが起きる状態）。
    """
    import re
    found = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            text = "".join(r.text for r in p.runs)
            m = re.search(r"<a:buAutoNum[^>]*>", p._p.xml)
            if not m:
                continue
            at = re.search(r'startAt="(\d+)"', m.group(0))
            found.append((text, int(at.group(1)) if at else None))
    return found


def test_the_second_column_continues_the_numbering(tmp_path):
    """右カラムが ``8.`` から続く（これが直したかったもの）．"""
    src = _FM + ("### シラバス\n\n"
                 "1. 第1回\n"
                 "2. 第2回\n"
                 "<!-- @col -->\n"
                 "8. 第8回\n"
                 "9. 第9回\n")
    slide, = _build(tmp_path, src).slides
    assert _numbers(slide) == [
        ("第1回", 1), ("第2回", 2), ("第8回", 8), ("第9回", 9)]


def test_only_the_first_number_counts(tmp_path):
    """先頭の番号だけが種になる——``1. 1. 1.`` は従来どおり 1・2・3．

    CommonMark と同じ規則。原稿の番号をそのまま全行へ書き込むと、この慣用的な
    書き方が 1・1・1 になって壊れる（``example.md`` もこの書き方をしている）。
    """
    src = _FM + "### x\n\n1. a\n1. b\n1. c\n"
    slide, = _build(tmp_path, src).slides
    assert _numbers(slide) == [("a", 1), ("b", 2), ("c", 3)]


def test_a_list_that_does_not_start_at_one(tmp_path):
    """先頭が ``3.`` なら 3 から始まる（続きを別スライドに分けたとき）．"""
    src = _FM + "### 続き\n\n3. 三番目\n4. 四番目\n"
    slide, = _build(tmp_path, src).slides
    assert _numbers(slide) == [("三番目", 3), ("四番目", 4)]


def test_parenthesised_numbers_carry_the_start_too(tmp_path):
    """``(1)`` 形式でも同じ．"""
    src = _FM + "### x\n\n- 左\n<!-- @col -->\n(5) 五\n(6) 六\n"
    slide, = _build(tmp_path, src).slides
    assert _numbers(slide) == [("五", 5), ("六", 6)]


def test_circled_numbers_carry_the_start_too(tmp_path):
    """丸数字でも同じ（``⑤`` から始めれば 5）．"""
    src = _FM + "### x\n\n- 左\n<!-- @col -->\n⑤ 五\n⑥ 六\n"
    slide, = _build(tmp_path, src).slides
    assert _numbers(slide) == [("五", 5), ("六", 6)]


def test_each_level_counts_on_its_own(tmp_path):
    """階層ごとに数える（入れ子の内側が外側の番号を食わない）．

    内側のリストは**外側の項目が変わっても数え直さない**。1 つの枠の中では
    (深さ, 形式) ごとに 1 本の連番、という規則にしている——数え直す条件を
    足すと「どこで区切れるか」を別に決めることになるので、規則は 1 つで足りる。
    """
    src = _FM + ("### x\n\n"
                 "1. 外1\n"
                 "  1. 内1\n"
                 "  2. 内2\n"
                 "2. 外2\n"
                 "  1. 内3\n")
    slide, = _build(tmp_path, src).slides
    assert _numbers(slide) == [
        ("外1", 1), ("内1", 1), ("内2", 2), ("外2", 2), ("内3", 3)]


def test_the_numbering_survives_a_figure_in_the_middle(tmp_path):
    """導入文と結論文に分かれても振り直されない．

    表・図のあるスライドでは地の文が帯の上下に分かれ、``_append_lines`` が
    同じ枠に対して 2 回呼ばれる。呼び出しごとに数え直すと、結論文の採番が
    1 に戻ってしまう。
    """
    src = _FM + ("### x\n\n"
                 "1. 導入1\n"
                 "2. 導入2\n"
                 "\n| a | b |\n|---|---|\n| 1 | 2 |\n\n"
                 "3. 結論1\n")
    slide, = _build(tmp_path, src).slides
    assert _numbers(slide) == [("導入1", 1), ("導入2", 2), ("結論1", 3)]


def test_every_numbered_paragraph_carries_its_number(tmp_path):
    """番号は全段落に書く——先頭だけでは PowerPoint が次から数え直す．

    最初はリストの先頭にだけ ``startAt`` を付けたが、実 PowerPoint では
    「8. 1. 2. 3. …」になった。これはその実装へ戻らないための固定。
    """
    src = _FM + "### x\n\n5. a\n6. b\n7. c\n"
    slide, = _build(tmp_path, src).slides
    assert all(n is not None for _, n in _numbers(slide))
