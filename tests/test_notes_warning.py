#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""発表者ノートを置けないときに黙らないことを固定するテスト（Issue #35 の途中で追加）．

``` ```note ``` で書いたノートは、テーマのノートマスターにノート本文用の
プレースホルダがあって初めて置ける．無いテーマでは置き場所が無い．

**止めはしない**——ノートは補助情報で、ここで落とすとスライド自体が出せなくなる
（「PDF が作れなくても pptx は成功」と同じ考え方）．**ただし黙らない**．
書いたノートが出ないことは pptx を開くまで分からず、開いても原因がテーマ側だとは
思い当たらない．本文プレースホルダについて Issue #67 で決めたのと同じ扱い．

テーマは python-pptx 同梱の既定テンプレートのノートマスターから BODY を外して
作るので、リポジトリの ``OfficeTheme.pptx`` にも実 PowerPoint にも依存しない．
"""
from __future__ import annotations

import pytest
from pptx import Presentation

from md2pptx import render
from md2pptx.ir import Deck, Line, Slide


def _theme(tmp_path, with_notes_placeholder=True):
    """ノート本文の置き場所の有無を選べるテーマ pptx を作って返す．"""
    prs = Presentation()
    if not with_notes_placeholder:
        for ph in list(prs.notes_master.placeholders):
            if str(ph.placeholder_format.type).startswith("BODY"):
                ph._element.getparent().remove(ph._element)
    path = tmp_path / f"theme-{with_notes_placeholder}.pptx"
    prs.save(str(path))
    return str(path)


def _build(theme, tmp_path, notes, name="out.pptx"):
    out = tmp_path / name
    r = render.Renderer(theme)
    r.render(Deck(slides=[
        Slide(title="見出し", blocks=[Line(text="本文")], notes=notes)]))
    r.save(str(out))
    return out


NOTES = "ここで一呼吸おいて、聴衆の反応を見る。"


def test_a_normal_theme_keeps_the_notes_and_says_nothing(tmp_path, capsys):
    """置き場所があれば普通に入り、警告は出ない．"""
    out = _build(_theme(tmp_path), tmp_path, NOTES)

    slide, = Presentation(str(out)).slides
    assert slide.notes_slide.notes_text_frame.text == NOTES
    assert capsys.readouterr().err == ""


def test_it_reports_notes_it_could_not_place(tmp_path, capsys):
    """置き場所が無いとノートは消える——**そのことを言う**．"""
    _build(_theme(tmp_path, with_notes_placeholder=False), tmp_path, NOTES)

    err = capsys.readouterr().err
    assert "no placeholder for notes text" in err
    assert NOTES[:20] in err, "どのノートが消えたか分からないと，直しようがない"


def test_a_slide_without_notes_stays_quiet(tmp_path, capsys):
    """ノートを書いていなければ言わない（置けなくても困らない）．

    ここで鳴らすと、ノートを使わない人がテーマを変えるたびに警告を浴びる．
    """
    _build(_theme(tmp_path, with_notes_placeholder=False), tmp_path, None)

    assert capsys.readouterr().err == ""


def test_it_does_not_turn_a_build_into_a_failure(tmp_path, capsys):
    """警告であって失敗ではない——スライド本体は出す．"""
    out = _build(_theme(tmp_path, with_notes_placeholder=False), tmp_path, NOTES)

    assert out.exists() and out.stat().st_size > 0
    slide, = Presentation(str(out)).slides
    assert slide.shapes.title.text == "見出し"
    assert "no placeholder for notes text" in capsys.readouterr().err
