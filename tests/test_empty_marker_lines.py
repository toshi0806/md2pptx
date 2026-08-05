#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""箇条書きマーカーだけの行が空段落になることを固定するテスト（Issue #82）．

表紙の著者欄やセクション扉のように**記号の出ない枠**では，行の塊を分けるのに
1 行空けたい．front matter の ``affiliation`` では ``- "{-2} "`` と書けば空段落が
残っていたのに，本文では捨てられていた——同じことを二重に実装した結果のずれで，
本文側だけが取りこぼしていた．

壊れ方は 2 通りあった．

- ``- {-2} `` … マーカーとトークンを除くと空になり，**行ごと落ちていた**
- ``- `` … ``strip()`` 後に ``"-"`` だけになり ``"- "`` 判定を外れて既定の箇条書きへ落ち、
  **文字の "-" が出ていた**

代用の ``- <br>`` は残るが**1 行多く空く**（段落 1 行＋``a:br`` の 2 行目）ので，
空の段落そのものを作れる必要がある．

**末尾に空白の無い形（``-``）を必ず受けること**が要点のひとつ．多くのエディタは
保存時に行末空白を除去するので，``- `` だけを受けると空行スペーサが保存した瞬間に
壊れる．テストで両方の形を押さえる．

空を段落として残すのは**箇条書きマーカーだけ**にしてある．採番行（``①`` / ``(1)``）を
空で残すと番号を 1 つ消費するだけで，空けたい 1 行は得られない．
"""
from __future__ import annotations

from pptx import Presentation

from md2pptx import render
from md2pptx.parser import parse

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _theme(tmp_path):
    path = tmp_path / "theme.pptx"
    Presentation().save(str(path))
    return str(path)


def _lines(body):
    """本文だけの Markdown を parse して Line の (本文, 段数, 種別) を返す．"""
    deck = parse("---\ntheme: t.pptx\n---\n\n## 見出し\n" + body)
    return [(b.text, b.size_delta, b.kind) for b in deck.slides[0].blocks]


def _paragraphs(tmp_path, src, name="out.pptx", ph_idx=1):
    """描画して保存し，本文枠の段落を (本文, サイズ pt, run 数, br 数) で返す．

    保存した pptx を読み直すのは，利用者が開くのがファイルのほうだから
    （test_text_language.py と同じ理由）．空段落は run を持たないので，
    ``run 数 == 0`` が「本当に空の段落」であることの確認になる．

    ``ph_idx=1`` は本文（タイトルレイアウトでは著者・所属）の枠．python-pptx
    同梱の既定テンプレートはレイアウト 0「Title Slide」もレイアウト 1
    「Title and Content」も idx 0（表題）と idx 1 を持つので，このテストが使う
    どちらのレイアウトでも解決できる．テンプレートが変わって見つからなくなったら
    ``assert`` で止まる——**その枠が無ければテストの前提が崩れている**ので，
    静かに読み飛ばすのではなく落ちるのが正しい．
    """
    out = tmp_path / name
    r = render.Renderer(_theme(tmp_path))
    r.render(parse(src))
    r.save(str(out))
    slide = Presentation(str(out)).slides[0]
    ph = next((s for s in slide.placeholders
               if s.placeholder_format.idx == ph_idx), None)
    assert ph is not None, f"placeholder idx={ph_idx} not found"
    got = []
    for para in ph.text_frame._txBody.findall(f"{_A}p"):
        text = "".join(t.text or "" for t in para.iter(f"{_A}t"))
        el = para.find(f"{_A}pPr/{_A}defRPr")
        sz = int(el.get("sz")) / 100.0 if el is not None and el.get("sz") else None
        got.append((text, sz,
                    len(para.findall(f"{_A}r")), len(para.findall(f"{_A}br"))))
    return got


def test_bullet_marker_alone_makes_an_empty_line():
    """``- `` も ``-`` も空の箇条書き行になる（行末空白の有無で変わらない）．

    エディタが行末空白を除去しても壊れないことがこの記法の前提．
    """
    for spacer in ("- ", "-", "* ", "*"):
        assert _lines(f"- A\n{spacer}\n- B\n") == [
            ("A", None, "bullet"),
            ("", None, "bullet"),
            ("B", None, "bullet"),
        ], f"spacer={spacer!r}"


def test_size_token_survives_on_an_empty_line():
    """空行にも相対サイズが効く（行高を周りに合わせるために要る）．"""
    assert _lines("- A\n- {-2}\n- B\n")[1] == ("", -2, "bullet")
    assert _lines("- A\n- {-2} \n- B\n")[1] == ("", -2, "bullet")


def test_numbered_markers_do_not_make_empty_lines():
    """採番マーカーだけの行は段落を作らない（番号を無駄に消費しないため）．"""
    assert _lines("① A\n① \n① B\n") == [("A", None, "autonum"),
                                        ("B", None, "autonum")]


def test_blank_source_line_is_still_a_separator():
    """素の空行は従来どおり段落を作らない（段落区切りのまま）．"""
    assert _lines("- A\n\n- B\n") == [("A", None, "bullet"),
                                      ("B", None, "bullet")]


def test_empty_line_is_rendered_as_a_runless_paragraph(tmp_path):
    """空行は run を持たない段落として出る（``- <br>`` の 2 行分にならない）．

    サイズは実装の式を写さず**数値で書く**．既定テンプレートの本文 lvl1 は 32pt で，
    1 段 ÷1.125 の 2 段下は 25pt（32 ÷ 1.125² = 25.28）．比率は DESIGN.md §5.8 で
    決めた仕様なので，変わればこのテストが落ちるのが正しい——実装から計算すると
    式を写すだけになり，比率が変わってもテストは通ってしまう．
    """
    got = _paragraphs(tmp_path, """---
theme: t.pptx
---

## 見出し

- A
- {-2}
- B
""")
    assert got == [("A", None, 1, 0),
                   ("", 25.0, 0, 0),
                   ("B", None, 1, 0)]


def test_body_spacer_matches_the_front_matter_one(tmp_path):
    """本文の空行スペーサが front matter の空 ``affiliation`` と同じ出力になる．

    これがこの修正の目的．同じ枠（タイトルレイアウトの idx 1）に同じ絵を出す 2 つの
    記法が食い違っていた．Issue #82 で front matter 側を非推奨にするとき，
    **見た目を変えずに移行できる**ことがここで担保される．
    """
    front = _paragraphs(tmp_path, """---
theme: t.pptx
title: 表題
author: 著者
affiliation:
  - "{-2} "
  - "{-2} 所属"
---
""", name="front.pptx")
    body = _paragraphs(tmp_path, """---
theme: t.pptx
---

## 表題
<!-- @layout: 0 -->

- 著者
- {-2}
- {-2} 所属
""", name="body.pptx")
    assert front == body
