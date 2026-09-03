#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""``` ```arrow ``` の横位置を帯の左端からの距離で書ける（Issue #180）．

`align` は left / center / right の3択で、元スライドの矢印は**そのどれでもない**
ことが多い。cn2025-01 s27 の下向き矢印は本文枠の左から 3.33in（`center` なら
5.17in、`left` なら 0）で、3択のどれとも 1in 以上ずれる。

大きさは既に「元スライドから実測して書く」運用（DESIGN.md §5.16）なので、
位置も同じ `Length` の語彙で書けるようにする。`left:` は `align` に優先する。
"""
from __future__ import annotations

import pytest

from md2pptx.parser import parse
from md2pptx.ir import Arrow

_FM = "---\ntheme: t.pptx\n---\n\n"


def _arrow(fence: str) -> Arrow:
    slide, = parse(_FM + "### x\n\n```arrow\n" + fence.strip() + "\n```\n").slides
    return [b for b in slide.blocks if isinstance(b, Arrow)][0]


def test_left_is_parsed_as_a_length():
    """``in`` は parser が EMU へ換算する——``width`` / ``height`` と同じ扱い．"""
    a = _arrow("direction: down\nleft: 3.33in")
    assert a.left is not None
    assert a.left.unit == "emu"
    assert a.left.value == pytest.approx(3.33 * 914400, rel=1e-6)


def test_left_defaults_to_none_so_align_still_decides():
    assert _arrow("direction: down").left is None
    assert _arrow("direction: down\nalign: right").left is None


def test_a_percentage_is_accepted_like_width():
    a = _arrow("direction: down\nleft: 25%")
    assert a.left is not None
    assert a.left.unit == "percent" and a.left.value == pytest.approx(25.0)


def test_a_bad_key_still_names_left_in_the_message():
    with pytest.raises(ValueError, match="left"):
        _arrow("direction: down\nsideways: 1in")


def _build(tmp_path, body):
    """既存の arrow テストと同じ組み立て（本文と結論文で帯を作る）．"""
    from pptx import Presentation
    from md2pptx import render
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    r = render.Renderer(str(theme))
    prs = r.render(parse(_FM + "### x\n\n- 上\n\n" + body + "\n\n→ 下\n"))
    return prs


def _down_arrow(prs):
    from pptx.enum.shapes import MSO_SHAPE
    got = [sh for sh in prs.slides[0].shapes
           if not sh.is_placeholder
           and getattr(sh, "auto_shape_type", None) == MSO_SHAPE.DOWN_ARROW]
    assert len(got) == 1, f"下向き矢印が {len(got)} 個"
    return got[0]


def _band_left(prs):
    body = next(sh for sh in prs.slides[0].shapes
                if sh.is_placeholder and sh.placeholder_format.idx == 1)
    return body.left


def test_left_places_the_arrow_there(tmp_path):
    """描かれた矢印の左端が、帯の左端から指定ぶんの位置に来る．"""
    from pptx.util import Emu
    prs = _build(tmp_path, "```arrow\ndirection: down\nwidth: 1in\nleft: 2in\n```")
    got = Emu(_down_arrow(prs).left - _band_left(prs)).inches
    assert got == pytest.approx(2.0, abs=0.02)


def test_left_beats_align(tmp_path):
    """``align`` と両方書いたら ``left`` が勝つ．"""
    only_align = _down_arrow(_build(
        tmp_path / "a", "```arrow\ndirection: down\nwidth: 1in\nalign: right\n```")).left
    both = _down_arrow(_build(
        tmp_path / "b",
        "```arrow\ndirection: down\nwidth: 1in\nalign: right\nleft: 0in\n```")).left
    assert both < only_align
