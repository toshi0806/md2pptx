#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""段落の高さの見積もりを固定するテスト（Issue #145）．

帯（表・図の置き場）の上端は「地の文が何 cm を占めるか」で決まる。そこを
小さく見積もると**図が地の文に重なる**。cn2026-01 p.17 で 0.27cm 足りず、
矢印が4行目に食い込んでいた。

取りこぼしは2つ。

1. テーマの**段落前アキ**（`spcBef`）を数えていない
2. どの行も lvl1 のサイズで数えている（実際はレベルごとに違う）
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from md2pptx import render
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"


def _theme(tmp_path, spc_pct=None, spc_pts=None, sizes=(3000, 2600, 2200)):
    """マスター本文スタイルに spcBef とレベル別サイズを持つテーマを作る．"""
    tmp_path.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    body = prs.slide_masters[0].element.find(
        qn("p:txStyles") + "/" + qn("p:bodyStyle"))
    for lvl, sz in enumerate(sizes, start=1):
        el = body.find(qn(f"a:lvl{lvl}pPr"))
        assert el is not None
        for tag in ("a:spcBef", "a:defRPr"):
            old = el.find(qn(tag))
            if old is not None:
                el.remove(old)
        if spc_pct is not None or spc_pts is not None:
            spc = el.makeelement(qn("a:spcBef"), {})
            inner = spc.makeelement(
                qn("a:spcPct") if spc_pct is not None else qn("a:spcPts"),
                {"val": str(spc_pct if spc_pct is not None else spc_pts)})
            spc.append(inner)
            el.insert(0, spc)
        rpr = el.makeelement(qn("a:defRPr"), {"sz": str(sz)})
        el.append(rpr)
    path = tmp_path / "theme.pptx"
    prs.save(str(path))
    return path


def _build(tmp_path, src, theme):
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


_SRC = (_FM + "### x\n\n- A\n  - a1\n  - a2\n  - a3\n\n"
        "```arrow\ndirection: down\n```\n\n- B\n  - b1\n")


def _body(slide):
    """本文プレースホルダ（idx==1）．スライド番号の枠を拾わない．"""
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == 1:
            return sh
    raise AssertionError("本文プレースホルダが無い")


def _arrow_top(prs):
    slide = prs.slides[-1]
    arrow, = [sh for sh in slide.shapes
              if not sh.is_placeholder
              and getattr(sh, "auto_shape_type", None) is not None]
    return arrow.top


def _text_bottom(prs, sizes=(3000, 2600, 2200), spc_pct=0):
    """地の文4段落が実際に占める高さの下端（EMU）を、OOXML の規則で数える．"""
    slide = prs.slides[-1]
    body = _body(slide)
    y = body.top + body.text_frame.margin_top
    for lvl in (0, 1, 1, 1):
        pt = sizes[lvl] / 100.0
        y += Pt(pt * 1.32) + Pt(pt * spc_pct / 100000.0)
    return y


# ---------------------------------------------------------------- spcBef

def test_the_band_clears_the_prose_with_space_before(tmp_path):
    """段落前アキのあるテーマでも、図は地の文の下から始まる．"""
    theme = _theme(tmp_path, spc_pct=20000)
    prs = _build(tmp_path, _SRC, theme)
    assert _arrow_top(prs) >= _text_bottom(prs, spc_pct=20000)


def test_space_before_pushes_the_band_down(tmp_path):
    """アキのあるテーマのほうが、図は下から始まる．"""
    none_ = _build(tmp_path / "a", _SRC, _theme(tmp_path / "a"))
    with_ = _build(tmp_path / "b", _SRC, _theme(tmp_path / "b", spc_pct=20000))
    assert _arrow_top(with_) > _arrow_top(none_)


def test_absolute_space_before_counts_too(tmp_path):
    """``spcPts``（絶対値）も数える．"""
    none_ = _build(tmp_path / "a", _SRC, _theme(tmp_path / "a"))
    with_ = _build(tmp_path / "b", _SRC, _theme(tmp_path / "b", spc_pts=1200))
    assert _arrow_top(with_) > _arrow_top(none_)


# ---------------------------------------------------------------- レベル別

def _blanks(prs):
    """帯を空けるために入れた空段落の数．帯の高さがそのまま出る．"""
    return sum(1 for para in _body(prs.slides[-1]).text_frame.paragraphs
               if not para.text.strip())


def test_deep_levels_are_measured_at_their_own_size(tmp_path):
    """深い階層は小さい字なので、地の文が低く、帯が広くなる．

    どの行も lvl1（30pt）で数えていると、浅い版と深い版で帯の広さが同じになる。
    帯の広さは**空段落の数**にそのまま出るので、そこを見る。
    """
    shallow = (_FM + "### x\n\n- A\n- B\n- C\n- D\n\n"
               "```arrow\ndirection: down\n```\n\n→ 結論\n")
    deep = (_FM + "### x\n\n- A\n  - b\n  - c\n  - d\n\n"
            "```arrow\ndirection: down\n```\n\n→ 結論\n")
    a = _blanks(_build(tmp_path / "a", shallow, _theme(tmp_path / "a")))
    b = _blanks(_build(tmp_path / "b", deep, _theme(tmp_path / "b")))
    assert b > a, f"浅い版 {a} 空行 / 深い版 {b} 空行——同じなら lvl1 で数えている"


# ---------------------------------------------------------------- 回帰

def test_a_theme_without_space_before_is_unchanged(tmp_path):
    """アキの無いテーマでは従来どおり（既定テーマの出力を変えない）．"""
    theme = _theme(tmp_path)
    prs = _build(tmp_path, _SRC, theme)
    assert _arrow_top(prs) >= _text_bottom(prs)


def test_the_band_still_fits_the_frame(tmp_path):
    """帯を下げすぎて図が枠から出ない．"""
    theme = _theme(tmp_path, spc_pct=20000)
    prs = _build(tmp_path, _SRC, theme)
    slide = prs.slides[-1]
    body = _body(slide)
    arrow, = [sh for sh in slide.shapes
              if not sh.is_placeholder
              and getattr(sh, "auto_shape_type", None) is not None]
    assert arrow.top + arrow.height <= body.top + body.height


# ---------------------------------------------------------------- @autofit

def test_autofit_frees_room_for_the_band(tmp_path):
    """``@autofit`` で字を縮めたら、そのぶん帯が広がる．

    縮小率を見ていないと、空いた場所を帯が使えず図が結論文へ食い込む。
    """
    plain = _SRC
    shrunk = plain.replace("### x\n", "### x\n<!-- @autofit: 70 -->\n", 1)
    t1, t2 = _theme(tmp_path / "a", spc_pct=20000), _theme(tmp_path / "b", spc_pct=20000)
    a = _blanks(_build(tmp_path / "a", plain, t1))
    b = _blanks(_build(tmp_path / "b", shrunk, t2))
    assert b > a, f"縮小なし {a} 空行 / 縮小あり {b} 空行"


def test_autofit_moves_the_band_up(tmp_path):
    """縮めた版のほうが、図は上から始まる（地の文が低くなるので）．"""
    plain = _SRC
    shrunk = plain.replace("### x\n", "### x\n<!-- @autofit: 70 -->\n", 1)
    a = _arrow_top(_build(tmp_path / "a", plain,
                          _theme(tmp_path / "a", spc_pct=20000)))
    b = _arrow_top(_build(tmp_path / "b", shrunk,
                          _theme(tmp_path / "b", spc_pct=20000)))
    assert b < a
