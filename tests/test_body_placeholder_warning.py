#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""本文の置き場所が無いときに黙らないことを固定するテスト（Issue #67）．

md2pptx は本文プレースホルダを **``idx == 1``** で探す．この番号は PowerPoint が
内部で振るもので**画面に出ず，書き換えることもできない**．標準レイアウトを動かす
分には 1 のままだが，枠を消して作り直すと別の番号になりうる——テーマを自作する人が
普通にやる操作で，そうなると**地の文がどこにも入らない**．

修正前は，それが**警告も無く，終了コード 0** で起きていた．出力を開くまで気づけず，
気づいてもテーマが原因だとは思い当たらない．

止めないことは変えていない（表・図とタイトルは出せるし，pptx の保存自体は成功する）．
「PDF が作れなくても pptx は成功」（Issue #39）と同じ考え方で，**出せるものは出す．
ただし黙らない**——ここが Issue #67 で足した部分．

テーマは python-pptx 同梱の既定テンプレートから作るので，リポジトリの
``OfficeTheme.pptx`` にも実 PowerPoint にも依存しない．
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Pt

from md2pptx import render
from md2pptx.ir import Deck, Line, Slide, Table

_PH = "{http://schemas.openxmlformats.org/presentationml/2006/main}ph"


def _theme(tmp_path, body_idx=1):
    """本文プレースホルダの ``idx`` を選べるテーマ pptx を作って返す．

    ``body_idx`` を 1 以外にしたものが，このテストで再現したい「枠を作り直した
    テーマ」に相当する（PowerPoint 上での見た目は変わらない）．
    """
    prs = Presentation()
    for ph in prs.slide_layouts[1].placeholders:      # Title and Content
        if ph.placeholder_format.idx == 1:
            ph._element.find(f".//{_PH}").set("idx", str(body_idx))
    path = tmp_path / f"theme-{body_idx}.pptx"
    prs.save(str(path))
    return str(path)


def _build(theme, tmp_path, blocks, name="out.pptx"):
    """1 枚だけのデッキを描いて，出力パスを返す．"""
    out = tmp_path / name
    r = render.Renderer(theme)
    r.render(Deck(slides=[Slide(title="見出し", blocks=list(blocks))]))
    r.save(str(out))
    return out


def _texts(path):
    """出力に入っている文字列（プレースホルダ・図形を問わず）を集める．"""
    found = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                found.append(shape.text_frame.text)
    return "\n".join(found)


PROSE = [Line(text="導入文です。"), Line(text="結論行です。", kind="plain")]
TABLE = Table(header=["項目", "値"], rows=[["A", "1"]])


def test_a_normal_theme_says_nothing(tmp_path, capsys):
    """標準の番号なら地の文は本文枠へ入り，警告は出ない．"""
    out = _build(_theme(tmp_path), tmp_path, PROSE)

    assert "導入文です。" in _texts(out)
    assert capsys.readouterr().err == ""


def test_a_renamed_body_placeholder_is_reported(tmp_path, capsys):
    """番号が違うと地の文は消える——**そのことを言う**．

    修正前はここが無言だった．消えた事実と，何が消えたかを出す．
    """
    out = _build(_theme(tmp_path, body_idx=7), tmp_path, PROSE)

    assert "導入文です。" not in _texts(out), "前提: この状況では実際に消える"

    err = capsys.readouterr().err
    assert "no body placeholder" in err
    assert "導入文です。" in err, "何が消えたのか分からないと，直しようがない"


def test_the_same_goes_for_a_slide_with_a_table(tmp_path, capsys):
    """表・図のあるスライドは別の描画経路を通るので，そちらも確かめる．

    ``_render_stacked_into`` は docstring に「body が None の場合は地の文を捨て」と
    書いてあるとおり，意図して捨てている．捨てるのはよいが黙るのはよくない．
    """
    out = _build(_theme(tmp_path, body_idx=7), tmp_path,
                 [PROSE[0], TABLE, PROSE[1]])

    assert "導入文です。" not in _texts(out)
    assert "no body placeholder" in capsys.readouterr().err


def test_a_slide_with_nothing_to_drop_stays_quiet(tmp_path, capsys):
    """捨てる地の文が無ければ言わない．

    表だけのスライドは本文枠を使わない（むしろ消してから帯全体に置く）ので，
    ここで警告を出すと**正常な運用で毎回鳴る**ことになる．
    """
    _build(_theme(tmp_path, body_idx=7), tmp_path, [TABLE])

    assert capsys.readouterr().err == ""


def test_it_does_not_turn_a_build_into_a_failure(tmp_path, capsys):
    """警告であって失敗ではない——出せるものは出す．

    Issue #39 の「PDF が作れなくても pptx は成功」と同じ扱い．例外にすると，
    テーマを直すまで**タイトルも表も**手に入らなくなる．
    """
    out = _build(_theme(tmp_path, body_idx=7), tmp_path, [PROSE[0], TABLE])

    assert out.exists() and out.stat().st_size > 0
    assert "見出し" in _texts(out), "タイトルは出ている"
    assert "no body placeholder" in capsys.readouterr().err
