#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""本文の流れを示す大きな矢印を固定するテスト（Issue #134）．

「左の列 → 右の列 → その下の結論」を矢印の図形で見せる版が cn2025 にある
（cn2025-01 s13「インターネット以前の通信技術」・s16「☆共通基盤」）。
`flow` は箱と矢印の図で、箇条書きの列の間には置けない。

2つの記法を足す。

- ``<!-- @col: arrow -->`` — カラムの区切りを右向きの矢印として描く
- ``↓`` だけの行 — 下向きの矢印（**オブジェクトブロック**なので座標は帯が決める）
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

from md2pptx import render
from md2pptx.ir import Arrow
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"


def _slide(body):
    slide, = parse(_FM + "### x\n\n" + body + "\n").slides
    return slide


def _build(tmp_path, src):
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _arrows(slide, shape_type):
    return [sh for sh in slide.shapes
            if not sh.is_placeholder
            and getattr(sh, "auto_shape_type", None) == shape_type]


def _rights(slide):
    return _arrows(slide, MSO_SHAPE.RIGHT_ARROW)


def _downs(slide):
    return _arrows(slide, MSO_SHAPE.DOWN_ARROW)


# ---------------------------------------------------------------- 下向き矢印

@pytest.mark.parametrize("glyph", ["↓", "⇓"])
def test_an_arrow_only_line_becomes_a_block(glyph):
    """矢印1文字だけの行はオブジェクトブロックになる．"""
    blocks = _slide(f"- 上\n\n{glyph}\n\n→ 下").blocks
    assert any(isinstance(b, Arrow) for b in blocks)


def test_an_arrow_with_text_stays_a_line():
    """``→ 結論`` は従来どおり地の文（既存の書き方とぶつけない）．"""
    blocks = _slide("→ 結論の行").blocks
    assert not any(isinstance(b, Arrow) for b in blocks)
    assert blocks[0].text == "→ 結論の行"


def test_a_right_arrow_glyph_alone_is_not_a_block():
    """横向きはカラム区切りの仕事なので、``→`` 単独は矢印図形にしない．"""
    blocks = _slide("→").blocks
    assert not any(isinstance(b, Arrow) for b in blocks)


def test_the_down_arrow_is_drawn(tmp_path):
    prs = _build(tmp_path, _FM + "### x\n\n- 上\n\n↓\n\n→ 下\n")
    assert len(_downs(prs.slides[-1])) == 1


def test_the_down_arrow_sits_between_the_prose(tmp_path):
    """矢印は導入文より下、結論文より上に来る．"""
    prs = _build(tmp_path, _FM + "### x\n\n- 上\n\n↓\n\n→ 下\n")
    slide = prs.slides[-1]
    arrow, = _downs(slide)
    body, = [sh for sh in slide.shapes
             if sh.is_placeholder and sh.has_text_frame
             and sh != slide.shapes.title and "上" in sh.text_frame.text]
    assert body.top < arrow.top
    assert arrow.top + arrow.height <= body.top + body.height


def test_the_down_arrow_works_inside_a_column(tmp_path):
    """右カラムの中に置ける（cn2025-01 s13 の形）．"""
    src = _FM + "### x\n\n- 左\n\n<!-- @col -->\n\n- 右\n\n↓\n\n→ 結論\n"
    prs = _build(tmp_path, src)
    slide = prs.slides[-1]
    arrow, = _downs(slide)
    right, = [sh for sh in slide.shapes
              if sh.is_placeholder and sh.placeholder_format.idx == 2]
    assert right.left <= arrow.left < right.left + right.width


# ---------------------------------------------------------------- カラム間の矢印

def test_col_arrow_draws_a_right_arrow(tmp_path):
    src = _FM + "### x\n\n- 左\n\n<!-- @col: arrow -->\n\n- 右\n"
    prs = _build(tmp_path, src)
    assert len(_rights(prs.slides[-1])) == 1


def test_col_arrow_never_touches_the_right_column(tmp_path):
    """矢印は左カラムの右寄りに置き、**右カラムの本文には掛けない**．

    テーマのカラム間のすき間は 0.5cm ほどしか無く、そこへ収めると矢印が
    糸のように細くなる。元の講義スライドも、すき間ではなく左カラムの
    右寄りに置いてある。
    """
    src = _FM + "### x\n\n- 左\n\n<!-- @col: arrow -->\n\n- 右\n"
    prs = _build(tmp_path, src)
    slide = prs.slides[-1]
    arrow, = _rights(slide)
    left, = [sh for sh in slide.shapes
             if sh.is_placeholder and sh.placeholder_format.idx == 1]
    right, = [sh for sh in slide.shapes
              if sh.is_placeholder and sh.placeholder_format.idx == 2]
    assert arrow.left + arrow.width <= right.left
    assert arrow.left >= left.left + left.width // 2   # 左カラムの右半分にある


def test_col_arrow_is_not_a_hairline(tmp_path):
    """すき間の幅に引きずられて矢印が細くならない．"""
    src = _FM + "### x\n\n- 左\n\n<!-- @col: arrow -->\n\n- 右\n"
    prs = _build(tmp_path, src)
    arrow, = _rights(prs.slides[-1])
    assert arrow.width >= 457200      # 0.5in 以上


def test_plain_col_draws_nothing(tmp_path):
    """``<!-- @col -->`` の出力は変わらない．"""
    src = _FM + "### x\n\n- 左\n\n<!-- @col -->\n\n- 右\n"
    prs = _build(tmp_path, src)
    assert _rights(prs.slides[-1]) == []


def test_col_still_splits_the_columns():
    """値を付けてもカラム区切りとしての働きは同じ．"""
    slide = _slide("- 左\n\n<!-- @col: arrow -->\n\n- 右")
    assert len(slide.columns) == 2


def test_an_unknown_col_value_stops():
    """``@col: 2`` のような値はエラーで止まる（従来どおり）．"""
    with pytest.raises(ValueError, match="@col"):
        _slide("- 左\n\n<!-- @col: 2 -->\n\n- 右")


def test_col_arrow_needs_two_columns(tmp_path, capsys):
    """カラムが 1 つしかなければ矢印は描かず、警告する．"""
    src = _FM + "### x\n\n- 左だけ\n"
    _build(tmp_path, src)
    assert "@col: arrow" not in capsys.readouterr().err
