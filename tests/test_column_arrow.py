#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""カラム区切りの矢印 ``<!-- @col: arrow -->`` を固定するテスト（Issue #134）．

「左の列 → 右の列」を大きな矢印で見せる版が cn2025 にある
（cn2025-01 s13「インターネット以前の通信技術」・s16「☆共通基盤」）。
`flow` は箱と矢印の図で、箇条書きの列の間には置けない。

カラム**内**に置く矢印は ``` ```arrow ``` フェンス（Issue #137）で、
そちらは `tests/test_arrow_block.py` が持つ。ここは区切りの矢印だけ。
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

from md2pptx import render
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"


def _slide(body):
    slide, = parse(_FM + "### x\n\n" + body + "\n").slides
    return slide


def _build(tmp_path, src):
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _arrows(slide, shape_type):
    return [sh for sh in slide.shapes
            if not sh.is_placeholder
            and getattr(sh, "auto_shape_type", None) == shape_type]


def _rights(slide):
    return _arrows(slide, MSO_SHAPE.RIGHT_ARROW)


# ---------------------------------------------------------------- カラム間の矢印

def test_col_arrow_draws_a_right_arrow(tmp_path):
    src = _FM + "### x\n\n- 左\n\n<!-- @col: arrow -->\n\n- 右\n"
    prs = _build(tmp_path, src)
    assert len(_rights(prs.slides[-1])) == 1


def test_col_arrow_never_touches_the_right_column(tmp_path):
    """矢印は左カラムの右寄りに置き、**右カラムの本文には掛けない**．

    テーマのカラム間のすき間は 0.5cm ほどしか無く、そこへ収めると矢印が
    糸のように細くなる。元の講義スライドも、すき間ではなく左カラムの
    右寄りに置いてある。
    """
    src = _FM + "### x\n\n- 左\n\n<!-- @col: arrow -->\n\n- 右\n"
    prs = _build(tmp_path, src)
    slide = prs.slides[-1]
    arrow, = _rights(slide)
    left, = [sh for sh in slide.shapes
             if sh.is_placeholder and sh.placeholder_format.idx == 1]
    right, = [sh for sh in slide.shapes
              if sh.is_placeholder and sh.placeholder_format.idx == 2]
    assert arrow.left + arrow.width <= right.left
    assert arrow.left >= left.left + left.width // 2   # 左カラムの右半分にある


def test_col_arrow_is_not_a_hairline(tmp_path):
    """すき間の幅に引きずられて矢印が細くならない．"""
    src = _FM + "### x\n\n- 左\n\n<!-- @col: arrow -->\n\n- 右\n"
    prs = _build(tmp_path, src)
    arrow, = _rights(prs.slides[-1])
    assert arrow.width >= 457200      # 0.5in 以上


def test_plain_col_draws_nothing(tmp_path):
    """``<!-- @col -->`` の出力は変わらない．"""
    src = _FM + "### x\n\n- 左\n\n<!-- @col -->\n\n- 右\n"
    prs = _build(tmp_path, src)
    assert _rights(prs.slides[-1]) == []


def test_col_still_splits_the_columns():
    """値を付けてもカラム区切りとしての働きは同じ．"""
    slide = _slide("- 左\n\n<!-- @col: arrow -->\n\n- 右")
    assert len(slide.columns) == 2


def test_an_unknown_col_value_stops():
    """``@col: 2`` のような値はエラーで止まる（従来どおり）．"""
    with pytest.raises(ValueError, match="@col"):
        _slide("- 左\n\n<!-- @col: 2 -->\n\n- 右")


def test_a_single_column_slide_is_untouched(tmp_path, capsys):
    """カラム区切りの無いスライドには何も起きない（警告も出さない）．

    ``@col: arrow`` はカラム区切りそのものなので、パーサを通る限りカラムは
    必ず 2 つ以上ある。描画側の ``ncols < 2`` は直接呼ばれたときの防御で、
    起きないことに警告を出しても読まれない。
    """
    prs = _build(tmp_path, _FM + "### x\n\n- 左だけ\n")
    assert _rights(prs.slides[-1]) == []
    assert capsys.readouterr().err == ""
