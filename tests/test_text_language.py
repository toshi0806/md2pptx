#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""描画した run に言語が付くことを固定するテスト（Issue #79）．

PowerPoint は**行分割の規則を run の言語で選ぶ**．python-pptx は
``paragraph.text = …`` で ``<a:r><a:t>…</a:t></a:r>`` を作り ``a:rPr`` を一切
書かないので，md2pptx が何もしないと言語不明のまま出力される．そうなると日本語の
禁則処理が適用されず，**行頭に「ー」や句読点が来る**．

文字列としては正しいので pptx を開くまで気づけず，開いても原因が md2pptx だとは
思い当たらない——だから「見た目が崩れる」ではなく「``lang`` が付いている」を固定する．

実 PowerPoint 変換で確かめた切り分けもここに残す．効くのは ``lang`` **だけ**で，

- ``kumimoji="1"`` のみでは禁則は効かない（縦書き中の数字の扱いであって行分割ではない）
- ``presentation.xml`` の ``<p:kinsoku>``（禁則文字の定義）を足しても効かない

という関係だった．テストからは実 PowerPoint を叩けないので，ここで押さえるのは
「``lang`` が全 run に付く」こと．これが崩れたら禁則も崩れる．

テーマは python-pptx 同梱の既定テンプレートから作るので，リポジトリの
``OfficeTheme.pptx`` にも実 PowerPoint にも依存しない．
"""
from __future__ import annotations

from pptx import Presentation

from md2pptx import render
from md2pptx.ir import Deck, Line, Slide, Table, TitleSlide
from md2pptx.parser import parse

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _theme(tmp_path):
    path = tmp_path / "theme.pptx"
    Presentation().save(str(path))
    return str(path)


def _theme_with_its_own_run(tmp_path):
    """番号プレースホルダに**言語付きの run** を持つテーマを作って返す．

    これが「テーマ由来の run」がスライドへ入る実際の経路．``add_slide_number`` は
    レイアウトの番号プレースホルダ（``idx == 12``）を ``deepcopy`` してスライドへ
    足すので，レイアウト側の run がそのまま出力に現れる．既定テンプレートの番号枠は
    ``a:fld``（番号フィールド）だけだが，「Page 3」のように語を添えるテーマでは
    ``a:r`` が並ぶ——そこに ``lang`` が付いていることも普通にある．
    """
    prs = Presentation()
    for lay in prs.slide_layouts:
        for ph in lay.placeholders:
            if ph.placeholder_format.idx != 12:
                continue
            p_el = ph.text_frame.paragraphs[0]._p
            fld = p_el.find(f"{_A}fld")
            if fld is None:
                # 番号フィールドが無い形の既定テンプレートに変わったら，run を
                # 置く位置が決まらない．黙って素通りしてよい——run が 1 つも
                # 入らなければテスト側の「前提」アサーションが落ちて気づける．
                continue
            # lang を持つ run と，altLang だけを持つ run の 2 つを入れる．
            # 後者は珍しいが，「決まっているものは触らない」を altLang についても
            # 効かせているかを見るために要る．
            for text, attrs in (("Page ", {"lang": "en-US"}),
                                ("/ ", {"altLang": "ko-KR"})):
                run = p_el.makeelement(f"{_A}r", {})
                run.append(run.makeelement(f"{_A}rPr", attrs))
                t = run.makeelement(f"{_A}t", {})
                t.text = text
                run.append(t)
                p_el.insert(list(p_el).index(fld), run)
    path = tmp_path / "theme-with-run.pptx"
    prs.save(str(path))
    return str(path)


def _build(tmp_path, deck, name="out.pptx", theme=None):
    out = tmp_path / name
    r = render.Renderer(theme or _theme(tmp_path))
    r.render(deck)
    r.save(str(out))
    return out


def _runs(path, attr="lang"):
    """出力の全 run を (本文, その属性値) で返す（ノートスライドも含む）．

    ``r.prs`` を直接見ずに**保存した pptx を読み直す**のは意図的．利用者が開くのは
    ファイルであって，メモリ上のツリーではない．保存の往復で落ちる属性があれば
    禁則も効かないので，そこまで含めて固定する．

    ``a:rPr`` ごと無い run は ``None`` を返す——**それも「言語が決まっていない」**で，
    呼び出し側は属性が空の場合と区別しなくてよい（どちらでも禁則は効かない）．
    """
    found = []
    for slide in Presentation(str(path)).slides:
        parts = [slide.element]
        if slide.has_notes_slide:
            parts.append(slide.notes_slide.element)
        for part in parts:
            for r in part.iter(f"{_A}r"):
                t = r.find(f"{_A}t")
                rPr = r.find(f"{_A}rPr")
                found.append((t.text if t is not None else "",
                              rPr.get(attr) if rPr is not None else None))
    return found


DECK = Deck(
    title_slide=TitleSlide(title="表題", author="著者", affiliation=["所属"]),
    slides=[Slide(title="見出し",
                  blocks=[Line(text="コンテンツ配信のためのネットワーク"),
                          Line(text="→ 結論行", kind="plain"),
                          Table(header=["項目", "値"], rows=[["表の中", "1"]])],
                  notes="ノートの本文")],
)


def test_every_run_gets_a_language(tmp_path):
    """タイトル・本文・表・ノートまで，md2pptx が作った run すべてに付く．

    run を作る経路は多い．どれか 1 つでも漏れるとその枠だけ禁則が効かない
    ——**枠ごとに崩れ方が違う**ので原因を追いにくい．
    """
    runs = _runs(_build(tmp_path, DECK))

    assert runs, "前提: run が 1 つも無ければこのテストは何も見ていない"

    # 網羅：言語の決まらない run を 1 つも残さない．テーマが何を持ち込んでも
    # 誤検知しない——テーマ由来の run も，言語を持たなければここで ja-JP が付く．
    # 残った None は md2pptx の取りこぼしに限られる．
    missing = [text for text, lang in runs if lang is None]
    assert missing == [], f"言語の付いていない run が残っている: {missing}"

    # 個別：描画経路ごとに名指しで ja-JP を確かめる．こちらもテーマ由来 run の
    # 有無に左右されない（既定テンプレートが将来 run を持っても影響しない）．
    by_text = {}
    for text, lang in runs:
        by_text.setdefault(text, []).append(lang)
    for expected in ("表題", "著者", "所属", "見出し",
                     "コンテンツ配信のためのネットワーク", "→ 結論行",
                     "表の中", "ノートの本文"):
        # 同じ本文が複数の run に現れても（表のセルと本文が同じ文字列など）
        # 落ちないよう，個数ではなく全要素を見る．空でないことは直前の in が担保する．
        assert expected in by_text, f"{expected!r} を通る経路が見えていない"
        assert all(lang == "ja-JP" for lang in by_text[expected]), \
            f"{expected!r} の run の言語: {by_text[expected]}"


def test_it_does_not_overwrite_a_language_already_set(tmp_path):
    """テーマが自分で持ってきた run の言語は書き換えない．

    未設定のものだけ埋めるので**何度通しても結果が変わらない**．run 単位で言語を
    決めるようになったときも，後から通るこの処理が決定を上書きしない．

    テーマの決めた言語を ja-JP で塗り潰すと，そのテーマだけ番号や見出しの語が
    別の言語として扱われる——**テーマに委ねる**という本ツールの方針にも反する．
    """
    out = _build(tmp_path, DECK, name="theme-run.pptx",
                 theme=_theme_with_its_own_run(tmp_path))
    runs = _runs(out)

    # dict 化すると同じ本文の run が複数あったとき後勝ちで別の run を見てしまう．
    # テーマ由来の run だけを名指しで取る．
    theirs = [lang for text, lang in runs if text == "Page "]
    assert theirs, "前提: テーマ由来の run が出力に入っていない"
    assert set(theirs) == {"en-US"}, f"テーマの決めた言語を書き換えた: {theirs}"

    # 同じ出力の中で，md2pptx が描いた run にはちゃんと ja-JP が付いている．
    assert next(lang for text, lang in runs if text == "見出し") == "ja-JP"

    # altLang も同じ規則．テーマが決めていれば残し，決めていなければ埋める．
    alt = _runs(out, attr="altLang")
    assert set(a for text, a in alt if text == "/ ") == {"ko-KR"}, \
        "テーマの決めた altLang を書き換えた"
    assert next(a for text, a in alt if text == "見出し") == "en-US"

    # altLang だけ持つ run にも lang は付く．**これが狙いどおり**——禁則を選ぶのは
    # lang で，altLang がいくら決まっていても lang が空なら行分割の規則が決まらない．
    assert set(lang for text, lang in runs if text == "/ ") == {"ja-JP"}


def test_br_in_front_matter_becomes_a_line_break(tmp_path):
    """front matter の ``<br>`` は行内改行になる（見出し・本文と同じ規則．Issue #79）．

    修正前はここだけ素通しで，``"<br>"`` の 4 文字がそのまま画面に出ていた．
    タイトルが枠幅に収まらないときの折り位置を著者が決める手段でもある．
    """
    deck = parse("---\n"
                 "title: DoHを用いた問い合わせ先<br>指定による\n"
                 "subtitle: 副題<br>つづき\n"
                 "author: 著者<br>ふりがな\n"
                 "affiliation:\n"
                 "  - 所属<br>2 行目\n"
                 "---\n")

    ts = deck.title_slide
    assert ts.title == "DoHを用いた問い合わせ先\v指定による"
    assert ts.subtitle == "副題\vつづき"
    assert ts.author == "著者\vふりがな"
    assert ts.affiliation == ["所属\v2 行目"]

    # \v は python-pptx が <a:br/> へ展開する（"<br>" の文字が残らないことまで見る）．
    out = _build(tmp_path, deck, name="br.pptx")
    slide = Presentation(str(out)).slides[0]
    body = "\n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    assert "<br>" not in body
    assert "問い合わせ先\v指定による" in body
    assert any(True for _ in slide.element.iter(f"{_A}br")), "<a:br/> が無い"


def test_a_size_token_still_works_next_to_br(tmp_path):
    """``{-1}`` の剥がしと ``<br>`` の変換は両立する（順序を取り違えない）．"""
    deck = parse("---\ntitle: 表題\nsubtitle: '{-1}副題<br>つづき'\n---\n")

    assert deck.title_slide.subtitle_delta == -1
    assert deck.title_slide.subtitle == "副題\vつづき"
