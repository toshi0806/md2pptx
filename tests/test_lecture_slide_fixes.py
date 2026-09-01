#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""cn2026 の講義スライドで足りなかった 3 つを固定するテスト．

いずれも「書けないので原稿の側を曲げていた」ものを、書けるようにしたもの:

- ``` ```arrow ``` の ``align:`` — 2 カラムの片側に置いた矢印が、左寄せの項目に
  対して中央のままで右へずれて見えた．
- ``> 本文`` — 記号の出ない本文行．図形の矢印の直下に ``→ 本文`` と書くと矢印が
  二重に見え、``- `` にするとビュレットが付いて項目と同列に見えた．
- ``@title-width`` — 章扉のタイトルが枠に収まらず 1 文字だけ次行へ落ちた．
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

from md2pptx import render
from md2pptx.ir import Arrow
from md2pptx.parser import parse, parse_content_line


_FM = "---\ntheme: t.pptx\n---\n\n"


def _theme(tmp_path):
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    if not theme.exists():
        Presentation().save(str(theme))
    return theme


def _build(tmp_path, src):
    theme = _theme(tmp_path)
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _fence(**kw):
    return "```arrow\n" + "\n".join(f"{k}: {v}" for k, v in kw.items()) + "\n```"


def _arrow_block(body):
    slide, = parse(_FM + "### x\n\n" + body + "\n").slides
    return [b for b in slide.blocks if isinstance(b, Arrow)][0]


def _arrow_shape(prs):
    got = [sh for sh in prs.slides[-1].shapes
           if not sh.is_placeholder
           and getattr(sh, "auto_shape_type", None) == MSO_SHAPE.DOWN_ARROW]
    assert len(got) == 1, f"矢印が {len(got)} 個"
    return got[0]


# ---------------------------------------------------------------- arrow の align

def test_align_defaults_to_center():
    """書かなければ従来どおり中央．"""
    assert _arrow_block(_fence(direction="down")).align == "center"


@pytest.mark.parametrize("value", ["left", "center", "right"])
def test_align_takes_the_image_vocabulary(value):
    """語彙は ``` ```image ``` と共通（left / center / right）．"""
    assert _arrow_block(_fence(direction="down", align=value)).align == value


def test_an_unknown_align_is_an_error():
    """知らない値はタイポとみなして止める（他のキーと同じ）．"""
    with pytest.raises(ValueError, match="align"):
        _arrow_block(_fence(direction="down", align="middle"))


def test_align_left_moves_the_arrow_to_the_left_edge(tmp_path):
    """``align: left`` は帯の左端に置く（中央より必ず左）．"""
    body = "### x\n\n- 上\n\n{}\n\n- 下\n"
    left = _arrow_shape(_build(
        tmp_path / "l", _FM + body.format(_fence(direction="down", align="left"))))
    center = _arrow_shape(_build(
        tmp_path / "c", _FM + body.format(_fence(direction="down"))))
    assert left.left < center.left
    assert left.width == center.width      # 寄せは大きさを変えない


def test_align_right_moves_the_arrow_to_the_right_edge(tmp_path):
    """``align: right`` は帯の右端に置く．"""
    body = "### x\n\n- 上\n\n{}\n\n- 下\n"
    right = _arrow_shape(_build(
        tmp_path / "r", _FM + body.format(_fence(direction="down", align="right"))))
    center = _arrow_shape(_build(
        tmp_path / "c", _FM + body.format(_fence(direction="down"))))
    assert right.left > center.left


# ---------------------------------------------------------------- 記号なしの行

def test_a_quoted_line_keeps_no_marker():
    """``> 本文`` は記号なしの段落．**記号は本文に残さない**．"""
    line = parse_content_line("> 全てアナログ情報")
    assert (line.kind, line.text) == ("plain", "全てアナログ情報")


def test_an_arrow_line_still_shows_its_arrow():
    """``→`` は従来どおり本文に残る（導線として見せる書き方）．"""
    line = parse_content_line("→ 残る矢印")
    assert (line.kind, line.text) == ("plain", "→ 残る矢印")


def test_a_marker_needs_its_space():
    """``>=`` のように記号が続く行は本文のまま（誤って記号を食わない）．

    空白を必須にしていないと、書いた覚えのない行で先頭の 1 文字が消える．
    """
    line = parse_content_line(">= 3 のとき")
    assert (line.kind, line.text) == ("bullet", ">= 3 のとき")


def test_a_bare_marker_is_an_empty_plain_line():
    """``>`` だけの行は記号の出ない空段落（``-`` 単独と同じくスペーサ）．"""
    line = parse_content_line(">")
    assert (line.kind, line.text) == ("plain", "")


def test_the_quoted_line_reaches_the_slide(tmp_path):
    """描画まで通る（本文プレースホルダの段落として出る）．"""
    prs = _build(tmp_path, _FM + "### x\n\n- 項目\n\n> 地の文\n")
    texts = [p.text for sh in prs.slides[-1].shapes if sh.has_text_frame
             for p in sh.text_frame.paragraphs]
    assert "地の文" in texts


# ---------------------------------------------------------------- @title-width

def _title(prs):
    return prs.slides[-1].shapes.title


def test_title_width_widens_the_title_placeholder(tmp_path):
    """``@title-width`` はタイトル枠を広げる（本文枠は動かさない）．"""
    src = _FM + "### 第2回：ネットワークコミュニケーション\n\n{}- 項目\n"
    plain = _build(tmp_path / "p", src.format(""))
    wide = _build(tmp_path / "w", src.format("<!-- @title-width: 120 -->\n\n"))
    assert _title(wide).width > _title(plain).width
    assert _title(wide).left == _title(plain).left      # 左端は動かさない


def test_title_width_does_not_touch_the_body(tmp_path):
    """本文プレースホルダは ``@widths`` の担当（混ざらない）．"""
    src = _FM + "### t\n\n{}- 項目\n"
    plain = _build(tmp_path / "p", src.format(""))
    wide = _build(tmp_path / "w", src.format("<!-- @title-width: 120 -->\n\n"))

    def body_width(prs):
        for sh in prs.slides[-1].shapes:
            if sh.is_placeholder and sh.placeholder_format.idx == 1:
                return sh.width
        raise AssertionError("本文プレースホルダが無い")

    assert body_width(wide) == body_width(plain)


def test_title_width_is_clamped_to_the_right_margin(tmp_path, capsys):
    """右余白を越える指定はクランプして警告する（``@widths`` と同じ規約）．"""
    prs = _build(tmp_path, _FM + "### t\n\n<!-- @title-width: 400 -->\n\n- 項目\n")
    err = capsys.readouterr().err
    assert "@title-width" in err and "clamping" in err
    assert _title(prs).left + _title(prs).width <= prs.slide_width


def test_a_bad_title_width_is_ignored_with_a_warning(tmp_path, capsys):
    """読めない値・非正の値は無視して警告（黙って別の幅にしない）．"""
    plain = _build(tmp_path / "p", _FM + "### t\n\n- 項目\n")
    bad = _build(tmp_path / "b",
                 _FM + "### t\n\n<!-- @title-width: -10 -->\n\n- 項目\n")
    assert "@title-width" in capsys.readouterr().err
    assert _title(bad).width == _title(plain).width
