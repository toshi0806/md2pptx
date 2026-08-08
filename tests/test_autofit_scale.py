#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""``default_autofit`` が実際に縮めることを固定するテスト（Issue #154）．

`auto_size = TEXT_TO_FIT_SHAPE` を設定すると pptx には
``<a:bodyPr><a:normAutofit/></a:bodyPr>`` と書かれるが、**PowerPoint は開いた
ときに自動調整を計算し直さない**。保存されている ``fontScale``（無ければ 100%）で
描く。つまり「自動調整の枠だ」と印を付けるだけでは、変換した PDF では何も起きず、
本文がプレースホルダを越えて下の罫線に掛かる（cn2026-02 p.17）。

md2pptx 自身が縮小率を出して焼き込む。同じ率は ``{box}`` の位置にも要る——
縮小率を知らないと枠が字とずれる。
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from md2pptx import render
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"


def _build(tmp_path, src, theme=None):
    tmp_path.mkdir(parents=True, exist_ok=True)
    if theme is None:
        theme = tmp_path / "theme.pptx"
        Presentation().save(str(theme))
    r = render.Renderer(str(theme))
    r.render(parse(src))
    out = tmp_path / "out.pptx"
    r.save(str(out))
    return Presentation(str(out)), r


def _body(prs, idx=1):
    for sh in prs.slides[-1].shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == idx:
            return sh
    raise AssertionError(f"idx={idx} のプレースホルダが無い")


def _font_scale(shape):
    """焼き込まれた縮小率（％）．無ければ None．"""
    na = shape.text_frame._txBody.bodyPr.find(qn("a:normAutofit"))
    if na is None or na.get("fontScale") is None:
        return None
    return int(na.get("fontScale")) / 1000.0


def _many(n, level=0):
    ind = "  " * level
    return "".join(f"- {ind}項目{i}\n" for i in range(n))


# ---------------------------------------------------------------- 縮める

def test_a_crowded_slide_gets_a_scale(tmp_path):
    """入り切らない枚数を書いたら、縮小率が焼き込まれる．"""
    prs, _ = _build(tmp_path, _FM + "### x\n\n" + _many(20))
    scale = _font_scale(_body(prs))
    assert scale is not None, "fontScale が焼き込まれていない"
    assert 10 < scale < 100


def test_a_roomy_slide_is_not_shrunk(tmp_path):
    """収まるスライドは縮めない（従来の見た目を変えない）．"""
    prs, _ = _build(tmp_path, _FM + "### x\n\n" + _many(3))
    assert _font_scale(_body(prs)) is None


def test_more_lines_shrink_more(tmp_path):
    a = _font_scale(_body(_build(tmp_path / "a", _FM + "### x\n\n" + _many(20))[0]))
    b = _font_scale(_body(_build(tmp_path / "b", _FM + "### x\n\n" + _many(40))[0]))
    assert b < a


def test_the_text_fits_after_shrinking(tmp_path):
    """縮めた率で数え直すと、プレースホルダに収まっている．"""
    src = _FM + "### x\n\n" + _many(20)
    prs, r = _build(tmp_path, src)
    from md2pptx.parser import parse as _parse
    body = _body(prs)
    scale = _font_scale(body) / 100.0
    tf = body.text_frame
    avail_h = body.height - tf.margin_top - tf.margin_bottom
    avail_w = body.width - tf.margin_left - tf.margin_right
    # 折り返しも含めて数え直す——md2pptx が実際に使っている式で確かめる
    lines = [b for b in _parse(src).slides[-1].blocks]
    need = r._text_height(lines, r._frame_font_levels(tf), None, avail_w, scale)
    assert need <= avail_h


# ---------------------------------------------------------------- 明示指定

def test_an_explicit_autofit_still_wins(tmp_path):
    """``@autofit: 60`` と書いたら、計算より書いたほうを採る．"""
    src = _FM + "### x\n<!-- @autofit: 60 -->\n\n" + _many(20)
    prs, _ = _build(tmp_path, src)
    assert _font_scale(_body(prs)) == pytest.approx(60.0)


def test_an_explicit_autofit_on_a_roomy_slide(tmp_path):
    """収まるスライドでも、書いてあれば縮める．"""
    src = _FM + "### x\n<!-- @autofit: 80 -->\n\n" + _many(3)
    prs, _ = _build(tmp_path, src)
    assert _font_scale(_body(prs)) == pytest.approx(80.0)


def test_turning_the_default_off(tmp_path):
    """``default_autofit: false`` なら計算もしない．"""
    src = ("---\ntheme: t.pptx\ndefault_autofit: false\n---\n\n"
           "### x\n\n" + _many(20))
    prs, _ = _build(tmp_path, src)
    assert _font_scale(_body(prs)) is None


# ---------------------------------------------------------------- カラム

def test_each_column_is_measured_on_its_own(tmp_path):
    """カラムごとに別の枠なので、率も別に出す．"""
    src = _FM + "### x\n\n" + _many(20) + "\n<!-- @col -->\n\n" + _many(2)
    prs, _ = _build(tmp_path, src)
    assert _font_scale(_body(prs, 1)) is not None
    assert _font_scale(_body(prs, 2)) is None


# ---------------------------------------------------------------- {box}

def _boxes(prs):
    return [sh for sh in prs.slides[-1].shapes
            if not sh.is_placeholder
            and getattr(sh, "auto_shape_type", None) == MSO_SHAPE.ROUNDED_RECTANGLE]


def test_the_box_follows_the_shrunken_text(tmp_path):
    """縮んだぶん、枠も上へ来る．

    縮小率を見ないと、枠だけが元の大きさの位置に取り残される（枠が下へ寄る）。
    """
    plain = _FM + "### x\n\n" + _many(3) + "- {box} 囲む行\n"
    crowd = _FM + "### x\n\n" + _many(3) + "- {box} 囲む行\n" + _many(20)
    a, _ = _build(tmp_path / "a", plain)
    b, _ = _build(tmp_path / "b", crowd)
    (ba,), (bb,) = _boxes(a), _boxes(b)
    assert _font_scale(_body(b)) is not None, "前提が崩れた（縮んでいない）"
    assert bb.top < ba.top
    assert bb.height < ba.height        # 字が小さくなれば枠も低い
