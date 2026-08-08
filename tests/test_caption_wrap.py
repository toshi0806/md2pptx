#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""折り返すキャプションの高さを固定するテスト（Issue #169）．

``render_image`` はキャプションの高さを**1行ぶん決め打ち**で確保していたのに、
``_draw_caption`` は ``word_wrap = True`` で描く。長いキャプションの 2 行目は
確保の外——図の下端より下、つまり罫線やページ番号の側へ出ていた（cn2026-09 p.6）。

キャプションはもともと図の取り分を削るので、**長く書けば図が小さくなる**。
それは書いた人に見える。罫線に潜り込むより分かりやすい壊れ方。
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from md2pptx import render
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"


def _fig(tmp_path, name="fig.png", size=(800, 600)):
    from PIL import Image
    tmp_path.mkdir(parents=True, exist_ok=True)
    p = tmp_path / name
    Image.new("RGB", size, "white").save(p)
    return p


def _build(tmp_path, caption, size=(800, 600)):
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    img = _fig(tmp_path, size=size)
    src = (_FM + "### x\n\n```image\nsrc: %s\nwidth: 100%%\ncaption: %s\n```\n"
           % (img.name, caption))
    r = render.Renderer(str(theme), base_dir=str(tmp_path))
    r.render(parse(src))
    out = tmp_path / "out.pptx"
    r.save(str(out))
    return Presentation(str(out)), r


def _pic(prs):
    pics = [sh for sh in prs.slides[-1].shapes
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pics, "図が置かれていない"
    return pics[0]


def _caption(prs):
    boxes = [sh for sh in prs.slides[-1].shapes
             if not sh.is_placeholder and sh.has_text_frame
             and sh.text_frame.text.strip()]
    assert boxes, "キャプションが無い"
    return boxes[0]


_SHORT = "短い説明"
_LONG = "とても長い説明の文がここに入ります" * 3      # 確実に折り返す


def test_a_wrapped_caption_gets_room_for_every_line(tmp_path):
    """折り返すキャプションは、行数ぶんの高さを確保する．"""
    prs, r = _build(tmp_path, _LONG)
    cap = _caption(prs)
    tf = cap.text_frame
    avail_pt = (cap.width - tf.margin_left - tf.margin_right) / 12700.0
    size = r._caption_size()
    n = r._wrapped_lines(_LONG, size, avail_pt)
    assert n > 1, "前提が崩れた（キャプションが折り返さない）"
    assert cap.height >= n * r._line_height(size)


def test_the_caption_stays_above_the_bottom(tmp_path):
    """キャプションの下端が、図の置き場の下端を越えない．"""
    prs, _ = _build(tmp_path, _LONG)
    cap = _caption(prs)
    body = [sh for sh in prs.slides[-1].shapes
            if sh.is_placeholder and sh.placeholder_format.idx == 1]
    limit = (body[0].top + body[0].height) if body else prs.slide_height
    assert cap.top + cap.height <= limit


def test_a_long_caption_shrinks_the_figure(tmp_path):
    """長いキャプションを書けば、そのぶん図が小さくなる（見えて分かる）．"""
    a, _ = _build(tmp_path / "a", _SHORT)
    b, _ = _build(tmp_path / "b", _LONG)
    assert _pic(b).height < _pic(a).height


def test_a_short_caption_is_unchanged(tmp_path):
    """1行に収まるキャプションでは従来どおり（回帰させない）．"""
    prs, r = _build(tmp_path, _SHORT)
    cap = _caption(prs)
    size = r._caption_size()
    assert cap.height < 2 * r._line_height(size)


def test_no_caption_gives_the_figure_everything(tmp_path):
    """キャプションが無ければ図が最大（回帰させない）．"""
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    img = _fig(tmp_path)
    src = _FM + "### x\n\n```image\nsrc: %s\nwidth: 100%%\n```\n" % img.name
    r = render.Renderer(str(theme), base_dir=str(tmp_path))
    r.render(parse(src))
    out = tmp_path / "out.pptx"
    r.save(str(out))
    plain = Presentation(str(out))
    withcap, _ = _build(tmp_path / "c", _SHORT)
    assert _pic(plain).height > _pic(withcap).height


def test_a_borderline_caption_is_counted_as_wrapping(tmp_path):
    """幅ぎりぎりのキャプションは「折り返す」側に数える．

    折り返し判定の許容幅（``_WRAP_SLACK``）は ``{box}`` のために「折り返さない」
    へ倒してある。キャプションは**逆向き**——1 行多く取っても図が少し小さく
    なるだけだが、足りないと 2 行目が罫線の下へ出る（Issue #169）。
    """
    prs, r = _build(tmp_path, _SHORT)
    cap = _caption(prs)
    tf = cap.text_frame
    avail_pt = (cap.width - tf.margin_left - tf.margin_right) / 12700.0
    size = r._caption_size()
    # 使える幅を少しだけ超える文字列を作る（許容幅 5% の内側に収まる長さ）
    unit = r._text_width_pt("あ", size)
    text = "あ" * int(avail_pt * 1.03 / unit)
    w = r._text_width_pt(text, size)
    assert avail_pt < w <= avail_pt * 1.05, "前提が崩れた（狙った幅にならない）"
    assert r._wrapped_lines(text, size, avail_pt) == 1, "前提が崩れた（許容幅が効かない）"
    assert r._wrapped_lines(text, size, avail_pt, slack=1.0) == 2
    assert r._caption_height(text, cap.width) >= 2 * r._line_height(size)
