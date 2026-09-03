#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""フェンスドコードブロックを固定するテスト（Issue #106）．

``flow`` / ``image`` / ``note`` 以外のフェンスは**跡形もなく捨てられていた**
（警告すら出なかった）。他の記法が全部「未知はタイポとみなしてエラー」なのと
非対称で、``SYNTAX.md`` にも「捨てる」とは書かれていない。

講義スライドでは telnet の HTTP / SMTP セッション、``dig`` の出力、2進数の
ネットマスク列といった「等幅で、書いたまま出したい塊」が要る。これらは
行頭マーカー記法では書けない（``-`` や ``1.`` に見える行が化ける）。

**構文強調はしない。** 色をテーマに委ねるという方針と噛み合わないし、
講義で見せたいのはセッションの往復であって色ではない。info string
（``` ```console ``` の ``console``）は書けるが、md2pptx は読み飛ばす。

行は**書いたまま**出す——行頭マーカーもサイズトークンも ``<br>`` も解釈しない。
コードブロックの中で ``- `` が箇条書きになったら、それはもうコードではない。
"""
from __future__ import annotations

import pytest
from pptx import Presentation

from md2pptx import render
from md2pptx.ir import Line
from md2pptx.parser import parse

_FM = "---\ntheme: t.pptx\n---\n\n"
_FENCE = "```"


def _build(tmp_path, src):
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    return Presentation(str(out))


def _body_paragraphs(slide):
    """本文枠の (テキスト, フォント名) を段落順に返す．"""
    for sh in slide.shapes:
        if sh.has_text_frame and sh != slide.shapes.title:
            return [("".join(r.text for r in p.runs),
                     next((r.font.name for r in p.runs if r.font.name), None))
                    for p in sh.text_frame.paragraphs]
    return []


def test_a_code_block_is_no_longer_dropped():
    """コードブロックが本文として残る（これが直したかったもの）．"""
    src = _FM + f"### HTTP\n\n{_FENCE}text\nGET / HTTP/1.1\nHost: example.com\n{_FENCE}\n"
    slide, = parse(src).slides
    assert [b.text for b in slide.blocks if isinstance(b, Line)] == [
        "GET / HTTP/1.1", "Host: example.com"]


def test_every_line_is_a_paragraph_of_kind_code():
    """1 行が 1 段落になり、種別は code（行頭記号を消して等幅にする目印）．"""
    src = _FM + f"### x\n\n{_FENCE}\na\nb\n{_FENCE}\n"
    slide, = parse(src).slides
    assert [(b.kind, b.text) for b in slide.blocks] == [("code", "a"), ("code", "b")]


def test_the_lines_are_verbatim():
    """行は書いたまま——行頭マーカーもサイズトークンも ``<br>`` も解釈しない．

    コードブロックの中で ``- `` が箇条書きになったら、それはもうコードではない。
    """
    src = _FM + (f"### x\n\n{_FENCE}text\n"
                 "- これは箇条書きではない\n"
                 "1. これも採番されない\n"
                 "{+1} サイズトークンも文字のまま\n"
                 "→ 矢印も文字のまま\n"
                 "a<br>b\n"
                 f"{_FENCE}\n")
    slide, = parse(src).slides
    assert [b.text for b in slide.blocks] == [
        "- これは箇条書きではない",
        "1. これも採番されない",
        "{+1} サイズトークンも文字のまま",
        "→ 矢印も文字のまま",
        "a<br>b",
    ]


def test_indentation_is_kept():
    """字下げを保つ（保たないとコードとして読めない）．"""
    src = _FM + f"### x\n\n{_FENCE}\ndef f():\n    return 1\n{_FENCE}\n"
    slide, = parse(src).slides
    assert [b.text for b in slide.blocks] == ["def f():", "    return 1"]


def test_blank_lines_inside_the_block_are_kept():
    """途中の空行も残す（セッションの区切りが消えると読めない）．"""
    src = _FM + f"### x\n\n{_FENCE}\na\n\nb\n{_FENCE}\n"
    slide, = parse(src).slides
    assert [b.text for b in slide.blocks] == ["a", "", "b"]


def test_the_info_string_is_free_and_ignored():
    """info string は何を書いてもよく、md2pptx は読み飛ばす（強調はしない）．"""
    for info in ("", "text", "console", "http", "python"):
        src = _FM + f"### x\n\n{_FENCE}{info}\nhello\n{_FENCE}\n"
        slide, = parse(src).slides
        assert [b.text for b in slide.blocks] == ["hello"], info


def test_reserved_fences_still_do_their_own_thing():
    """``flow`` / ``image`` / ``note`` は従来どおり（コードにはならない）．"""
    src = _FM + (f"### x\n\n{_FENCE}flow\n[a] -> [b]\n{_FENCE}\n\n"
                 f"{_FENCE}note\nノート本文\n{_FENCE}\n")
    slide, = parse(src).slides
    assert not [b for b in slide.blocks if isinstance(b, Line)]
    assert slide.notes == "ノート本文"


def test_code_is_rendered_in_a_monospace_font(tmp_path):
    """描画は等幅フォントで、行頭記号は付かない．

    フォント名は front matter の ``mono_font`` で変えられる。等幅かどうかは
    見た目の好みではなく**桁が揃うか**という機能なので、ここだけはテーマに
    委ねきれない（既定を持ち、変えたい人は front matter で変える）。
    """
    src = _FM + f"### x\n\n{_FENCE}\nGET / HTTP/1.1\n{_FENCE}\n"
    slide, = _build(tmp_path, src).slides
    text, font = _body_paragraphs(slide)[0]
    assert text == "GET / HTTP/1.1"
    assert font == "Consolas"


def test_the_monospace_font_can_be_changed(tmp_path):
    """``mono_font:`` で等幅フォントを差し替えられる．"""
    src = ("---\ntheme: t.pptx\nmono_font: Menlo\n---\n\n"
           f"### x\n\n{_FENCE}\nls -l\n{_FENCE}\n")
    slide, = _build(tmp_path, src).slides
    assert _body_paragraphs(slide)[0] == ("ls -l", "Menlo")


def test_code_sits_between_the_prose(tmp_path):
    """導入文・結論文と一緒に置ける（地の文と同じ枠に流れる）．

    **等幅になるのはコードの行だけ**——同じ枠に流し込む以上、地の文まで
    等幅にしてしまう取り違えは起こりうるし、桁の揃った本文は読みにくい。
    """
    src = _FM + (f"### x\n\n"
                 "以下のようにリクエストする\n"
                 f"{_FENCE}\nGET /\n{_FENCE}\n"
                 "→ 応答が返る\n")
    slide, = _build(tmp_path, src).slides
    assert _body_paragraphs(slide) == [
        ("以下のようにリクエストする", None),
        ("GET /", "Consolas"),
        ("→ 応答が返る", None),
    ]


def test_an_unclosed_fence_stops(tmp_path):
    """閉じ忘れたフェンスは行番号付きで止める．

    黙って末尾まで飲み込むと、以降のスライドが丸ごと消えたデッキが出る。
    """
    src = _FM + f"### x\n\n{_FENCE}\nGET /\n\n### y\n\n- b\n"
    with pytest.raises(ValueError, match="unclosed code fence at line"):
        parse(src)
