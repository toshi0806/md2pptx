#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""表・図の帯が結論文に食い込まないことを固定するテスト（Issue #131）．

結論文が実際に始まる位置は、本文プレースホルダに流し込んだ**空行の数**で決まる。
帯の高さをそれとは別の値（``band_h``）から取ると、空行数の切り捨てぶんだけ
帯のほうが下まで伸び、表が結論文の上に重なる。**警告も出ない**ので、
pptx を開くまで気づけない。

実測（cn2026-05「まとめ」右カラム）では、表の下端 5.99in に対して
結論文の開始位置が約 5.65in——0.34in 重なっていた。

結論文の y は python-pptx から直接は読めない（実際の行組みは PowerPoint がやる）ので、
**render と同じ規則**——枠の上端から段落 1 行ぶん（行送り＋段落前アキ）を積む——で求める。
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Emu, Pt

from md2pptx import render
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"

_TABLE = """| 型 | プロトコル | 使う場所 |
|:--|:--|:--|
| 距離ベクトル | RIP | AS 内（小規模） |
| リンク状態 | OSPF | AS 内（大規模） |
| パスベクトル | BGP | AS 間 |
"""


def _build(tmp_path, src):
    tmp_path.mkdir(parents=True, exist_ok=True)
    theme = tmp_path / "theme.pptx"
    Presentation().save(str(theme))
    out = tmp_path / "out.pptx"
    r = render.Renderer(str(theme))
    r.render(parse(src))
    r.save(str(out))
    # 段落 1 行ぶんの高さ．render と同じ規則で数える——行送りに加えて、
    # テーマの段落前アキ（``spcBef``）も含む（Issue #145）．
    return Presentation(str(out)), r._para_height(0, r._body_font_size())


def _objects(slide):
    """表・図（プレースホルダでないシェイプ）を返す．"""
    return [sh for sh in slide.shapes if not sh.is_placeholder]


def _only_object(slide):
    """表・図がちょうど1つあることを確かめて返す．"""
    objs = _objects(slide)
    assert len(objs) == 1, f"表・図が1つのはずが {len(objs)} 個: {objs}"
    return objs[0]


def _body_of(slide, needle="結論文"):
    """``needle`` を含む段落を持つ本文プレースホルダと、その段落番号を返す．"""
    for sh in slide.shapes:
        if not sh.has_text_frame or sh == slide.shapes.title:
            continue
        for i, p in enumerate(sh.text_frame.paragraphs):
            if needle in p.text:
                return sh, i
    pytest.fail(f"結論文 {needle!r} が本文に無い")


def _conclusion_top(slide, line_h, needle="結論文"):
    """結論文の描き始めの y（EMU）．枠の上端＋段落番号×行高．

    プレースホルダ自身の上マージン（既定 0.05in）は**わざと足さない**。
    足せば結論文の位置が下がり、この判定は緩くなる。回帰テストとしては
    厳しいほうへ倒しておきたい——実際の描画はここで求めた位置より
    わずかに下から始まる。
    """
    sh, i = _body_of(slide, needle)
    return sh.top + i * line_h


def _assert_above(slide, line_h):
    objs = _objects(slide)
    assert objs, "表・図が描かれていない"
    concl_top = _conclusion_top(slide, line_h)
    bottom = max(o.top + o.height for o in objs)
    assert bottom <= concl_top, (
        f"帯の下端 {Emu(bottom).inches:.2f}in が "
        f"結論文 {Emu(concl_top).inches:.2f}in に食い込んでいる")


# ---------------------------------------------------------------- 単一カラム

def test_table_stops_above_the_conclusion(tmp_path):
    """表の下端が、結論文の描き始めより上にある．"""
    prs, line_h = _build(
        tmp_path, _FM + "### 表\n\n導入文\n\n" + _TABLE + "\n→ 結論文\n")
    _assert_above(prs.slides[-1], line_h)


def test_figure_stops_above_the_conclusion(tmp_path):
    """フロー図でも同じ（帯の高さの求め方は表と共通）．"""
    prs, line_h = _build(
        tmp_path,
        _FM + "### 図\n\n導入文\n\n```flow\n[#a A] -> [#b B]\n```\n\n→ 結論文\n")
    _assert_above(prs.slides[-1], line_h)


def test_no_intro_line(tmp_path):
    """導入文が無く結論文だけの場合（``nb == 0``）も食い込まない．"""
    prs, line_h = _build(tmp_path, _FM + "### 表\n\n" + _TABLE + "\n→ 結論文\n")
    _assert_above(prs.slides[-1], line_h)


# ---------------------------------------------------------------- 多カラム

def test_table_in_a_column_stops_above_the_conclusion(tmp_path):
    """2カラムの右側に置いた表でも食い込まない（Issue #131 の再現形）．"""
    prs, line_h = _build(
        tmp_path,
        _FM + "### まとめ\n\n- 左の話\n\n<!-- @col -->\n\n" + _TABLE + "\n→ 結論文\n")
    _assert_above(prs.slides[-1], line_h)


_TALL_TABLE = "| 行 | 値 |\n|:--|:--|\n" + "".join(
    f"| 行{i} | 値{i} |\n" for i in range(1, 13))


# ---------------------------------------------------------------- 削りすぎない

def test_the_band_keeps_most_of_the_frame(tmp_path):
    """食い込みを直すために帯を削りすぎていない．

    許容は **3行ぶん**。原稿の地の文が導入文と結論文の 2 行あり、そこへ
    空行数の切り捨てで最大 1 行が加わる（2 + 1）。この修正で減るのはその
    1 行ぶんだけで、それ以上痩せていたら削りすぎ。

    帯の高さは**帯いっぱいになる表**で測る。短い表は帯を埋めない（#165）ので、
    表の高さを帯の高さの代わりに使えるのは、収まりきらない表だけ。
    """
    prs, line_h = _build(
        tmp_path, _FM + "### 表\n\n導入文\n\n" + _TALL_TABLE + "\n→ 結論文\n")
    slide = prs.slides[-1]
    tbl = _only_object(slide)
    body, _ = _body_of(slide)
    assert tbl.height >= body.height - 3 * line_h


# ---------------------------------------------------------------- overflow

def test_overflow_still_extends_below_the_conclusion(tmp_path):
    """``@overflow: true`` は従来どおり結論文より下まで伸びる（対象外）．

    はみ出しは**指定した人が選んだ挙動**なので、この修正で塞いではいけない。
    帯に収まらない高さの表で、指定あり／なしの差をそのまま見る。
    """
    body_src = "### 表\n{d}\n導入文\n\n" + _TALL_TABLE + "\n→ 結論文\n"
    plain, line_h = _build(tmp_path / "a", _FM + body_src.format(d=""))
    over, _ = _build(tmp_path / "b",
                     _FM + body_src.format(d="<!-- @overflow: true -->\n"))

    plain_tbl = _only_object(plain.slides[-1])
    over_tbl = _only_object(over.slides[-1])
    concl_top = _conclusion_top(plain.slides[-1], line_h)

    assert plain_tbl.top + plain_tbl.height <= concl_top
    assert over_tbl.top + over_tbl.height > concl_top


# ---------------------------------------------------------------- 最小高の警告

def test_warns_when_the_band_hits_its_minimum(tmp_path, capsys):
    """地の文が多すぎて帯が最小高（0.8in）に張り付いたら警告する．

    下限は図を読める大きさに保つためのもので、外すわけにはいかない。
    ただしその結果として結論文へ食い込むので、**黙って重ねない**
    （気づけないのが Issue #131 の本体）。
    """
    prose = "".join(f"- 地の文 {i}\n" for i in range(1, 12))
    _build(tmp_path, _FM + "### 表\n\n" + prose + "\n" + _TABLE + "\n→ 結論文\n")
    assert "band hit its minimum height" in capsys.readouterr().err


def test_no_minimum_warning_for_an_ordinary_slide(tmp_path, capsys):
    """ふつうの分量では警告しない（出しすぎると読まれなくなる）．"""
    _build(tmp_path, _FM + "### 表\n\n導入文\n\n" + _TABLE + "\n→ 結論文\n")
    assert "band hit its minimum height" not in capsys.readouterr().err


# ------------------------------------------------------- 折り返す地の文（#158）

def test_a_wrapped_intro_pushes_the_band_down(tmp_path):
    """導入文が折り返したら、そのぶん帯が下がる（結論文がはみ出さない）．

    帯の高さを「地の文は1行ずつ」で数えていた頃は、折り返した1行ぶん帯が
    上へずれ、結論文が下の罫線を越えていた（cn2026-02 p.34）。
    """
    intro = "とても長い導入文の続きです" * 4      # 全角52字
    short = _FM + "### x\n\n- 短い導入\n\n" + _TABLE + "\n→ 結論\n"
    long_ = _FM + "### x\n\n- " + intro + "\n\n" + _TABLE + "\n→ 結論\n"
    a, _ = _build(tmp_path / "a", short)
    b, _ = _build(tmp_path / "b", long_)
    r = render.Renderer(str(tmp_path / "b" / "theme.pptx"))
    # 前提：長いほうが本当に折り返している（既定テンプレートの本文枠は
    # 8.5in ≈ 612pt、lvl1 は 18pt なので 30字ほどで折り返る）
    body, = [sh for sh in b.slides[-1].shapes
             if sh.is_placeholder and sh.placeholder_format.idx == 1]
    tf = body.text_frame
    avail_pt = (body.width - tf.margin_left - tf.margin_right) / 12700.0
    assert r._wrapped_lines(intro, r._body_font_size(), avail_pt) > 1, \
        "前提が崩れた（導入文が折り返らない）"
    assert _only_object(b.slides[-1]).top > _only_object(a.slides[-1]).top


# --------------------------------------------- 短い表を引き伸ばさない（#165）

def test_a_short_table_is_not_stretched(tmp_path):
    """行数の少ない表は、帯いっぱいに広げない．

    ``add_table`` に帯の高さをそのまま渡すと、PowerPoint がそれを行へ配分して
    1 行が 2.5cm ほどになる（cn2026-05 p.17）。
    """
    src = _FM + "### x\n\n| a | b |\n|:--|:--|\n| 1 | 2 |\n"
    prs, line_h = _build(tmp_path, src)
    tbl = _only_object(prs.slides[-1])
    # 1 行の高さで見る——帯いっぱいに広げると 1 行 2.5cm 級になる。
    # 1 行の文字は本文標準サイズなので、その 2 倍を超えたら広げすぎ。
    assert tbl.height / 2 < 2 * line_h


def test_a_short_table_is_centred_in_its_band(tmp_path):
    """余った空きは上下へ均等に——帯の上端に貼りつけない．"""
    src = _FM + "### x\n\n- 導入\n\n| a | b |\n|:--|:--|\n| 1 | 2 |\n\n→ 結論\n"
    prs, line_h = _build(tmp_path, src)
    slide = prs.slides[-1]
    tbl = _only_object(slide)
    body, = [sh for sh in slide.shapes
             if sh.is_placeholder and sh.placeholder_format.idx == 1]
    above = tbl.top - (body.top + line_h)          # 導入文の下から表の上まで
    below = (body.top + body.height) - (tbl.top + tbl.height) - line_h
    assert above > 0 and below > 0
    assert abs(above - below) < line_h


def test_a_tall_table_still_uses_the_whole_band(tmp_path):
    """帯に収まらない表は従来どおり（回帰させない）．"""
    rows = "".join(f"| {i} | {i} |\n" for i in range(14))
    src = _FM + "### x\n\n| a | b |\n|:--|:--|\n" + rows
    prs, _ = _build(tmp_path, src)
    slide = prs.slides[-1]
    tbl = _only_object(slide)
    body_h = max(sh.height for sh in slide.shapes if sh.is_placeholder)
    assert tbl.height > body_h * 0.8
