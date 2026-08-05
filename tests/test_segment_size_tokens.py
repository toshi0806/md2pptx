#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""``<br>`` セグメントごとの相対サイズを固定するテスト（Issue #82）．

``<br>`` は段落を分けずに行だけを折る．python-pptx はこれを ``\\v`` として受け取り，
セグメントごとに ``a:r`` を作って ``a:br`` でつなぐので，**run と IR のセグメントが
順に 1 対 1 で対応する**．そこにサイズを書けば，段落を分けずに一部だけ大きさを
変えられる．

これが要るのはタイトル枠に副題を収めるため．テーマはタイトル枠と副題枠の間に
罫線を引いていることがあり（実測），副題を副題枠へ移すと**著者情報の塊に混ざる**．
副題はタイトル枠の中に，主題より少し小さく置く必要がある．

    # 主題<br>{-2} 副題

格納先が本文行とタイトルで違うのは，書き込む先が違うため．

- 本文行の行頭トークンは ``Line.size_delta`` で，段落の既定文字書式（``defRPr``）へ
  書く．そうしないと**ビュレットや採番記号のサイズが本文とずれる**．
  2 番目以降は ``Line.seg_deltas``（``[0]`` は常に None）で run へ書く．
- タイトルにはその記号が無いので ``Slide.title_deltas`` だけ．``[0]`` も有効．

テーマは python-pptx 同梱の既定テンプレートから作るので，リポジトリの
``OfficeTheme.pptx`` にも実 PowerPoint にも依存しない．このテンプレートは
タイトル・本文とも lvl1 をレイアウトで上書きしないので，基点はマスターの値
（表題 44pt / 本文 32pt）になる．そこから丸めた実サイズは次のとおり．

- 表題 44pt … ``{-2}`` → 35pt（44 ÷ 1.125² ＝ 34.77）／``{-1}`` → 39pt
- 本文 32pt … ``{-2}`` → 25pt（32 ÷ 1.125² ＝ 25.28）／``{+1}`` → 36pt

``_apply_size_delta`` / ``_apply_segment_deltas`` が pt 単位で ``round`` してから
``Pt()`` へ渡すので，``sz`` に端数は出ない．
"""
from __future__ import annotations

from pptx import Presentation

from md2pptx import render
from md2pptx.ir import Line, Slide
from md2pptx.parser import parse


def _theme(tmp_path):
    path = tmp_path / "theme.pptx"
    Presentation().save(str(path))
    return str(path)


def _runs(tmp_path, src, ph_idx, name="out.pptx"):
    """描画して保存し，先頭スライドの指定枠の run を (本文, サイズ pt) で返す．

    保存した pptx を読み直すのは，利用者が開くのがファイルのほうだから
    （test_text_language.py と同じ理由）．サイズ未指定（テーマ継承）は None．
    """
    out = tmp_path / name
    r = render.Renderer(_theme(tmp_path))
    r.render(parse(src))
    r.save(str(out))
    slide = Presentation(str(out)).slides[0]
    ph = next((s for s in slide.placeholders
               if s.placeholder_format.idx == ph_idx), None)
    assert ph is not None, (
        f"placeholder idx={ph_idx} not found "
        f"(layout {slide.slide_layout.name!r})")
    return [(run.text, run.font.size.pt if run.font.size else None)
            for para in ph.text_frame.paragraphs for run in para.runs]


def _para_size(tmp_path, src, ph_idx, name="para.pptx"):
    """先頭スライドの指定枠の，先頭段落の**段落既定サイズ**（pt）を返す．

    ``_runs`` が見るのは run の ``rPr`` で，行頭トークンが書かれる
    ``pPr/defRPr`` は映らない．役割が分かれていることを確かめるには両方要る．
    """
    out = tmp_path / name
    r = render.Renderer(_theme(tmp_path))
    r.render(parse(src))
    r.save(str(out))
    slide = Presentation(str(out)).slides[0]
    ph = next((s for s in slide.placeholders
               if s.placeholder_format.idx == ph_idx), None)
    assert ph is not None, f"placeholder idx={ph_idx} not found"
    size = ph.text_frame.paragraphs[0].font.size
    return size.pt if size else None


# ------------------------------------------------------------------ parser
def test_parser_splits_tokens_off_each_segment():
    """2 番目以降のセグメント先頭のトークンが seg_deltas へ移る．"""
    deck = parse("---\ntheme: t.pptx\n---\n\n## 見出し\n- {+1} A<br>{-2} B<br>C\n")
    line = deck.slides[0].blocks[0]
    assert (line.text, line.size_delta, line.seg_deltas) == (
        "A\vB\vC", 1, [None, -2, None])


def test_parser_takes_the_first_title_segment_too():
    """タイトルは先頭セグメントの段数も取る（本文行と違い格納先が 1 つ）．"""
    deck = parse("---\ntheme: t.pptx\n---\n\n"
                 "# 主題<br>{-2} 副題\n\n## {+1} 大見出し<br>{-1} 小さく\n")
    assert deck.slides[0].title_deltas == [None, -2]
    assert deck.slides[1].title_deltas == [1, -1]


def test_arrow_lines_take_segment_tokens():
    """``→`` 行でもセグメントの段数を取れる（行種で挙動を変えない）．"""
    deck = parse("---\ntheme: t.pptx\n---\n\n## 見出し\n→ {-1} 結論<br>{+2} 強め\n")
    line = deck.slides[0].blocks[0]
    assert (line.text, line.size_delta, line.seg_deltas) == (
        "→ 結論\v強め", -1, [None, 2])


def test_plain_br_keeps_meaning_nothing():
    """トークンの無い ``<br>`` は従来どおり段数を持たない．"""
    deck = parse("---\ntheme: t.pptx\n---\n\n## 見出し<br>続き\n- A<br>B\n")
    assert deck.slides[0].title_deltas == [None, None]
    assert deck.slides[0].blocks[0].seg_deltas == [None, None]


def test_whitespace_around_br_is_eaten():
    """``<br>`` の前後の空白は従来どおり消える（この PR で変えていない）．

    ``_RE_BR`` は ``\\s*<br\\s*/?>\\s*`` で前後の空白ごと置換するので，
    先頭セグメントの**末尾**空白も落ちる．トークンを剥がすようになっても
    そこは変わらない——``A <br>{-2} B`` は ``"A\\vB"`` であって ``"A \\vB"`` ではない．
    行末に空白を残すエディタ設定でも表示がずれない，という既存の性質に乗っている．
    """
    for src in ("A <br>{-2} B", "A<br> {-2} B", "A  <br>  {-2}  B"):
        deck = parse(f"---\ntheme: t.pptx\n---\n\n## 見出し\n- {src}\n")
        line = deck.slides[0].blocks[0]
        assert (line.text, line.seg_deltas) == ("A\vB", [None, -2]), src
    # タイトルでも同じ（通す場所は _split_br の 1 つ）．
    deck = parse("---\ntheme: t.pptx\n---\n\n## A <br>{-2} B\n")
    assert (deck.slides[0].title, deck.slides[0].title_deltas) == (
        "A\vB", [None, -2])


# ---------------------------------------------------------------------- IR
def test_ir_normalises_the_parallel_lists():
    """並行配列はセグメント数へ揃う（render が添字で run に対応付けるため）．

    長さがずれると**別のセグメントにサイズが付く**——見た目の崩れと違い、
    1 つずれたサイズは正しく見えてしまうので、構築時に潰しておく．
    """
    assert Line(text="A\vB\vC").seg_deltas == [None, None, None]
    assert Line(text="A\vB", seg_deltas=[9, 9, 9, 9]).seg_deltas == [None, 9]
    # Line.seg_deltas[0] は常に None（先頭は size_delta が持つ）．
    assert Line(text="A\vB", seg_deltas=[5, -1]).seg_deltas == [None, -1]
    assert Slide(title="X\vY").title_deltas == [None, None]
    assert Slide(title=None).title_deltas == []


# ------------------------------------------------------------------ render
def test_render_sizes_the_title_segments(tmp_path):
    """副題がタイトル枠の中に，主題の 2 段下で出る（この機能の目的）．"""
    assert _runs(tmp_path, """---
theme: t.pptx
---

# 主題<br>{-2} 副題
<!-- @layout: 0 -->
""", ph_idx=0) == [("主題", None), ("副題", 35.0)]


def test_render_sizes_the_first_title_segment(tmp_path):
    """タイトルは先頭セグメントにも段数を書ける（``Line`` との非対称の実物）．

    44pt × 1.125 ＝ 49.5 を丸めて 50pt．run へ入るので段落既定は空のまま．
    """
    src = """---
theme: t.pptx
---

## {+1} 大見出し<br>{-1} 小さく
"""
    assert _runs(tmp_path, src, ph_idx=0) == [("大見出し", 50.0), ("小さく", 39.0)]
    assert _para_size(tmp_path, src, ph_idx=0) is None


def test_render_sizes_the_body_segments(tmp_path):
    """本文行でも 2 番目以降の run だけサイズが変わる．"""
    assert _runs(tmp_path, """---
theme: t.pptx
---

## 見出し

- A<br>{-2} 小さい続き
""", ph_idx=1) == [("A", None), ("小さい続き", 25.0)]


def test_segment_base_ignores_the_line_token(tmp_path):
    """セグメントの段数は行の段数から数え直さない．

    ``{-2}`` は行が ``{+1}`` でも「テーマ既定の 2 段下」＝ 25pt のまま．
    記法の意味を段で変えないほうが，書く側が結果を予測できる．
    """
    src = """---
theme: t.pptx
---

## 見出し

- {+1} 大きい行<br>{-2} 小さい続き
"""
    assert _runs(tmp_path, src, ph_idx=1) == [("大きい行", None),
                                              ("小さい続き", 25.0)]
    # 行頭の {+1} は run ではなく**段落の既定文字書式**へ入る（32 × 1.125 ＝ 36pt）．
    # そうでないとビュレットが本文と違う大きさで出る．先頭 run が None なのは
    # 「サイズ未指定」であって「調整されていない」ではない——段落側から継承する．
    assert _para_size(tmp_path, src, ph_idx=1) == 36.0
