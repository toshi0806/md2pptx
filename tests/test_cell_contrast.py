#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""塗ったセルの文字色を固定するテスト（Issue #148）．

セルの個別着色（#126）で**濃い色を指定すると、文字が黒のままで読めなかった**。
cn2026-theme の `accent2` は #3B812F（濃い緑）、`bg2` は #5F5F5F（グレー）で、
経路表や NAT 変換表がまるごと読めない状態になっていた。

見出し行はもともと「塗って白文字」にしていた。塗った以上は文字色も塗りに
合わせるという同じ規則を、全てのセルへ広げる——**見出しだけの特別扱いをやめる**
ので、明るいアクセント色のテーマでは見出しが黒文字になる（そちらが読める）。
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn

from md2pptx import render
from md2pptx.parser import parse


_FM = "---\ntheme: t.pptx\n---\n\n"


def _theme(tmp_path, **colors):
    """clrScheme を差し替えたテーマを作る（``accent2="3B812F"`` のように渡す）．"""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.oxml import parse_xml
    tmp_path.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    # テーマは python-pptx では XML パートではなく生のバイト列なので、
    # 読んで書き換えて戻す（`part._element` は無い）。
    part = prs.slide_masters[0].part.part_related_by(RT.THEME)
    root = parse_xml(part.blob)
    scheme = root.find(qn("a:themeElements") + "/" + qn("a:clrScheme"))
    for name, rgb in colors.items():
        el = scheme.find(qn("a:" + name))
        assert el is not None, f"{name} が clrScheme に無い"
        for child in list(el):
            el.remove(child)
        el.append(el.makeelement(qn("a:srgbClr"), {"val": rgb}))
    part._blob = root.xml.encode("utf-8")
    path = tmp_path / "theme.pptx"
    prs.save(str(path))
    return path


def _build(tmp_path, src, theme):
    r = render.Renderer(str(theme))
    r.render(parse(src))
    out = tmp_path / "out.pptx"
    r.save(str(out))
    return Presentation(str(out))


def _table(prs):
    for sh in prs.slides[-1].shapes:
        if sh.has_table:
            return sh.table
    raise AssertionError("表が無い")


def _run_color(cell):
    """セル 1 つめの run の文字色（テーマ色名 or None）．"""
    runs = cell.text_frame.paragraphs[0].runs
    assert runs, f"run が無い: {cell.text!r}"
    return runs[0].font.color


def _src(cells, header="| 見出し | 値 |\n|:--|:--|\n"):
    return _FM + "### x\n\n" + header + cells


# ---------------------------------------------------------------- 本文セル

def test_a_dark_fill_gets_light_text(tmp_path):
    """濃い塗り（#3B812F）のセルは文字を背景色にする．"""
    theme = _theme(tmp_path, accent2="3B812F")
    prs = _build(tmp_path, _src("| a {accent2} | b |\n"), theme)
    cell = _table(prs).cell(1, 0)
    assert _run_color(cell).theme_color == MSO_THEME_COLOR.BACKGROUND_1


def test_a_light_fill_keeps_the_default_text(tmp_path):
    """明るい塗りでは文字色に触らない（テーマの本文色のまま）．"""
    theme = _theme(tmp_path, accent2="EEEEEE")
    prs = _build(tmp_path, _src("| a {accent2} | b |\n"), theme)
    cell = _table(prs).cell(1, 0)
    assert _run_color(cell).type is None


def test_an_unfilled_cell_is_untouched(tmp_path):
    """塗っていないセルは従来どおり（回帰させない）．"""
    theme = _theme(tmp_path, accent2="3B812F")
    prs = _build(tmp_path, _src("| a {accent2} | b |\n"), theme)
    assert _run_color(_table(prs).cell(1, 1)).type is None


@pytest.mark.parametrize("hexval,light", [("000000", True), ("FFFFFF", False)])
def test_a_hex_fill_is_judged_too(tmp_path, hexval, light):
    """16進で書いた塗りも同じ規則で判定する．"""
    theme = _theme(tmp_path)
    prs = _build(tmp_path, _src("| a {#%s} | b |\n" % hexval), theme)
    color = _run_color(_table(prs).cell(1, 0))
    if light:
        assert color.theme_color == MSO_THEME_COLOR.BACKGROUND_1
    else:
        assert color.type is None


def test_bg2_goes_through_the_colour_map(tmp_path):
    """``bg2`` はテーマの ``lt2``——``p:clrMap`` を踏まないと別の色を見てしまう．"""
    theme = _theme(tmp_path, lt2="5F5F5F")
    prs = _build(tmp_path, _src("| a {bg2} | b |\n"), theme)
    assert _run_color(_table(prs).cell(1, 0)).theme_color \
        == MSO_THEME_COLOR.BACKGROUND_1


# ---------------------------------------------------------------- 見出し行

def test_a_dark_header_keeps_light_text(tmp_path):
    """既定の見出し（accent2 で塗る）が濃ければ、従来どおり白文字．"""
    theme = _theme(tmp_path, accent2="3B812F")
    prs = _build(tmp_path, _src("| a | b |\n"), theme)
    assert _run_color(_table(prs).cell(0, 0)).theme_color \
        == MSO_THEME_COLOR.BACKGROUND_1


def test_a_light_header_gets_dark_text(tmp_path):
    """明るいアクセント色のテーマでは、見出しも黒文字になる．

    白文字の決め打ちだと、ここが**白地に白**になっていた。
    """
    theme = _theme(tmp_path, accent2="EEEEEE")
    prs = _build(tmp_path, _src("| a | b |\n"), theme)
    assert _run_color(_table(prs).cell(0, 0)).type is None


def test_the_header_stays_bold(tmp_path):
    """文字色の規則を変えても、見出しの太字は残す．"""
    theme = _theme(tmp_path, accent2="EEEEEE")
    prs = _build(tmp_path, _src("| a | b |\n"), theme)
    assert _table(prs).cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold


def test_an_overridden_header_fill_is_judged(tmp_path):
    """見出しセルを個別に塗ったら、その色で判定する（既定色ではなく）．"""
    theme = _theme(tmp_path, accent2="EEEEEE", accent6="35742A")
    prs = _build(tmp_path,
                 _src("| a | b |\n", header="| 見出し {accent6} | 値 |\n|:--|:--|\n"),
                 theme)
    tbl = _table(prs)
    assert _run_color(tbl.cell(0, 0)).theme_color == MSO_THEME_COLOR.BACKGROUND_1
    assert _run_color(tbl.cell(0, 1)).type is None       # 既定の accent2 は明るい


# ---------------------------------------------------------------- 判定そのもの

@pytest.mark.parametrize("rgb,dark", [
    ("000000", True), ("FFFFFF", False),
    ("3B812F", True),      # cn2026-theme の accent2
    ("5F5F5F", True),      # cn2026-theme の lt2（bg2）
    ("E2CAAA", False),     # cn2026-theme の accent5
    ("CC9900", False),     # cn2026-theme の accent1（金）
])
def test_the_luminance_threshold(rgb, dark):
    assert render.is_dark(rgb) is dark


def test_an_unreadable_colour_falls_back_to_light_text(tmp_path):
    """テーマから色を引けなければ、塗った以上は白文字にしておく．

    塗ったのに黒のままだと**読めない**——読めるかもしれない側へ倒す。
    """
    assert render.is_dark(None) is True
