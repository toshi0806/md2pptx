#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""ラダー図の文字が箱からはみ出さないことを固定するテスト（Issue #167）．

`plan_seq` は文字の幅を ``LABEL_PT``（16pt）で見積もって箱を決めるのに、
`render_seq` は**本文標準サイズ**で描いていた。cn2026-theme の lvl1 は 30pt なので
見積もりの 1.9 倍の字が入り、注記が `wrap=False` のまま**スライドの外へ**出ていた
（cn2026-04 p.32 の `MTU 1500 を超える`）。

python-pptx の既定テーマは本文 18pt でほぼ一致するため、それだけでは出ない。
**本文の大きいテーマを組み立てて確かめる。**
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

from md2pptx import render
from md2pptx.parser import parse
from md2pptx.seq import LABEL_PT


_FM = "---\ntheme: t.pptx\n---\n\n"

_SRC = (_FM + "### x\n\n```seq\nlifelines: 送信元, ルータ, 宛先\n"
        "送信元 -> ルータ: 4000 バイトのパケット\n"
        "note: MTU 1500 を超える\n"
        "ルータ -> 宛先: 断片1 (1480)\n```\n")


def _theme(tmp_path, size=3000):
    """本文 lvl1 が ``size``（1/100 pt）のテーマを作る．"""
    tmp_path.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    body = prs.slide_masters[0].element.find(
        qn("p:txStyles") + "/" + qn("p:bodyStyle"))
    for lvl in range(1, 10):
        el = body.find(qn(f"a:lvl{lvl}pPr"))
        if el is None:
            continue
        old = el.find(qn("a:defRPr"))
        if old is not None:
            el.remove(old)
        el.append(el.makeelement(qn("a:defRPr"), {"sz": str(size)}))
    path = tmp_path / "theme.pptx"
    prs.save(str(path))
    return path


def _build(tmp_path, src, size=3000):
    r = render.Renderer(str(_theme(tmp_path, size)))
    r.render(parse(src))
    out = tmp_path / "out.pptx"
    r.save(str(out))
    return Presentation(str(out)), r


def _texts(prs):
    """図として置かれたテキストボックス（プレースホルダでないもの）．"""
    return [sh for sh in prs.slides[-1].shapes
            if not sh.is_placeholder and sh.has_text_frame and sh.text_frame.text]


def _sizes(prs):
    out = set()
    for sh in _texts(prs):
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    out.add(r.font.size)
    return out


# ---------------------------------------------------------------- 大きさ

def test_the_text_is_not_drawn_larger_than_the_plan_assumes(tmp_path):
    """本文 30pt のテーマでも、図の文字は 16pt を超えない．"""
    prs, _ = _build(tmp_path, _SRC, size=3000)
    assert _sizes(prs), "図のテキストが無い（前提が崩れている）"
    assert max(_sizes(prs)) <= Pt(LABEL_PT)


def test_a_small_body_theme_is_unchanged(tmp_path):
    """本文が小さいテーマでは従来どおり本文サイズで描く（回帰させない）．"""
    prs, _ = _build(tmp_path / "s", _SRC, size=1200)
    assert max(_sizes(prs)) == Pt(12.0)


# ---------------------------------------------------------------- はみ出し

def test_nothing_runs_off_the_slide(tmp_path):
    """注記もラベルも、スライドの外へ出ない．

    ``wrap=False`` で描くので、箱が足りなければそのまま外へ出る。箱の右端が
    スライド内にあることを見る。
    """
    prs, _ = _build(tmp_path, _SRC, size=3000)
    for sh in _texts(prs):
        assert sh.left >= 0, sh.text_frame.text
        assert sh.left + sh.width <= prs.slide_width, sh.text_frame.text


def test_the_note_fits_the_box_it_was_planned_for(tmp_path):
    """注記の実幅が、plan の見積もった箱に収まる．"""
    prs, r = _build(tmp_path, _SRC, size=3000)
    note = [sh for sh in _texts(prs)
            if "MTU" in sh.text_frame.text]
    assert note, "注記が見つからない"
    sh = note[0]
    size = max(rr.font.size for p in sh.text_frame.paragraphs for rr in p.runs)
    assert r._text_width_pt(sh.text_frame.text, size / 12700) <= sh.width / 12700
